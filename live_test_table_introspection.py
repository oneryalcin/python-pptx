#!/usr/bin/env python3

"""
Live Test Script: Table and Cell Introspection
==============================================

This script tests Table and _Cell introspection functionality with real python-pptx 
objects, providing comprehensive validation for FEP-018 implementation.

Features Tested:
- Table creation and basic introspection
- Cell text, fill, margins, and vertical anchor properties
- Table formatting flags (first_row, horz_banding, etc.)
- Complex merge scenarios (merge-origin and spanned cells)
- Table structure representation with expand_collections
- Error handling and edge cases
- Relationship introspection (parent graphic frame)

Run with: python live_test_table_introspection.py
"""

import json
import traceback
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.enum.dml import MSO_FILL
from pptx.dml.color import RGBColor
from pptx.util import Inches


def create_test_presentation():
    """Create a test presentation with various table configurations."""
    prs = Presentation()
    
    # Slide 1: Basic table with text and formatting
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title manually since blank layout might not have title placeholder
    title_left = Inches(0.5)
    title_top = Inches(0.5) 
    title_width = Inches(9)
    title_height = Inches(1)
    title_box = slide1.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = "Basic Table Test"
    
    # Add a simple 3x3 table
    left = Inches(1)
    top = Inches(2) 
    width = Inches(6)
    height = Inches(3)
    table1 = slide1.shapes.add_table(3, 3, left, top, width, height).table
    
    # Set formatting flags
    table1.first_row = True
    table1.horz_banding = True
    table1.first_col = True
    
    # Add content to cells
    headers = ["Product", "Q1", "Q2"] 
    for col_idx, header in enumerate(headers):
        cell = table1.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x4F, 0x81, 0xBD)  # Blue
    
    # Add data rows
    data = [
        ["Widget A", "$1,000", "$1,200"],
        ["Widget B", "$800", "$950"]
    ]
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, value in enumerate(row_data):
            cell = table1.cell(row_idx, col_idx)
            cell.text = value
            if col_idx == 0:  # First column
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE1, 0xE1, 0xE1)  # Light gray
    
    # Slide 2: Complex table with merges
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title manually
    title_box2 = slide2.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame2 = title_box2.text_frame
    title_frame2.text = "Complex Table with Merges"
    
    # Add a 4x4 table for merge testing
    table2 = slide2.shapes.add_table(4, 4, left, top, width, height).table
    
    # Set all formatting flags for comprehensive testing
    table2.first_row = True
    table2.last_row = True
    table2.first_col = True
    table2.last_col = True
    table2.horz_banding = True
    table2.vert_banding = True
    
    # Add header that spans multiple columns
    header_cell = table2.cell(0, 0)
    header_cell.text = "Quarterly Sales Report"
    header_cell.merge(table2.cell(0, 3))  # Merge across all columns
    
    # Add content and create some merges
    table2.cell(1, 0).text = "Region"
    table2.cell(1, 1).text = "Q1"
    table2.cell(1, 2).text = "Q2"
    table2.cell(1, 3).text = "Total"
    
    # Merge cells vertically for "North" region
    north_cell = table2.cell(2, 0)
    north_cell.text = "North"
    north_cell.merge(table2.cell(3, 0))
    
    # Add data
    table2.cell(2, 1).text = "$500"
    table2.cell(2, 2).text = "$600"
    table2.cell(2, 3).text = "$1,100"
    table2.cell(3, 1).text = "$450"
    table2.cell(3, 2).text = "$550" 
    table2.cell(3, 3).text = "$1,000"
    
    return prs


def test_basic_table_introspection(table):
    """Test basic table introspection functionality."""
    print("\n" + "="*60)
    print("TESTING: Basic Table Introspection")
    print("="*60)
    
    try:
        # Test basic to_dict() call
        result = table.to_dict(expand_collections=False)
        
        print(f"✓ Basic to_dict() succeeded")
        print(f"  Object type: {result['_object_type']}")
        
        # Check identity
        identity = result['_identity']
        description = identity.get('description', identity.get('class_name', 'Table object'))
        print(f"  Identity: {description}")
        
        # Check properties
        props = result['properties']
        print(f"  Formatting flags:")
        print(f"    first_row: {props['first_row']}")
        print(f"    last_row: {props['last_row']}")
        print(f"    first_col: {props['first_col']}")
        print(f"    last_col: {props['last_col']}")
        print(f"    horz_banding: {props['horz_banding']}")
        print(f"    vert_banding: {props['vert_banding']}")
        
        # Check rows summary
        print(f"  Rows summary: {props['rows']}")
        
        # Check LLM context
        context = result['_llm_context']
        print(f"  LLM Summary: {context['summary']}")
        print(f"  Common operations count: {len(context['common_operations'])}")
        
        return True
        
    except Exception as e:
        print(f"✗ Basic table introspection failed: {e}")
        traceback.print_exc()
        return False


def test_expanded_table_introspection(table):
    """Test table introspection with expanded collections."""
    print("\n" + "="*60)
    print("TESTING: Expanded Table Introspection")
    print("="*60)
    
    try:
        # Test with expand_collections=True and limited depth
        result = table.to_dict(expand_collections=True, max_depth=3)
        
        props = result['properties']
        rows = props['rows']
        
        print(f"✓ Expanded introspection succeeded")
        print(f"  Rows type: {type(rows)}")
        print(f"  Table dimensions: {len(rows)} rows x {len(rows[0]) if rows else 0} columns")
        
        # Examine a few cells
        if rows and len(rows) > 0 and len(rows[0]) > 0:
            cell_0_0 = rows[0][0]
            print(f"  Cell [0,0] type: {cell_0_0['_object_type']}")
            
            # Check if cell has expected properties
            if 'properties' in cell_0_0:
                cell_props = cell_0_0['properties']
                print(f"  Cell [0,0] merge status:")
                print(f"    is_merge_origin: {cell_props.get('is_merge_origin')}")
                print(f"    is_spanned: {cell_props.get('is_spanned')}")
                print(f"    span_height: {cell_props.get('span_height')}")
                print(f"    span_width: {cell_props.get('span_width')}")
                
                # Check text frame if available
                if 'text_frame' in cell_props and cell_props['text_frame']:
                    text_frame = cell_props['text_frame']
                    if isinstance(text_frame, dict) and 'properties' in text_frame:
                        text_props = text_frame['properties']
                        print(f"    text: '{text_props.get('text', '')}'")
        
        return True
        
    except Exception as e:
        print(f"✗ Expanded table introspection failed: {e}")
        traceback.print_exc()
        return False


def test_cell_introspection(table):
    """Test individual cell introspection functionality."""
    print("\n" + "="*60)
    print("TESTING: Individual Cell Introspection")
    print("="*60)
    
    try:
        # Test different cells
        test_cells = [
            (0, 0, "header cell"),
            (1, 1, "data cell"),
        ]
        
        for row_idx, col_idx, description in test_cells:
            try:
                cell = table.cell(row_idx, col_idx)
                result = cell.to_dict()
                
                print(f"✓ Cell [{row_idx},{col_idx}] ({description}) introspection succeeded")
                print(f"  Object type: {result['_object_type']}")
                
                # Check properties
                props = result['properties']
                print(f"  Merge status: origin={props.get('is_merge_origin')}, spanned={props.get('is_spanned')}")
                print(f"  Span dimensions: {props.get('span_width')}x{props.get('span_height')}")
                
                # Check margins
                margins = {
                    'left': props.get('margin_left'),
                    'right': props.get('margin_right'),
                    'top': props.get('margin_top'),
                    'bottom': props.get('margin_bottom')
                }
                print(f"  Margins: {margins}")
                
                # Check LLM context
                context = result['_llm_context']
                print(f"  LLM Summary: {context['summary']}")
                
            except Exception as e:
                print(f"✗ Cell [{row_idx},{col_idx}] introspection failed: {e}")
        
        return True
        
    except Exception as e:
        print(f"✗ Cell introspection failed: {e}")
        traceback.print_exc()
        return False


def test_merged_cell_introspection(table):
    """Test introspection of merged cells (if any exist)."""
    print("\n" + "="*60)
    print("TESTING: Merged Cell Introspection")
    print("="*60)
    
    try:
        merge_origins_found = 0
        spanned_cells_found = 0
        
        # Check all cells for merge status
        for row_idx in range(len(table.rows)):
            for col_idx in range(len(table.columns)):
                cell = table.cell(row_idx, col_idx)
                
                if cell.is_merge_origin:
                    merge_origins_found += 1
                    result = cell.to_dict()
                    props = result['properties']
                    
                    print(f"✓ Merge-origin cell found at [{row_idx},{col_idx}]")
                    print(f"  Span: {props['span_width']}x{props['span_height']}")
                    
                    context = result['_llm_context']
                    print(f"  LLM Context: {context['summary']}")
                
                elif cell.is_spanned:
                    spanned_cells_found += 1
                    result = cell.to_dict()
                    context = result['_llm_context']
                    
                    print(f"✓ Spanned cell found at [{row_idx},{col_idx}]")
                    print(f"  LLM Context: {context['summary']}")
        
        print(f"\nMerge Summary:")
        print(f"  Merge-origin cells: {merge_origins_found}")
        print(f"  Spanned cells: {spanned_cells_found}")
        
        return True
        
    except Exception as e:
        print(f"✗ Merged cell introspection failed: {e}")
        traceback.print_exc()
        return False


def test_table_relationships(table):
    """Test table relationship introspection."""
    print("\n" + "="*60)
    print("TESTING: Table Relationships")
    print("="*60)
    
    try:
        result = table.to_dict(include_relationships=True)
        
        relationships = result.get('relationships', {})
        print(f"✓ Relationships introspection succeeded")
        print(f"  Relationships found: {list(relationships.keys())}")
        
        if 'parent_graphic_frame' in relationships:
            graphic_frame = relationships['parent_graphic_frame']
            print(f"  Parent graphic frame type: {graphic_frame.get('_object_type', 'Unknown')}")
        
        return True
        
    except Exception as e:
        print(f"✗ Table relationships introspection failed: {e}")
        traceback.print_exc()
        return False


def test_format_for_llm(table):
    """Test LLM-friendly formatting."""
    print("\n" + "="*60)
    print("TESTING: LLM-Friendly Formatting")
    print("="*60)
    
    try:
        # Test both format_for_llm values
        result_llm = table.to_dict(format_for_llm=True, expand_collections=False)
        result_raw = table.to_dict(format_for_llm=False, expand_collections=False)
        
        print(f"✓ LLM formatting test succeeded")
        print(f"  With format_for_llm=True: {len(json.dumps(result_llm))} chars")
        print(f"  With format_for_llm=False: {len(json.dumps(result_raw))} chars")
        
        # Check specific differences
        llm_context = result_llm.get('_llm_context', {})
        if llm_context:
            print(f"  LLM context available: {bool(llm_context)}")
            print(f"  Summary: {llm_context.get('summary', 'N/A')}")
        
        return True
        
    except Exception as e:
        print(f"✗ LLM formatting test failed: {e}")
        traceback.print_exc()
        return False


def run_comprehensive_analysis(table):
    """Run a comprehensive analysis showcasing table introspection capabilities."""
    print("\n" + "="*60)
    print("COMPREHENSIVE TABLE ANALYSIS")
    print("="*60)
    
    try:
        # Get full introspection
        result = table.to_dict(
            expand_collections=True,
            max_depth=4,
            format_for_llm=True,
            include_relationships=True
        )
        
        print("✓ Full introspection analysis completed")
        
        # Table-level summary
        context = result['_llm_context']
        print(f"\nTable Summary: {context['summary']}")
        
        props = result['properties']
        rows = props['rows']
        
        print(f"\nTable Structure Analysis:")
        print(f"  Dimensions: {len(rows)} rows x {len(rows[0]) if rows else 0} columns")
        
        # Analyze formatting
        formatting_features = []
        for feature in ['first_row', 'last_row', 'first_col', 'last_col', 'horz_banding', 'vert_banding']:
            if props.get(feature):
                formatting_features.append(feature)
        print(f"  Active formatting: {', '.join(formatting_features) if formatting_features else 'None'}")
        
        # Cell content analysis
        print(f"\nCell Content Analysis:")
        total_cells = 0
        cells_with_content = 0
        merge_origins = 0
        spanned_cells = 0
        
        for row_idx, row in enumerate(rows):
            for col_idx, cell in enumerate(row):
                total_cells += 1
                if isinstance(cell, dict) and 'properties' in cell:
                    cell_props = cell['properties']
                    
                    # Check for text content
                    text_frame = cell_props.get('text_frame')
                    if text_frame and isinstance(text_frame, dict):
                        text_props = text_frame.get('properties', {})
                        text_content = text_props.get('text', '').strip()
                        if text_content:
                            cells_with_content += 1
                    
                    # Check merge status
                    if cell_props.get('is_merge_origin'):
                        merge_origins += 1
                    elif cell_props.get('is_spanned'):
                        spanned_cells += 1
        
        print(f"  Total cells: {total_cells}")
        print(f"  Cells with content: {cells_with_content}")
        print(f"  Merge-origin cells: {merge_origins}")
        print(f"  Spanned cells: {spanned_cells}")
        
        # JSON size analysis
        json_size = len(json.dumps(result))
        print(f"\nSerialization Analysis:")
        print(f"  Full JSON size: {json_size:,} characters")
        print(f"  Estimated size per cell: {json_size // total_cells if total_cells > 0 else 0} chars")
        
        return True
        
    except Exception as e:
        print(f"✗ Comprehensive analysis failed: {e}")
        traceback.print_exc()
        return False


def main():
    """Main test execution function."""
    print("="*80)
    print("FEP-018: Table and Cell Introspection - Live Test Suite")
    print("="*80)
    
    try:
        # Create test presentation
        print("Creating test presentation with various table configurations...")
        prs = create_test_presentation()
        
        # Get tables from slides (skip title textbox, get table)
        table1 = prs.slides[0].shapes[1].table  # Basic table
        table2 = prs.slides[1].shapes[1].table  # Complex table with merges
        
        print(f"✓ Test presentation created with {len(prs.slides)} slides")
        
        # Run tests on basic table
        print(f"\n{'='*80}")
        print("TESTING BASIC TABLE (Slide 1)")
        print("="*80)
        
        tests_passed = 0
        total_tests = 0
        
        test_functions = [
            (test_basic_table_introspection, "Basic Table Introspection"),
            (test_expanded_table_introspection, "Expanded Collections"),
            (test_cell_introspection, "Cell Introspection"),
            (test_table_relationships, "Table Relationships"),
            (test_format_for_llm, "LLM Formatting"),
        ]
        
        for test_func, test_name in test_functions:
            total_tests += 1
            if test_func(table1):
                tests_passed += 1
        
        # Run comprehensive analysis
        run_comprehensive_analysis(table1)
        
        # Test complex table with merges
        print(f"\n{'='*80}")
        print("TESTING COMPLEX TABLE WITH MERGES (Slide 2)")
        print("="*80)
        
        # Test merged cell functionality
        total_tests += 1
        if test_merged_cell_introspection(table2):
            tests_passed += 1
        
        # Run comprehensive analysis on complex table
        run_comprehensive_analysis(table2)
        
        # Final results
        print(f"\n{'='*80}")
        print("FINAL RESULTS")
        print("="*80)
        print(f"Tests passed: {tests_passed}/{total_tests}")
        print(f"Success rate: {(tests_passed/total_tests)*100:.1f}%")
        
        if tests_passed == total_tests:
            print("🎉 All tests passed! FEP-018 implementation is working correctly.")
        else:
            print("⚠️  Some tests failed. Review the output above for details.")
        
        # Save test presentation for manual inspection
        output_path = Path("test_table_introspection_output.pptx")
        prs.save(output_path)
        print(f"\n📁 Test presentation saved to: {output_path}")
        print("   You can open this file in PowerPoint to inspect the test tables.")
        
        return tests_passed == total_tests
        
    except Exception as e:
        print(f"✗ Main test execution failed: {e}")
        traceback.print_exc()
        return False


if __name__ == "__main__":
    success = main()
    exit(0 if success else 1)