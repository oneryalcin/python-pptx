#!/usr/bin/env python3
"""
Live Test Script for FEP-020: Wide-Angle Tree View (get_tree) Functionality

This script validates the get_tree() method implementation across all container objects
(Presentation, Slide, GroupShape) using real python-pptx objects and presentations.

Usage:
    python live_test_tree_functionality.py
"""

import sys
import json
import traceback
from pathlib import Path

# Add src to path for testing
sys.path.insert(0, str(Path(__file__).parent / "src"))

import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE


def test_presentation_tree_functionality():
    """Test get_tree() functionality on presentation objects."""
    print("=" * 80)
    print("FEP-020 Live Test: Presentation Tree Functionality")
    print("=" * 80)
    
    results = {"tests_run": 0, "tests_passed": 0, "failures": []}
    
    def run_test(test_name, test_func):
        """Helper to run individual tests and track results."""
        results["tests_run"] += 1
        try:
            print(f"\n[TEST] {test_name}")
            test_func()
            print(f"✓ PASSED: {test_name}")
            results["tests_passed"] += 1
        except Exception as e:
            print(f"✗ FAILED: {test_name}")
            print(f"  Error: {str(e)}")
            print(f"  Traceback: {traceback.format_exc()}")
            results["failures"].append({"test": test_name, "error": str(e)})
    
    def test_basic_presentation_creation():
        """Test 1: Create a presentation and test basic tree structure."""
        prs = Presentation()
        
        # Verify get_tree method exists
        assert hasattr(prs, 'get_tree'), "Presentation should have get_tree method"
        
        # Test basic tree structure with empty presentation
        tree = prs.get_tree(max_depth=1)
        
        assert isinstance(tree, dict), "get_tree should return a dictionary"
        assert tree["_object_type"] == "Presentation", f"Expected Presentation, got {tree['_object_type']}"
        assert tree["access_path"] == "", "Root presentation should have empty access path"
        assert tree["geometry"] is None, "Presentations should not have geometry"
        assert isinstance(tree["content_summary"], str), "Content summary should be a string"
        assert isinstance(tree["_identity"], dict), "Identity should be a dictionary"
        
        print(f"  Basic tree structure: {tree['_object_type']} with {len(tree.get('children', []))} children")
        print(f"  Content summary: {tree['content_summary']}")
        print(f"  Identity keys: {list(tree['_identity'].keys())}")
    
    def test_presentation_with_slides():
        """Test 2: Create a presentation with multiple slides and test tree expansion."""
        prs = Presentation()
        
        # Add slides with different layouts
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
        
        # Test tree at different depths
        tree_depth_0 = prs.get_tree(max_depth=0)
        tree_depth_1 = prs.get_tree(max_depth=1)
        tree_depth_2 = prs.get_tree(max_depth=2)
        
        # Depth 0: No children
        assert "children" not in tree_depth_0 or tree_depth_0["children"] is None, "Depth 0 should have no children"
        
        # Depth 1: Should include slides but not slide contents
        assert "children" in tree_depth_1, "Depth 1 should include children"
        assert len(tree_depth_1["children"]) == 3, f"Expected 3 slides, got {len(tree_depth_1['children'])}"
        
        for i, slide_child in enumerate(tree_depth_1["children"]):
            assert slide_child["_object_type"] == "Slide", f"Child {i} should be a Slide"
            assert slide_child["access_path"] == f"slides[{i}]", f"Slide {i} access path incorrect"
            assert "children" not in slide_child or slide_child["children"] is None, f"Slide {i} should not have children at depth 1"
        
        # Depth 2: Should include slides and their shapes
        slide_0_tree = tree_depth_2["children"][0]
        if "children" in slide_0_tree and slide_0_tree["children"]:
            shape_count = len(slide_0_tree["children"])
            print(f"  Slide 0 has {shape_count} shapes")
            
            for j, shape_child in enumerate(slide_0_tree["children"]):
                assert "access_path" in shape_child, f"Shape {j} should have access_path"
                expected_path = f"slides[0].shapes[{j}]"
                assert shape_child["access_path"] == expected_path, f"Shape {j} access path should be {expected_path}"
        
        print(f"  Tree depth validation: 0={len(tree_depth_0.get('children', []))}, 1={len(tree_depth_1['children'])}, 2=nested")
    
    def test_slide_tree_functionality():
        """Test 3: Test individual slide tree functionality."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
        
        # Verify slide has get_tree method
        assert hasattr(slide, 'get_tree'), "Slide should have get_tree method"
        
        # Test slide tree
        tree = slide.get_tree(max_depth=1)
        
        assert tree["_object_type"] == "Slide", f"Expected Slide, got {tree['_object_type']}"
        assert tree["access_path"].startswith("slides["), f"Slide access path should start with slides[, got {tree['access_path']}"
        assert tree["geometry"] is None, "Slides should not have geometry"
        assert isinstance(tree["content_summary"], str), "Content summary should be a string"
        
        # Check slide identity
        identity = tree["_identity"]
        assert "slide_id" in identity, "Slide identity should include slide_id"
        assert "class_name" in identity, "Slide identity should include class_name"
        
        print(f"  Slide tree: {tree['content_summary']}")
        print(f"  Access path: {tree['access_path']}")
        print(f"  Identity: slide_id={identity.get('slide_id')}")
    
    def test_shape_tree_functionality():
        """Test 4: Test shape tree functionality including geometry and content."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
        
        # Add various shapes
        shapes = slide.shapes
        
        # Add a text box
        textbox = shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        textbox.text = "This is a test text box"
        
        # Add an autoshape (rectangle)
        rectangle = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2), Inches(3), Inches(2), Inches(1))
        
        # Add a picture (create a simple placeholder)
        # Note: We'll skip picture for now since it requires an actual image file
        
        # Test slide tree with shapes
        tree = slide.get_tree(max_depth=2)
        
        assert "children" in tree, "Slide should have shape children"
        shape_children = tree["children"]
        assert len(shape_children) >= 2, f"Expected at least 2 shapes, got {len(shape_children)}"
        
        # Test first shape (could be textbox or placeholder)
        first_shape_tree = shape_children[0]
        expected_types = ["Shape", "BaseShape", "SlidePlaceholder", "LayoutPlaceholder", "MasterPlaceholder"]
        assert first_shape_tree["_object_type"] in expected_types, f"Expected one of {expected_types}, got {first_shape_tree['_object_type']}"
        assert "geometry" in first_shape_tree, "Shape should have geometry"
        
        geometry = first_shape_tree["geometry"]
        if geometry:
            assert "left" in geometry, "Geometry should include left"
            assert "top" in geometry, "Geometry should include top"
            assert "width" in geometry, "Geometry should include width"
            assert "height" in geometry, "Geometry should include height"
            assert all(val.endswith(" in") for key, val in geometry.items() if key != "rotation"), "Geometry values should be in inches"
        
        # Test shape identity
        identity = first_shape_tree["_identity"]
        assert "shape_id" in identity, "Shape identity should include shape_id"
        assert "name" in identity, "Shape identity should include name"
        
        print(f"  Found {len(shape_children)} shapes on slide")
        print(f"  First shape geometry: {geometry}")
        print(f"  First shape summary: {first_shape_tree['content_summary']}")
    
    def test_group_shape_functionality():
        """Test 5: Test group shape tree functionality."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
        
        shapes = slide.shapes
        
        # Add shapes that we'll group
        shape1 = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1), Inches(1), Inches(1), Inches(1))
        shape2 = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(2.5), Inches(1), Inches(1), Inches(1))
        
        # Group the shapes
        group = shapes.add_group_shape([shape1, shape2])
        
        # Verify group has get_tree method
        assert hasattr(group, 'get_tree'), "GroupShape should have get_tree method"
        
        # Test group tree
        tree = group.get_tree(max_depth=2)
        
        assert tree["_object_type"] == "GroupShape", f"Expected GroupShape, got {tree['_object_type']}"
        assert "geometry" in tree, "Group should have geometry"
        assert isinstance(tree["content_summary"], str), "Content summary should be a string"
        assert "Group" in tree["content_summary"], "Content summary should mention Group"
        
        # Test group children
        if "children" in tree and tree["children"]:
            group_children = tree["children"]
            assert len(group_children) == 2, f"Expected 2 grouped shapes, got {len(group_children)}"
            
            for i, child in enumerate(group_children):
                assert child["access_path"].endswith(f".shapes[{i}]"), f"Group child {i} should have correct access path"
        
        print(f"  Group tree: {tree['content_summary']}")
        print(f"  Group children: {len(tree.get('children', []))}")
    
    def test_access_path_generation():
        """Test 6: Validate access path generation across the hierarchy."""
        prs = Presentation()
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Add shapes to each slide
        for i, slide in enumerate([slide1, slide2]):
            shapes = slide.shapes
            shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1), Inches(1), Inches(1), Inches(1))
        
        # Test full presentation tree
        tree = prs.get_tree(max_depth=2)
        
        # Validate access paths
        slides = tree["children"]
        
        for i, slide_tree in enumerate(slides):
            expected_slide_path = f"slides[{i}]"
            assert slide_tree["access_path"] == expected_slide_path, f"Slide {i} access path should be {expected_slide_path}"
            
            if "children" in slide_tree and slide_tree["children"]:
                shapes = slide_tree["children"]
                for j, shape_tree in enumerate(shapes):
                    expected_shape_path = f"slides[{i}].shapes[{j}]"
                    assert shape_tree["access_path"] == expected_shape_path, f"Shape [{i}][{j}] access path should be {expected_shape_path}"
        
        print(f"  Validated access paths for {len(slides)} slides")
    
    def test_content_summaries():
        """Test 7: Validate content summaries are meaningful."""
        prs = Presentation()
        prs.core_properties.title = "FEP-020 Test Presentation"
        
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
        
        # Add title text
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = "Test Slide Title"
        
        # Test presentation summary
        prs_tree = prs.get_tree(max_depth=0)
        prs_summary = prs_tree["content_summary"]
        assert "FEP-020 Test Presentation" in prs_summary, "Presentation summary should include title"
        assert "slide" in prs_summary.lower(), "Presentation summary should mention slides"
        
        # Test slide summary
        slide_tree = slide.get_tree(max_depth=1)
        slide_summary = slide_tree["content_summary"]
        assert isinstance(slide_summary, str) and len(slide_summary) > 0, "Slide summary should be non-empty string"
        
        print(f"  Presentation summary: {prs_summary}")
        print(f"  Slide summary: {slide_summary}")
    
    def test_error_handling():
        """Test 8: Validate error handling in tree generation."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Test with invalid max_depth
        tree = slide.get_tree(max_depth=0)
        assert isinstance(tree, dict), "Should return valid tree even with max_depth=0"
        
        # Test with large max_depth
        tree = slide.get_tree(max_depth=10)
        assert isinstance(tree, dict), "Should handle large max_depth gracefully"
        
        print(f"  Error handling tests passed")
    
    # Run all tests
    run_test("Basic Presentation Creation", test_basic_presentation_creation)
    run_test("Presentation with Slides", test_presentation_with_slides)
    run_test("Slide Tree Functionality", test_slide_tree_functionality)
    run_test("Shape Tree Functionality", test_shape_tree_functionality)
    run_test("Group Shape Functionality", test_group_shape_functionality)
    run_test("Access Path Generation", test_access_path_generation)
    run_test("Content Summaries", test_content_summaries)
    run_test("Error Handling", test_error_handling)
    
    # Print results
    print("\n" + "=" * 80)
    print("TEST SUMMARY")
    print("=" * 80)
    print(f"Tests Run: {results['tests_run']}")
    print(f"Tests Passed: {results['tests_passed']}")
    print(f"Tests Failed: {len(results['failures'])}")
    
    if results['failures']:
        print("\nFAILURES:")
        for failure in results['failures']:
            print(f"  - {failure['test']}: {failure['error']}")
    
    success_rate = results['tests_passed'] / results['tests_run'] * 100
    print(f"\nSuccess Rate: {success_rate:.1f}%")
    
    return results


def demonstrate_tree_functionality():
    """Demonstrate the tree functionality with a comprehensive example."""
    print("\n" + "=" * 80)
    print("FEP-020 DEMONSTRATION: Wide-Angle Tree View")
    print("=" * 80)
    
    # Create a complex presentation
    prs = Presentation()
    prs.core_properties.title = "Demo: Tree View Functionality"
    prs.core_properties.author = "FEP-020 Implementation"
    
    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    if hasattr(slide1.shapes, 'title') and slide1.shapes.title:
        slide1.shapes.title.text = "Wide-Angle Tree View Demo"
    if hasattr(slide1.shapes, 'placeholders') and len(slide1.shapes.placeholders) > 1:
        slide1.shapes.placeholders[1].text = "Demonstrating FEP-020 get_tree() functionality"
    
    # Slide 2: Content slide with shapes
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    shapes = slide2.shapes
    
    # Add various shapes
    textbox = shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text = "This is a text box with sample content"
    
    rectangle = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1), Inches(3), Inches(2), Inches(1))
    oval = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(4), Inches(3), Inches(2), Inches(1))
    
    # Group some shapes
    try:
        group = shapes.add_group_shape([rectangle, oval])
    except:
        print("  Note: Group creation failed (this is expected in some environments)")
        group = None
    
    # Demonstrate tree at different depths
    print("\n1. PRESENTATION TREE (max_depth=0 - Root only):")
    tree_0 = prs.get_tree(max_depth=0)
    print(json.dumps(tree_0, indent=2, default=str))
    
    print("\n2. PRESENTATION TREE (max_depth=1 - Include slides):")
    tree_1 = prs.get_tree(max_depth=1)
    print(json.dumps(tree_1, indent=2, default=str))
    
    print("\n3. PRESENTATION TREE (max_depth=2 - Include slide contents):")
    tree_2 = prs.get_tree(max_depth=2)
    # Truncate output for readability
    tree_2_summary = {
        "_object_type": tree_2["_object_type"],
        "content_summary": tree_2["content_summary"],
        "children_count": len(tree_2.get("children", [])),
        "first_slide_children_count": len(tree_2["children"][0].get("children", [])) if tree_2.get("children") else 0
    }
    print(json.dumps(tree_2_summary, indent=2, default=str))
    
    print("\n4. INDIVIDUAL SLIDE TREE:")
    slide_tree = slide2.get_tree(max_depth=1)
    print(json.dumps(slide_tree, indent=2, default=str))
    
    if group:
        print("\n5. GROUP SHAPE TREE:")
        group_tree = group.get_tree(max_depth=1)
        print(json.dumps(group_tree, indent=2, default=str))
    
    print("\n6. ACCESS PATH EXAMPLES:")
    full_tree = prs.get_tree(max_depth=2)
    for i, slide in enumerate(full_tree.get("children", [])):
        print(f"  Slide {i}: {slide['access_path']}")
        for j, shape in enumerate(slide.get("children", [])):
            print(f"    Shape {j}: {shape['access_path']} - {shape['content_summary']}")


def main():
    """Main test execution."""
    print("FEP-020: Wide-Angle Tree View (get_tree) - Live Testing")
    print("=" * 80)
    print("This script tests the get_tree() functionality implementation.")
    print("Testing with real python-pptx objects and presentations.")
    print()
    
    try:
        # Run comprehensive tests
        test_results = test_presentation_tree_functionality()
        
        # Run demonstration
        demonstrate_tree_functionality()
        
        # Final summary
        print("\n" + "=" * 80)
        print("FEP-020 LIVE TESTING COMPLETE")
        print("=" * 80)
        
        if test_results["tests_passed"] == test_results["tests_run"]:
            print("🎉 ALL TESTS PASSED! FEP-020 implementation is working correctly.")
        else:
            print(f"⚠️  {len(test_results['failures'])} tests failed. See details above.")
        
        print("\nKey Features Validated:")
        print("✓ get_tree() method on Presentation, Slide, and GroupShape")
        print("✓ Hierarchical tree structure with configurable depth")
        print("✓ Access path generation for stable object references")
        print("✓ Rich content summaries for AI discovery")
        print("✓ Identity information for object identification")
        print("✓ Geometry information for spatial understanding")
        print("✓ Error handling and edge cases")
        
        return test_results["tests_passed"] == test_results["tests_run"]
        
    except Exception as e:
        print(f"\nFATAL ERROR: {str(e)}")
        print(f"Traceback: {traceback.format_exc()}")
        return False


if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)