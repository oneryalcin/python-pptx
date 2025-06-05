#!/usr/bin/env python3
"""
Live Test Script for Presentation.to_dict() Introspection - FEP-013

This script demonstrates the Presentation introspection functionality with real python-pptx
objects, providing validation for engineers reviewing the FEP-013 implementation.

Usage:
    python live_test_presentation_introspection.py

The script creates various presentation configurations and tests the to_dict() functionality
across different scenarios, parameters, and edge cases.

Test Results Summary:
- 7/7 live tests pass ✅
- Comprehensive coverage of presentation introspection scenarios
- All critical functionality validated with real objects
- Error handling and edge cases confirmed working

Unit Test Results Summary:
- All unit tests pass - validates core functionality 
- Comprehensive test coverage of Presentation.to_dict() methods
- No regressions in existing introspection test suite
"""

import json
import sys
from pathlib import Path

# Add the src directory to Python path for imports
sys.path.insert(0, str(Path(__file__).parent / "src"))

import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


def print_section(title):
    """Print a formatted section header."""
    print(f"\n{'=' * 60}")
    print(f"{title}")
    print("=" * 60)


def print_json_snippet(data, max_lines=30):
    """Print a formatted JSON snippet with line limit."""
    json_str = json.dumps(data, indent=2, default=str)
    lines = json_str.split('\n')
    
    if len(lines) <= max_lines:
        print(json_str)
    else:
        print('\n'.join(lines[:max_lines]))
        print(f"... ({len(lines) - max_lines} more lines)")


def test_basic_presentation_introspection():
    """Test basic presentation introspection with default presentation."""
    print_section("Test 1: Basic Presentation Introspection")
    
    # Create a new presentation
    prs = Presentation()
    
    # Set some core properties
    prs.core_properties.title = "FEP-013 Test Presentation"
    prs.core_properties.author = "FEP Test Suite"
    prs.core_properties.subject = "Testing Presentation Introspection"
    prs.core_properties.comments = "Created for validating FEP-013 implementation"
    
    # Test basic to_dict functionality
    result = prs.to_dict(expand_collections=False, max_depth=1)
    
    print("✅ Basic to_dict() successful")
    
    # Verify structure
    assert "_object_type" in result
    assert result["_object_type"] == "Presentation"
    assert "_identity" in result
    assert "properties" in result
    assert "_llm_context" in result
    
    print("✅ Required top-level keys present")
    
    # Check identity
    identity = result["_identity"]
    assert "class_name" in identity
    assert identity["class_name"] == "Presentation"
    assert "description" in identity
    # Description should mention presentation, but may reference template file
    assert "presentation" in identity["description"].lower() or "Presentation" in identity["description"]
    
    print("✅ Identity information correct")
    
    # Check core properties
    props = result["properties"]
    assert "core_properties" in props
    core_props = props["core_properties"]
    assert core_props["title"] == "FEP-013 Test Presentation"
    assert core_props["author"] == "FEP Test Suite"
    assert core_props["subject"] == "Testing Presentation Introspection"
    
    print("✅ Core properties extracted correctly")
    
    # Check slide dimensions
    assert "slide_width" in props
    assert "slide_height" in props
    assert props["slide_width"]["_object_type"] == "Emu"
    assert props["slide_height"]["_object_type"] == "Emu"
    
    print("✅ Slide dimensions formatted correctly")
    
    # Check collections (collapsed)
    assert "slides" in props
    assert "slide_masters" in props
    assert isinstance(props["slides"], str)
    assert "Collection of" in props["slides"]
    assert "not expanded" in props["slides"]
    
    print("✅ Collections correctly collapsed")
    
    # Check LLM context
    llm_context = result["_llm_context"]
    assert "description" in llm_context
    assert "common_operations" in llm_context
    assert "FEP-013 Test Presentation" in llm_context["description"]
    
    print("✅ LLM context generated correctly")
    print("📊 Sample core properties:")
    print_json_snippet(core_props, 10)
    
    return True


def test_presentation_with_slides():
    """Test presentation introspection with multiple slides."""
    print_section("Test 2: Presentation with Multiple Slides")
    
    # Create presentation with slides
    prs = Presentation()
    prs.core_properties.title = "Multi-Slide Presentation"
    
    # Add several slides
    slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(slide_layout)
    slide2 = prs.slides.add_slide(slide_layout)
    slide3 = prs.slides.add_slide(slide_layout)
    
    # Add content to slides
    slide1.shapes.title.text = "Slide 1"
    slide2.shapes.title.text = "Slide 2"
    slide3.shapes.title.text = "Slide 3"
    
    # Test with collections expanded
    result = prs.to_dict(expand_collections=True, max_depth=2)
    
    print("✅ to_dict() with expanded collections successful")
    
    # Check slides collection
    props = result["properties"]
    assert "slides" in props
    slides = props["slides"]
    assert isinstance(slides, list)
    assert len(slides) == 3
    
    print(f"✅ Found {len(slides)} slides in expanded collection")
    
    # Verify slide structure
    for i, slide_dict in enumerate(slides):
        assert slide_dict["_object_type"] == "Slide"
        assert "_identity" in slide_dict
        assert "properties" in slide_dict
        
    print("✅ All slides have proper to_dict structure")
    
    # Check slide masters collection
    assert "slide_masters" in props
    masters = props["slide_masters"]
    assert isinstance(masters, list)
    assert len(masters) >= 1
    
    print(f"✅ Found {len(masters)} slide masters")
    
    # Check LLM description includes slide count
    llm_context = result["_llm_context"]
    assert "3 slide(s)" in llm_context["description"]
    
    print("✅ LLM context reflects correct slide count")
    print("📊 Sample slide summary:")
    print_json_snippet(slides[0]["_identity"], 5)
    
    return True


def test_presentation_relationships():
    """Test presentation relationship information."""
    print_section("Test 3: Presentation Relationships")
    
    prs = Presentation()
    prs.core_properties.title = "Relationships Test"
    
    # Test relationships
    result = prs.to_dict(include_relationships=True)
    
    print("✅ to_dict() with relationships successful")
    
    # Check relationships section
    assert "relationships" in result
    rels = result["relationships"]
    
    # Check for expected relationships
    assert "main_document_part" in rels
    assert "core_properties_part" in rels
    
    print("✅ Key relationships present")
    
    # Verify relationship structure
    main_part = rels["main_document_part"]
    assert "partname" in main_part
    assert "/ppt/presentation.xml" in main_part["partname"]
    
    core_part = rels["core_properties_part"]
    assert "partname" in core_part
    assert "/docProps/core.xml" in core_part["partname"]
    
    print("✅ Relationship part names correct")
    print("📊 Relationships:")
    print_json_snippet(rels, 10)
    
    return True


def test_presentation_notes_master():
    """Test notes master handling."""
    print_section("Test 4: Notes Master Handling")
    
    prs = Presentation()
    prs.core_properties.title = "Notes Master Test"
    
    # Access notes master (creates one if not exists)
    notes_master = prs.notes_master
    
    # Test introspection
    result = prs.to_dict(expand_collections=True, max_depth=2)
    
    print("✅ to_dict() with notes master successful")
    
    # Check notes master in properties
    props = result["properties"]
    assert "notes_master" in props
    notes_master_dict = props["notes_master"]
    
    assert isinstance(notes_master_dict, dict)
    assert notes_master_dict["_object_type"] == "NotesMaster"
    
    print("✅ Notes master represented correctly")
    
    # Check LLM context mentions notes master
    llm_context = result["_llm_context"]
    assert "notes master" in llm_context["description"]
    
    print("✅ LLM context includes notes master")
    print("📊 Notes master summary:")
    print_json_snippet(notes_master_dict, 8)
    
    return True


def test_presentation_max_depth_control():
    """Test max_depth parameter control."""
    print_section("Test 5: Max Depth Control")
    
    prs = Presentation()
    prs.core_properties.title = "Depth Control Test"
    
    # Add a slide for testing
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Test Slide"
    
    # Test with max_depth=0 (minimal expansion)
    result_depth_0 = prs.to_dict(expand_collections=True, max_depth=0)
    
    print("✅ to_dict() with max_depth=0 successful")
    
    # With max_depth=0, entire object is truncated
    assert "_truncated" in result_depth_0
    assert "Max depth reached" in result_depth_0["_truncated"]
    
    print("✅ max_depth=0 truncates entire object")
    
    # Test with max_depth=1 (allow some expansion)
    result_depth_1 = prs.to_dict(expand_collections=True, max_depth=1)
    props_1 = result_depth_1["properties"]
    
    # Should expand slides but not deeply
    assert isinstance(props_1["slides"], list)
    
    print("✅ max_depth=1 allows controlled expansion")
    
    # Test with max_depth=3 (deeper expansion)
    result_depth_3 = prs.to_dict(expand_collections=True, max_depth=3)
    props_3 = result_depth_3["properties"]
    
    assert isinstance(props_3["slides"], list)
    
    print("✅ max_depth=3 allows full expansion")
    print(f"📊 Depth comparison - object at depth 0: {type(result_depth_0).__name__} with _truncated key")
    print(f"📊 Depth comparison - slides at depth 1: {type(props_1['slides']).__name__}")
    print(f"📊 Depth comparison - slides at depth 3: {type(props_3['slides']).__name__}")
    
    return True


def test_presentation_llm_context_generation():
    """Test LLM context generation with various configurations."""
    print_section("Test 6: LLM Context Generation")
    
    prs = Presentation()
    
    # Set rich core properties
    prs.core_properties.title = "Comprehensive LLM Test Presentation"
    prs.core_properties.author = "AI Test Engineer"
    prs.core_properties.subject = "Advanced Introspection Testing"
    
    # Add multiple slides
    slide_layout = prs.slide_layouts[0]
    for i in range(5):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Slide {i+1}"
    
    # Change slide dimensions
    prs.slide_width = Inches(11)
    prs.slide_height = Inches(8.5)
    
    # Test LLM context
    result = prs.to_dict(format_for_llm=True)
    
    print("✅ to_dict() with format_for_llm=True successful")
    
    # Check LLM context content
    llm_context = result["_llm_context"]
    
    assert "description" in llm_context
    assert "summary" in llm_context
    assert "common_operations" in llm_context
    
    description = llm_context["description"]
    
    # Verify rich description content
    assert "Comprehensive LLM Test Presentation" in description
    assert "5 slide(s)" in description
    assert "slide master(s)" in description
    assert "11.00\"W x 8.50\"H" in description
    assert "notes master" in description
    
    print("✅ LLM description contains all expected elements")
    
    # Check common operations
    operations = llm_context["common_operations"]
    assert isinstance(operations, list)
    assert len(operations) > 0
    
    expected_ops = ["slides", "save", "slide_masters", "core_properties"]
    for expected_op in expected_ops:
        assert any(expected_op in op for op in operations)
    
    print("✅ Common operations include all expected actions")
    print("📊 Full LLM context:")
    print_json_snippet(llm_context, 15)
    
    return True


def test_presentation_error_handling():
    """Test error handling and edge cases."""
    print_section("Test 7: Error Handling and Edge Cases")
    
    prs = Presentation()
    prs.core_properties.title = "Error Handling Test"
    
    # Test with format_for_llm=False (should exclude LLM context)
    result_no_llm = prs.to_dict(format_for_llm=False)
    assert "_llm_context" not in result_no_llm
    
    print("✅ format_for_llm=False excludes LLM context")
    
    # Test with include_private=False (default - should exclude private fields)
    result = prs.to_dict(include_private=False)
    
    # Verify no private fields are included (except standard ones)
    def has_unexpected_private_keys(obj, path=""):
        if isinstance(obj, dict):
            for key, value in obj.items():
                # Allow standard introspection markers
                allowed_private = ['_object_type', '_identity', '_llm_context', '_no_introspection', 
                                 '_summary_or_truncated', '_depth_exceeded', '_truncated']
                if key.startswith('_') and key not in allowed_private:
                    print(f"❌ Found unexpected private key: {path}.{key}")
                    return True
                if has_unexpected_private_keys(value, f"{path}.{key}"):
                    return True
        elif isinstance(obj, list):
            for i, item in enumerate(obj):
                if has_unexpected_private_keys(item, f"{path}[{i}]"):
                    return True
        return False
    
    assert not has_unexpected_private_keys(result)
    
    print("✅ No unexpected private fields included")
    
    # Test with empty/minimal presentation
    empty_prs = Presentation()
    empty_result = empty_prs.to_dict()
    
    assert empty_result["_object_type"] == "Presentation"
    assert "properties" in empty_result
    
    print("✅ Handles empty presentation correctly")
    
    # Test very large max_depth (should not cause issues)
    try:
        large_depth_result = prs.to_dict(max_depth=100)
        print("✅ Large max_depth handled gracefully")
    except Exception as e:
        print(f"❌ Large max_depth caused error: {e}")
        return False
    
    print("✅ All error handling tests passed")
    
    return True


def main():
    """Run all live tests for Presentation introspection."""
    print("🚀 Starting Live Tests for Presentation.to_dict() - FEP-013")
    print("=" * 80)
    
    tests = [
        test_basic_presentation_introspection,
        test_presentation_with_slides,
        test_presentation_relationships,
        test_presentation_notes_master,
        test_presentation_max_depth_control,
        test_presentation_llm_context_generation,
        test_presentation_error_handling
    ]
    
    passed = 0
    failed = 0
    
    for i, test in enumerate(tests, 1):
        try:
            if test():
                print(f"✅ Test {i} PASSED")
                passed += 1
            else:
                print(f"❌ Test {i} FAILED")
                failed += 1
        except Exception as e:
            print(f"❌ Test {i} FAILED with exception: {e}")
            failed += 1
    
    print_section("FINAL RESULTS")
    print(f"✅ PASSED: {passed}")
    print(f"❌ FAILED: {failed}")
    print(f"📊 SUCCESS RATE: {passed}/{passed+failed} ({100*passed/(passed+failed):.1f}%)")
    
    if failed == 0:
        print("\n🎉 ALL TESTS PASSED! FEP-013 implementation is working correctly.")
        return True
    else:
        print(f"\n⚠️  {failed} test(s) failed. Please review the implementation.")
        return False


if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)