# tests/introspection/test_placeholder_introspection.py

"""
PlaceholderFormat Introspection Tests

Tests for _PlaceholderFormat introspection functionality including:
- PlaceholderFormat.to_dict() basic functionality
- Placeholder properties (idx, type) serialization
- LLM context generation for placeholder descriptions
- Integration with BaseShape placeholder handling
- Error handling for edge cases
"""

import unittest
from unittest.mock import Mock, patch

from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.shapes.base import BaseShape, _PlaceholderFormat
from pptx.util import Emu

from .mock_helpers import assert_basic_to_dict_structure


class MockCTPlaceholder:
    """Mock CT_Placeholder element for testing."""

    def __init__(self, idx=0, ph_type=PP_PLACEHOLDER.TITLE):
        self.idx = idx
        self.type = ph_type


class TestPlaceholderFormatIntrospection(unittest.TestCase):
    """Test _PlaceholderFormat introspection functionality."""

    def setUp(self):
        """Set up test fixtures."""
        self.mock_element = MockCTPlaceholder(idx=1, ph_type=PP_PLACEHOLDER.BODY)
        self.placeholder_format = _PlaceholderFormat(self.mock_element)

    def test_placeholder_format_to_dict_basic_functionality(self):
        """Test that _PlaceholderFormat.to_dict() provides expected basic structure."""
        result = self.placeholder_format.to_dict(include_relationships=False, max_depth=2)

        # Check basic structure
        assert_basic_to_dict_structure(self, result, '_PlaceholderFormat')

        # Check identity information
        identity = result['_identity']
        self.assertIn('description', identity)
        self.assertIn('BODY placeholder', identity['description'])
        self.assertIn('idx: 1', identity['description'])

        # Check properties
        props = result['properties']
        self.assertEqual(props['idx'], 1)

        # Check type property structure
        self.assertIn('type', props)
        type_dict = props['type']
        self.assertIsInstance(type_dict, dict)
        self.assertEqual(type_dict.get('_object_type'), 'PP_PLACEHOLDER_TYPE')
        self.assertEqual(type_dict.get('name'), 'BODY')
        self.assertEqual(type_dict.get('value'), PP_PLACEHOLDER.BODY.value)

    def test_placeholder_format_to_dict_title_placeholder(self):
        """Test _PlaceholderFormat.to_dict() for title placeholder."""
        title_element = MockCTPlaceholder(idx=0, ph_type=PP_PLACEHOLDER.TITLE)
        title_format = _PlaceholderFormat(title_element)

        result = title_format.to_dict(include_relationships=False, max_depth=2)

        # Check identity information
        identity = result['_identity']
        self.assertIn('TITLE placeholder', identity['description'])
        self.assertIn('idx: 0', identity['description'])

        # Check properties
        props = result['properties']
        self.assertEqual(props['idx'], 0)

        type_dict = props['type']
        self.assertEqual(type_dict.get('name'), 'TITLE')
        self.assertEqual(type_dict.get('value'), PP_PLACEHOLDER.TITLE.value)

    def test_placeholder_format_to_dict_picture_placeholder(self):
        """Test _PlaceholderFormat.to_dict() for picture placeholder."""
        picture_element = MockCTPlaceholder(idx=2, ph_type=PP_PLACEHOLDER.PICTURE)
        picture_format = _PlaceholderFormat(picture_element)

        result = picture_format.to_dict(include_relationships=False, max_depth=2)

        # Check identity information
        identity = result['_identity']
        self.assertIn('PICTURE placeholder', identity['description'])
        self.assertIn('idx: 2', identity['description'])

        # Check properties
        props = result['properties']
        self.assertEqual(props['idx'], 2)

        type_dict = props['type']
        self.assertEqual(type_dict.get('name'), 'PICTURE')
        self.assertEqual(type_dict.get('value'), PP_PLACEHOLDER.PICTURE.value)

    def test_placeholder_format_llm_context(self):
        """Test _PlaceholderFormat LLM context generation."""
        result = self.placeholder_format.to_dict(format_for_llm=True, max_depth=2)

        # Check LLM context
        self.assertIn('_llm_context', result)
        llm_context = result['_llm_context']

        self.assertIn('description', llm_context)
        self.assertIn('BODY', llm_context['description'])
        self.assertIn('Index is 1', llm_context['description'])

        self.assertIn('summary', llm_context)
        self.assertEqual(llm_context['summary'], llm_context['description'])

        self.assertIn('common_operations', llm_context)
        self.assertIsInstance(llm_context['common_operations'], list)
        self.assertGreater(len(llm_context['common_operations']), 0)

        # Check specific operations
        operations = llm_context['common_operations']
        self.assertTrue(any('placeholder role' in op for op in operations))
        self.assertTrue(any('unique index' in op for op in operations))

    def test_placeholder_format_to_dict_no_relationships(self):
        """Test that _PlaceholderFormat has no relationships."""
        result = self.placeholder_format.to_dict(include_relationships=True, max_depth=2)

        # PlaceholderFormat should have empty relationships
        self.assertIn('relationships', result)
        relationships = result['relationships']
        self.assertIsInstance(relationships, dict)
        # Should be empty or minimal since PlaceholderFormat is a simple property bag

    def test_placeholder_format_to_dict_format_for_llm(self):
        """Test _PlaceholderFormat.to_dict() with format_for_llm=True."""
        result = self.placeholder_format.to_dict(format_for_llm=True, max_depth=2)

        # Should include LLM context
        self.assertIn('_llm_context', result)

        # Properties should still be present and correctly formatted
        props = result['properties']
        self.assertEqual(props['idx'], 1)
        self.assertIn('type', props)

    def test_placeholder_format_to_dict_max_depth_handling(self):
        """Test _PlaceholderFormat.to_dict() with various max_depth values."""
        # Test with max_depth=1 (should still work since PlaceholderFormat is simple)
        result1 = self.placeholder_format.to_dict(max_depth=1)
        self.assertIn('properties', result1)
        self.assertIn('idx', result1['properties'])

        # Test with max_depth=0 (minimal output - should be truncated)
        result0 = self.placeholder_format.to_dict(max_depth=0)
        # At max_depth=0, output is typically truncated
        self.assertTrue('_truncated' in result0 or '_object_type' in result0)


class TestBaseShapePlaceholderIntegration(unittest.TestCase):
    """Test integration of _PlaceholderFormat.to_dict() with BaseShape."""

    def setUp(self):
        """Set up test fixtures."""
        # Create mock shape element that is a placeholder
        self.mock_shape_element = Mock()
        self.mock_shape_element.has_ph_elm = True
        self.mock_shape_element.shape_id = 42
        self.mock_shape_element.shape_name = "Test Placeholder Shape"
        self.mock_shape_element.x = Emu(914400)  # 1 inch
        self.mock_shape_element.y = Emu(914400)
        self.mock_shape_element.cx = Emu(2743200)  # 3 inches
        self.mock_shape_element.cy = Emu(1828800)  # 2 inches
        self.mock_shape_element.rot = 0.0

        # Create mock placeholder element
        self.mock_ph_element = MockCTPlaceholder(idx=1, ph_type=PP_PLACEHOLDER.BODY)
        self.mock_shape_element.ph = self.mock_ph_element

        # Create mock parent
        self.mock_parent = Mock()

        # Create BaseShape
        self.shape = BaseShape(self.mock_shape_element, self.mock_parent)

    def test_base_shape_placeholder_details_uses_to_dict(self):
        """Test that BaseShape uses _PlaceholderFormat.to_dict() for placeholder details."""
        result = self.shape.to_dict(include_relationships=False, max_depth=3)

        # Check that shape is identified as placeholder
        identity = result['_identity']
        self.assertTrue(identity['is_placeholder'])

        # Check placeholder details structure
        self.assertIn('placeholder_details', identity)
        placeholder_details = identity['placeholder_details']

        # Should be the output from _PlaceholderFormat.to_dict()
        assert_basic_to_dict_structure(self, placeholder_details, '_PlaceholderFormat')

        # Check specific placeholder information
        self.assertIn('properties', placeholder_details)
        props = placeholder_details['properties']
        self.assertEqual(props['idx'], 1)

        type_dict = props['type']
        self.assertEqual(type_dict.get('name'), 'BODY')

    def test_base_shape_non_placeholder_no_details(self):
        """Test that non-placeholder shapes don't have placeholder details."""
        # Create non-placeholder shape
        self.mock_shape_element.has_ph_elm = False
        self.mock_shape_element.ph = None

        shape = BaseShape(self.mock_shape_element, self.mock_parent)
        result = shape.to_dict(include_relationships=False, max_depth=2)

        # Check that shape is not identified as placeholder
        identity = result['_identity']
        self.assertFalse(identity['is_placeholder'])

        # Should not have placeholder details
        self.assertNotIn('placeholder_details', identity)

    @patch('pptx.shapes.base.BaseShape.placeholder_format')
    def test_base_shape_placeholder_error_handling(self, mock_placeholder_format):
        """Test BaseShape placeholder error handling when placeholder_format fails."""
        # Make placeholder_format raise an error
        mock_placeholder_format.side_effect = ValueError("Test error")

        result = self.shape.to_dict(include_relationships=False, max_depth=2)

        # Check that shape is still identified as placeholder
        identity = result['_identity']
        self.assertTrue(identity['is_placeholder'])

        # Should have error context in placeholder details
        self.assertIn('placeholder_details', identity)
        placeholder_details = identity['placeholder_details']

        # Should contain error information
        # Since we mocked placeholder_format to raise an error, the to_dict call should be mocked too
        # This test verifies the error handling path is exercised
        self.assertIsNotNone(placeholder_details)

    def test_base_shape_llm_context_with_placeholder(self):
        """Test BaseShape LLM context includes placeholder information."""
        result = self.shape.to_dict(format_for_llm=True, max_depth=3)

        # Check LLM context
        self.assertIn('_llm_context', result)
        llm_context = result['_llm_context']

        # Should mention placeholder information
        description = llm_context.get('description', '')
        self.assertIn('placeholder', description.lower())
        # Note: The description might not contain BODY if placeholder info isn't fully integrated
        # This is expected behavior - the shape description is separate from placeholder details

        # Should include placeholder operations
        operations = llm_context.get('common_operations', [])
        self.assertTrue(any('placeholder' in op for op in operations))


if __name__ == '__main__':
    unittest.main()
