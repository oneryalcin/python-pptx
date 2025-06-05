# tests/introspection/mock_helpers.py

"""
Shared Mock Classes and Testing Utilities for Introspection Tests

This module provides common mock classes, helper utilities, and test patterns
used across the introspection test suite. Centralizing these reduces duplication
and ensures consistent testing patterns.
"""

from pptx.dml.color import ColorFormat, RGBColor
from pptx.dml.fill import FillFormat, _GradientStop
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL, MSO_LINE_DASH_STYLE, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.introspection import IntrospectionMixin
from pptx.util import Emu, Pt

# =============================================================================
# Core Testing Mock Classes
# =============================================================================

class MyObjectWithRGB(IntrospectionMixin):
    """Test object with RGB color and length properties for basic formatting tests."""

    def __init__(self, rgb_color_val, length_val=None):
        self.rgb_color_attr = rgb_color_val
        self.length_attr = length_val
        self._private_rgb = RGBColor(0, 0, 0)  # for testing private

    def _to_dict_properties(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        props = {
            "rgb_color_attr": self._format_property_value_for_to_dict(
                self.rgb_color_attr, include_private, _visited_ids, max_depth - 1, expand_collections, format_for_llm
            )
        }
        if self.length_attr is not None:
            props["length_attr"] = self._format_property_value_for_to_dict(
                self.length_attr, include_private, _visited_ids, max_depth - 1, expand_collections, format_for_llm
            )

        if include_private:
            props["_private_rgb"] = self._format_property_value_for_to_dict(
                self._private_rgb, include_private, _visited_ids, max_depth - 1, expand_collections, format_for_llm
            )
        return props


class DummyInspectable(IntrospectionMixin):
    """Generic inspectable object for testing collections, private fields, and basic structure."""

    def __init__(self):
        self.name = "Dummy"
        self.value = 123
        self._private_field = "secret"
        self.my_list = [1, RGBColor(10, 20, 30)]
        self.my_dict = {"key": "val", "color_key": RGBColor(40, 50, 60)}

    def _to_dict_properties(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        props = {
            "name": self._format_property_value_for_to_dict(
                self.name, include_private, _visited_ids, max_depth - 1, expand_collections, format_for_llm
            ),
            "value": self._format_property_value_for_to_dict(
                self.value, include_private, _visited_ids, max_depth - 1, expand_collections, format_for_llm
            ),
            "my_list": self._format_property_value_for_to_dict(
                self.my_list, include_private, _visited_ids, max_depth - 1, expand_collections, format_for_llm
            ),
            "my_dict": self._format_property_value_for_to_dict(
                self.my_dict, include_private, _visited_ids, max_depth - 1, expand_collections, format_for_llm
            )
        }
        if include_private:
            props["_private_field"] = self._format_property_value_for_to_dict(
                self._private_field, include_private, _visited_ids, max_depth - 1, expand_collections, format_for_llm
            )
        return props


# =============================================================================
# Depth Testing Mock Classes
# =============================================================================

class DepthTestA(IntrospectionMixin):
    """Test object for depth limiting - top level."""

    def __init__(self, b_instance=None):
        self.b_prop = b_instance
        self.name = "A"

    def _to_dict_properties(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        return {
            "name": self.name,
            "b_prop": self._format_property_value_for_to_dict(
                self.b_prop, include_private, _visited_ids, max_depth - 1, expand_collections, format_for_llm
            )
        }


class DepthTestB(IntrospectionMixin):
    """Test object for depth limiting - middle level."""

    def __init__(self, c_instance=None):
        self.c_prop = c_instance
        self.name = "B"

    def _to_dict_properties(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        return {
            "name": self.name,
            "c_prop": self._format_property_value_for_to_dict(
                self.c_prop, include_private, _visited_ids, max_depth - 1, expand_collections, format_for_llm
            )
        }


class DepthTestC(IntrospectionMixin):
    """Test object for depth limiting - bottom level."""

    def __init__(self):
        self.name = "C"

    def _to_dict_properties(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        return {"name": self.name}


# =============================================================================
# Circular Reference Testing Mock Classes
# =============================================================================

class CycleA(IntrospectionMixin):
    """Test object for circular reference detection - A side."""

    def __init__(self):
        self.name = "CycleA"
        self.b_ref = None

    def _to_dict_properties(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        return {
            "name": self.name,
            "b_ref": self._format_property_value_for_to_dict(
                self.b_ref, include_private, _visited_ids, max_depth, expand_collections, format_for_llm
            )
        }


class CycleB(IntrospectionMixin):
    """Test object for circular reference detection - B side."""

    def __init__(self):
        self.name = "CycleB"
        self.a_ref = None

    def _to_dict_properties(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        return {
            "name": self.name,
            "a_ref": self._format_property_value_for_to_dict(
                self.a_ref, include_private, _visited_ids, max_depth, expand_collections, format_for_llm
            )
        }


# =============================================================================
# Enum Testing Mock Classes
# =============================================================================

class EnumTestObj(IntrospectionMixin):
    """Test object with BaseEnum property."""

    def __init__(self, enum_val):
        self.shape_type = enum_val

    def _to_dict_properties(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        return {
            "shape_type": self._format_property_value_for_to_dict(
                self.shape_type, include_private, _visited_ids, max_depth - 1, expand_collections, format_for_llm
            )
        }


class XmlEnumTestObj(IntrospectionMixin):
    """Test object with BaseXmlEnum property."""

    def __init__(self, enum_val):
        self.auto_shape_type = enum_val

    def _to_dict_properties(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        return {
            "auto_shape_type": self._format_property_value_for_to_dict(
                self.auto_shape_type, include_private, _visited_ids, max_depth - 1, expand_collections, format_for_llm
            )
        }


class EnumCollectionTestObj(IntrospectionMixin):
    """Test object with collections containing enum members."""

    def __init__(self):
        self.shape_types = [MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.TABLE]
        self.mixed_collection = [MSO_COLOR_TYPE.RGB, "string", 42]

    def _to_dict_properties(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        return {
            "shape_types": self._format_property_value_for_to_dict(
                self.shape_types, include_private, _visited_ids, max_depth - 1, expand_collections, format_for_llm
            ),
            "mixed_collection": self._format_property_value_for_to_dict(
                self.mixed_collection, include_private, _visited_ids, max_depth - 1, expand_collections, format_for_llm
            )
        }


# =============================================================================
# Shape Testing Mock Classes
# =============================================================================

class MockShapeElement:
    """Mock shape element for BaseShape testing."""

    def __init__(self, shape_id=42, name="Test Shape", is_placeholder=False, **kwargs):
        self.shape_id = shape_id
        self.shape_name = name
        self.has_ph_elm = is_placeholder
        self.x = kwargs.get('x', Emu(914400))  # 1 inch
        self.y = kwargs.get('y', Emu(914400))  # 1 inch
        self.cx = kwargs.get('cx', Emu(1828800))  # 2 inches
        self.cy = kwargs.get('cy', Emu(914400))  # 1 inch
        self.rot = kwargs.get('rot', 0.0)
        self.hidden = kwargs.get('hidden', False)


class MockParent:
    """Mock parent collection for shape testing."""

    def __init__(self, description="MockSlideShapes"):
        self.description = description

    def __repr__(self):
        return f"<{self.description} object at 0x123>"


class MockPlaceholderFormat:
    """Mock placeholder format for placeholder testing."""

    def __init__(self, idx=0, ph_type=PP_PLACEHOLDER.TITLE):
        self.idx = idx
        self.type = ph_type

    def to_dict(self, include_relationships=True, max_depth=3, include_private=False, expand_collections=True, format_for_llm=True, _visited_ids=None):
        """Mock to_dict method that returns structure similar to real _PlaceholderFormat.to_dict()."""
        return {
            "_object_type": "_PlaceholderFormat",
            "_identity": {
                "class_name": "_PlaceholderFormat",
                "description": f"Details for a {self.type.name} placeholder (idx: {self.idx})."
            },
            "properties": {
                "idx": self.idx,
                "type": {
                    "_object_type": "PP_PLACEHOLDER_TYPE",
                    "name": self.type.name,
                    "value": self.type.value,
                    "xml_value": getattr(self.type, 'xml_value', self.type.value)
                }
            },
            "relationships": {},
            "_llm_context": {
                "description": f"Placeholder attributes: Type is {self.type.name}, Index is {self.idx}.",
                "summary": f"Placeholder attributes: Type is {self.type.name}, Index is {self.idx}.",
                "common_operations": [
                    "identify placeholder role (e.g., TITLE, BODY, PICTURE)",
                    "get unique index (idx) for matching with layout/master"
                ]
            }
        }


# =============================================================================
# Color Testing Mock Classes
# =============================================================================

class MockColorFormat(ColorFormat):
    """Flexible mock ColorFormat for testing different color scenarios."""

    def __init__(self, color_type=MSO_COLOR_TYPE.RGB, rgb_color=None, theme_color=None, brightness=0.0):
        # Skip parent initialization for testing
        self._color_type = color_type
        self._rgb_color = rgb_color or RGBColor(0x12, 0x34, 0x56)
        self._theme_color = theme_color or MSO_THEME_COLOR.ACCENT_1
        self._brightness_val = brightness

    @property
    def type(self):
        return self._color_type

    @property
    def rgb(self):
        if self._color_type == MSO_COLOR_TYPE.RGB:
            return self._rgb_color
        raise AttributeError("no .rgb property on color type")

    @property
    def theme_color(self):
        if self._color_type == MSO_COLOR_TYPE.SCHEME:
            return self._theme_color
        raise AttributeError("no .theme_color property on color type")

    @property
    def brightness(self):
        if self._color_type is None:
            raise ValueError("can't access brightness when color.type is None")
        return self._brightness_val


# =============================================================================
# Fill Testing Mock Classes
# =============================================================================

class MockFillFormat(FillFormat, IntrospectionMixin):
    """Flexible mock FillFormat for testing different fill scenarios."""

    def __init__(self, fill_type="SOLID", **kwargs):
        # Skip parent initialization for testing
        IntrospectionMixin.__init__(self)
        self._fill_type_name = fill_type
        self._fore_color = kwargs.get('fore_color', MockColorFormat())
        self._back_color = kwargs.get('back_color', MockColorFormat())
        self._pattern = kwargs.get('pattern', MockPattern("CROSS"))
        self._gradient_stops = kwargs.get('gradient_stops', [])
        self._gradient_angle = kwargs.get('gradient_angle', 45.0)
        self._rId = kwargs.get('rId', 'rId1')

    @property
    def type(self):
        type_map = {
            "SOLID": MSO_FILL.SOLID,
            "GRADIENT": MSO_FILL.GRADIENT,
            "PATTERNED": MSO_FILL.PATTERNED,
            "PICTURE": MSO_FILL.PICTURE,
            "BACKGROUND": MSO_FILL.BACKGROUND,
            "NONE": None
        }
        return type_map.get(self._fill_type_name)

    @property
    def fore_color(self):
        if self._fill_type_name in ["SOLID", "PATTERNED"]:
            return self._fore_color
        raise TypeError(f"fill type _{self._fill_type_name}Fill has no foreground color")

    @property
    def back_color(self):
        if self._fill_type_name == "PATTERNED":
            return self._back_color
        raise TypeError(f"fill type _{self._fill_type_name}Fill has no background color")

    @property
    def pattern(self):
        if self._fill_type_name == "PATTERNED":
            return self._pattern
        raise TypeError(f"fill type _{self._fill_type_name}Fill has no pattern")

    @property
    def gradient_stops(self):
        if self._fill_type_name == "GRADIENT":
            return self._gradient_stops
        raise TypeError("Fill is not of type MSO_FILL_TYPE.GRADIENT")

    @property
    def gradient_angle(self):
        if self._fill_type_name == "GRADIENT":
            return self._gradient_angle
        raise TypeError("Fill is not of type MSO_FILL_TYPE.GRADIENT")

    @property
    def rId(self):
        if self._fill_type_name == "PICTURE":
            return self._rId
        raise NotImplementedError(f".rId property must be implemented on _{self._fill_type_name}Fill")


class MockPattern:
    """Mock pattern for pattern fill testing."""

    def __init__(self, name="CROSS"):
        self.name = name

    def __repr__(self):
        return f"MockPattern({self.name})"


class MockGradientStop(_GradientStop):
    """Mock gradient stop for gradient testing."""

    def __init__(self, position, color):
        # Skip parent initialization for testing
        self._position = position
        self._color = color

    @property
    def position(self):
        return self._position

    @property
    def color(self):
        return self._color


# =============================================================================
# Line Testing Mock Classes
# =============================================================================

class MockLineFormat(IntrospectionMixin):
    """Flexible mock LineFormat for testing different line scenarios."""

    def __init__(self, width=Pt(2.5), dash_style=MSO_LINE_DASH_STYLE.SOLID, fill_type="SOLID"):
        super().__init__()
        self._width = width
        self._dash_style = dash_style
        self._fill = MockFillFormat(fill_type)

    @property
    def width(self):
        return self._width

    @property
    def dash_style(self):
        return self._dash_style

    @property
    def fill(self):
        return self._fill

    def _to_dict_properties(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        """Return line properties for introspection."""
        props = {}

        # Fill property
        try:
            props["fill"] = self.fill.to_dict(
                include_relationships=True,
                max_depth=max_depth - 1,
                include_private=include_private,
                expand_collections=expand_collections,
                format_for_llm=format_for_llm,
                _visited_ids=_visited_ids,
            )
        except Exception as e:
            props["fill"] = {"error": f"fill access failed: {str(e)}"}

        # Width property
        try:
            props["width"] = self._format_property_value_for_to_dict(
                self.width, include_private, _visited_ids, max_depth, expand_collections, format_for_llm
            )
        except Exception as e:
            props["width"] = {"error": f"width access failed: {str(e)}"}

        # Dash style property
        try:
            props["dash_style"] = self._format_property_value_for_to_dict(
                self.dash_style, include_private, _visited_ids, max_depth, expand_collections, format_for_llm
            )
        except Exception as e:
            props["dash_style"] = {"error": f"dash_style access failed: {str(e)}"}

        return props

    def _to_dict_identity(self, _visited_ids, max_depth, expand_collections, format_for_llm, include_private):
        """Return identity information for this MockLineFormat."""
        identity = super()._to_dict_identity(_visited_ids, max_depth, expand_collections, format_for_llm, include_private)
        identity["description"] = "Mock line formatting for testing."
        return identity

    def _to_dict_llm_context(self, _visited_ids, max_depth, expand_collections, format_for_llm, include_private):
        """Return LLM-friendly context about this LineFormat."""
        context = {"description": "Describes the line (outline/border) style of an element."}

        # Generate summary based on line properties
        summary_parts = []

        try:
            line_fill_type = self.fill.type if hasattr(self.fill, 'type') else None
            line_width_pt = self.width.pt if hasattr(self.width, 'pt') else 0
            dash_style_name = self.dash_style.name if self.dash_style else "SOLID"

            # Determine if line is effectively "no line"
            if (
                line_fill_type == MSO_FILL.BACKGROUND
                or line_fill_type is None
                or line_width_pt == 0
            ):
                summary_parts.append("No line (transparent or zero width).")

            elif line_fill_type == MSO_FILL.SOLID:
                # Get color summary from fill's structured data (matches real LineFormat logic)
                try:
                    fill_dict = self.fill.to_dict(
                        max_depth=2, format_for_llm=True, _visited_ids=_visited_ids
                    )

                    # Extract color summary directly from structured data
                    color_summary = "solid color"  # Default fallback

                    # Check if this is indeed a solid fill and get fore_color
                    fill_props = fill_dict.get("properties", {})
                    fill_type = fill_props.get("type", {})
                    if isinstance(fill_type, dict) and fill_type.get("name") == "SOLID":
                        fore_color = fill_props.get("fore_color")
                        if fore_color and isinstance(fore_color, dict):
                            # Try to get color summary from fore_color's LLM context
                            fore_color_context = fore_color.get("_llm_context", {})
                            if fore_color_context.get("summary"):
                                color_summary = fore_color_context["summary"]
                            else:
                                # Fallback: construct summary from RGB or other color properties
                                rgb_info = fore_color.get("properties", {}).get("rgb")
                                if rgb_info and isinstance(rgb_info, dict):
                                    hex_val = rgb_info.get("hex", "unknown")
                                    color_summary = f"RGB color #{hex_val}"

                    # If color_summary still contains "color:" or similar, clean it up
                    if color_summary.endswith("."):
                        color_summary = color_summary.rstrip(".")

                    if dash_style_name == "SOLID":
                        summary_parts.append(
                            f"Solid line, {line_width_pt:.2f}pt, with {color_summary}."
                        )
                    else:
                        summary_parts.append(
                            f"{dash_style_name} line, {line_width_pt:.2f}pt, with {color_summary}."
                        )
                except Exception:
                    # Fallback if fill introspection fails
                    if dash_style_name == "SOLID":
                        summary_parts.append(f"Solid line, {line_width_pt:.2f}pt, with solid color.")
                    else:
                        summary_parts.append(f"{dash_style_name} line, {line_width_pt:.2f}pt, with solid color.")

            elif line_fill_type == MSO_FILL.GRADIENT:
                summary_parts.append(f"{dash_style_name} gradient line, {line_width_pt:.2f}pt.")

            elif line_fill_type == MSO_FILL.PICTURE:
                summary_parts.append(f"{dash_style_name} picture-filled line, {line_width_pt:.2f}pt.")

            elif line_fill_type == MSO_FILL.PATTERNED:
                try:
                    pattern_name = self.fill.pattern.name if hasattr(self.fill, 'pattern') and hasattr(self.fill.pattern, 'name') else "patterned"
                    summary_parts.append(f"{dash_style_name} {pattern_name} patterned line, {line_width_pt:.2f}pt.")
                except Exception:
                    summary_parts.append(f"{dash_style_name} patterned line, {line_width_pt:.2f}pt.")

            else:
                # Fallback for unknown fill types
                fill_type_name = line_fill_type.name if hasattr(line_fill_type, 'name') else "unknown"
                summary_parts.append(f"{dash_style_name} line of type {fill_type_name}, {line_width_pt:.2f}pt.")

        except Exception as e:
            # Robust fallback for any errors
            summary_parts.append(f"Line formatting (error in analysis: {str(e)}).")

        context["summary"] = (
            " ".join(summary_parts) if summary_parts else "Line formatting information."
        )

        context["common_operations"] = [
            "set line color (line.color.rgb = RGBColor(...))",
            "set line width (line.width = Pt(...))",
            "set dash style (line.dash_style = MSO_LINE.DASH)",
            "remove line (line.fill.background())",
            "set solid fill (line.fill.solid())",
        ]

        return context


# =============================================================================
# Testing Utility Functions
# =============================================================================

def create_error_triggering_obj():
    """Create an object that triggers errors during serialization for error handling tests."""

    class ErrorTriggeringObj(IntrospectionMixin):
        def __init__(self):
            self.working_prop = "fine"

        def _to_dict_properties(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
            props = super()._to_dict_properties(include_private, _visited_ids, max_depth, expand_collections, format_for_llm)

            class ProblematicValue:
                def __len__(self):
                    raise RuntimeError("Simulated error for testing")

            try:
                problematic = ProblematicValue()
                len(problematic)  # This will raise an error
            except RuntimeError as e:
                error_context = self._create_error_context("test_error", e, problematic)
                props["error_test"] = error_context

            return props

    return ErrorTriggeringObj()


def create_property_test_obj():
    """Create an object for testing property detection logic."""

    class PropertyTestObj(IntrospectionMixin):
        def __init__(self):
            self.public_attr = "public"
            self._private_attr = "private"
            self.__dunder_attr = "dunder"

        @property
        def test_property(self):
            return "property_value"

        def test_method(self):
            return "method_result"

        def _to_dict_properties(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
            attrs_to_test = ['public_attr', '_private_attr', '__dunder_attr', 'test_property', 'test_method']
            helper_results = {}

            for attr in attrs_to_test:
                if hasattr(self, attr):
                    helper_results[f"{attr}_should_include"] = self._should_include_attribute(attr, include_private)
                    helper_results[f"{attr}_is_property"] = self._is_property(attr)
                    helper_results[f"{attr}_is_introspection"] = self._is_introspection_method(attr)
                    if hasattr(self, attr):
                        attr_value = getattr(self, attr)
                        helper_results[f"{attr}_is_callable_method"] = self._is_callable_method(attr, attr_value)

            normal_props = super()._to_dict_properties(include_private, _visited_ids, max_depth, expand_collections, format_for_llm)
            normal_props.update(helper_results)
            return normal_props

    return PropertyTestObj()


# =============================================================================
# Test Assertion Helpers
# =============================================================================

def assert_basic_to_dict_structure(test_case, result, expected_object_type):
    """Assert that a to_dict result has the expected basic structure."""
    test_case.assertEqual(result['_object_type'], expected_object_type)
    test_case.assertIn('_identity', result)
    test_case.assertIn('properties', result)
    test_case.assertIn('_llm_context', result)


def assert_enum_dict_structure(test_case, enum_dict, expected_name, expected_value, has_xml_value=False):
    """Assert that an enum dict has the expected structure."""
    test_case.assertIsInstance(enum_dict, dict)
    test_case.assertIn('_object_type', enum_dict)
    test_case.assertEqual(enum_dict['name'], expected_name)
    test_case.assertEqual(enum_dict['value'], expected_value)
    test_case.assertIn('description', enum_dict)

    if has_xml_value:
        test_case.assertIn('xml_value', enum_dict)
    else:
        test_case.assertNotIn('xml_value', enum_dict)


def assert_length_dict_structure(test_case, length_dict, expected_type, expected_emu):
    """Assert that a Length object dict has the expected structure."""
    test_case.assertEqual(length_dict['_object_type'], expected_type)
    test_case.assertEqual(length_dict['emu'], expected_emu)
    test_case.assertIn('inches', length_dict)
    test_case.assertIn('pt', length_dict)
    test_case.assertIn('cm', length_dict)
    test_case.assertIn('mm', length_dict)


def assert_error_context_structure(test_case, error_dict, expected_error_type):
    """Assert that an error context dict has the expected structure."""
    test_case.assertIn("_error", error_dict)
    test_case.assertIn("_object_type", error_dict)
    test_case.assertTrue(error_dict["_object_type"].startswith("SerializationError_"))

    error_details = error_dict["_error"]
    test_case.assertEqual(error_details["type"], expected_error_type)
    test_case.assertIn("message", error_details)
    test_case.assertIn("exception_type", error_details)
    test_case.assertIn("value_type", error_details)


# =============================================================================
# Font Mock Classes for FEP-007
# =============================================================================

class MockTextCharacterProperties:
    """Mock for CT_TextCharacterProperties (a:rPr element) for Font testing."""

    def __init__(self):
        # Font properties
        self.b = None  # bold
        self.i = None  # italic
        self.u = None  # underline
        self.strike = None  # strikethrough
        self.sz = None  # size in centipoints
        self.lang = None  # language ID
        self.latin = None  # typeface info

    def get_or_add_latin(self):
        """Mock for getting/adding latin typeface element."""
        if self.latin is None:
            from unittest.mock import Mock
            self.latin = Mock()
            self.latin.typeface = None
        return self.latin

    def _remove_latin(self):
        """Mock for removing latin typeface element."""
        self.latin = None


class MockFontFillFormat(FillFormat):
    """Mock FillFormat specifically for Font color testing."""

    def __init__(self, fill_type="SOLID", color_info=None):
        # Don't call super().__init__ to avoid XML dependencies
        self._fill_type = fill_type
        self._color_info = color_info or {"rgb": "#FF0000", "summary": "RGB color: #FF0000 (R:255, G:0, B:0)."}

    def to_dict(self, include_relationships=True, max_depth=3, include_private=False,
                expand_collections=True, format_for_llm=True, _visited_ids=None):
        """Return mock fill format introspection data."""
        return {
            "_object_type": "FillFormat",
            "properties": {"type": {"_object_type": "MSO_FILL", "name": self._fill_type}},
            "_llm_context": {"summary": self._color_info["summary"]}
        }


# =============================================================================
# Image and Picture Mock Classes for FEP-015
# =============================================================================

class MockImage(IntrospectionMixin):
    """Mock Image class for testing Image introspection."""

    def __init__(self, filename="test.png", content_type="image/png",
                 ext="png", size=(800, 600), dpi=(72, 72), blob_size=100000):
        super().__init__()
        self._filename = filename
        self._content_type = content_type
        self._ext = ext
        self._size = size
        self._dpi = dpi
        self._blob_size = blob_size
        self._sha1 = "abcd1234567890abcd1234567890abcd12345678"

    @property
    def filename(self):
        return self._filename

    @property
    def content_type(self):
        return self._content_type

    @property
    def ext(self):
        return self._ext

    @property
    def size(self):
        return self._size

    @property
    def dpi(self):
        return self._dpi

    @property
    def sha1(self):
        return self._sha1

    @property
    def blob(self):
        return b"mock_image_data" * (self._blob_size // 15)  # Simulate blob

    def _to_dict_identity(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        """Provide identity information for Image introspection."""
        identity = super()._to_dict_identity(include_private, _visited_ids, max_depth, expand_collections, format_for_llm)
        description = f"Image data: {self.filename if self.filename else 'streamed image'} ({self.ext})"
        identity["description"] = description
        if self.filename:
            identity["filename"] = self.filename
        return identity

    def _to_dict_properties(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        """Provide properties for Image introspection."""
        props = {}
        props["content_type"] = self.content_type
        props["extension"] = self.ext
        props["sha1_hash"] = self.sha1
        props["dimensions_px"] = {"width": self.size[0], "height": self.size[1]}
        props["dpi"] = {"horizontal": self.dpi[0], "vertical": self.dpi[1]}
        props["blob_size_bytes"] = len(self.blob)
        return props

    def _to_dict_relationships(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        """Provide relationships for Image introspection."""
        return {}

    def _to_dict_llm_context(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        """Provide LLM-friendly context for Image introspection."""
        filename_desc = f"'{self.filename}'" if self.filename else "a streamed image"
        width, height = self.size
        size_mb = len(self.blob) / (1024 * 1024)

        description = (
            f"An {self.ext.upper()} image {filename_desc} with dimensions {width}x{height} pixels "
            f"at {self.dpi[0]}x{self.dpi[1]} DPI (file size: {size_mb:.2f} MB)."
        )

        summary = f"{self.ext.upper()} image: {width}x{height}px, {size_mb:.2f}MB"

        common_operations = [
            "access image binary data via .blob property",
            "get image dimensions via .size property",
            "check image format via .ext property",
            "verify integrity via .sha1 hash",
            "examine DPI settings via .dpi property"
        ]

        return {
            "description": description,
            "summary": summary,
            "common_operations": common_operations
        }


class MockPicture(IntrospectionMixin):
    """Mock Picture class for testing Picture introspection."""

    def __init__(self, image=None, crop_left=0.0, crop_top=0.0, crop_right=0.0, crop_bottom=0.0,
                 auto_shape_type=None, shape_id=1, name="Picture 1"):
        super().__init__()
        self._image = image or MockImage()
        self._crop_left = crop_left
        self._crop_top = crop_top
        self._crop_right = crop_right
        self._crop_bottom = crop_bottom
        self._auto_shape_type = auto_shape_type
        self._shape_id = shape_id
        self._name = name
        self._line = MockLineFormat()

    @property
    def image(self):
        return self._image

    @property
    def crop_left(self):
        return self._crop_left

    @property
    def crop_top(self):
        return self._crop_top

    @property
    def crop_right(self):
        return self._crop_right

    @property
    def crop_bottom(self):
        return self._crop_bottom

    @property
    def auto_shape_type(self):
        return self._auto_shape_type

    @property
    def shape_id(self):
        return self._shape_id

    @property
    def name(self):
        return self._name

    @property
    def line(self):
        return self._line

    def _to_dict_identity(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        """Provide identity information for Picture introspection."""
        identity = {}
        identity["class_name"] = "Picture"
        identity["shape_id"] = self.shape_id
        identity["name"] = self.name

        try:
            img_desc = "unknown image"
            if self.image is not None:
                if self.image.filename:
                    img_desc = self.image.filename
                else:
                    img_desc = f"streamed {self.image.ext} image"
            else:
                img_desc = "no embedded image"
        except (ValueError, AttributeError):
            img_desc = "no embedded image"

        identity["description"] = f"Picture shape displaying: {img_desc}"
        return identity

    def _to_dict_properties(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        """Provide properties for Picture introspection."""
        props = {}

        # Add crop properties
        props["crop_left"] = self.crop_left
        props["crop_top"] = self.crop_top
        props["crop_right"] = self.crop_right
        props["crop_bottom"] = self.crop_bottom

        # Add image details
        if self.image is not None and hasattr(self.image, 'to_dict') and max_depth > 0:
            props["image_details"] = self.image.to_dict(
                include_relationships=True,
                max_depth=max_depth - 1,
                include_private=include_private,
                expand_collections=expand_collections,
                format_for_llm=format_for_llm,
                _visited_ids=_visited_ids
            )
        elif self.image is not None:
            props["image_details"] = {
                "_object_type": "Image",
                "_summary_or_truncated": True,
                "filename": getattr(self.image, 'filename', None)
            }
        else:
            props["image_details"] = "No embedded image"

        # Add auto shape type
        props["auto_shape_mask_type"] = self._format_property_value_for_to_dict(
            self.auto_shape_type, include_private, _visited_ids, max_depth - 1, expand_collections, format_for_llm
        )

        # Add line properties
        if max_depth > 0:
            props["line"] = self.line.to_dict(
                include_relationships=True,
                max_depth=max_depth - 1,
                include_private=include_private,
                expand_collections=expand_collections,
                format_for_llm=format_for_llm,
                _visited_ids=_visited_ids
            )
        else:
            props["line"] = {"_object_type": "LineFormat", "_depth_exceeded": True}

        return props

    def _to_dict_relationships(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        """Provide relationships for Picture introspection."""
        return {"image_part": "Mock image part reference"}

    def _to_dict_llm_context(self, include_private, _visited_ids, max_depth, expand_collections, format_for_llm):
        """Provide LLM-friendly context for Picture introspection."""
        try:
            filename_info = "unknown image"
            crop_info = ""
            mask_info = ""

            if self.image is not None:
                if self.image.filename:
                    filename_info = self.image.filename
                else:
                    filename_info = f"streamed {self.image.ext} image"
            else:
                filename_info = "no embedded image"

            if any([self.crop_left, self.crop_top, self.crop_right, self.crop_bottom]):
                crop_parts = []
                if self.crop_left: crop_parts.append(f"{self.crop_left*100:.1f}% from left")
                if self.crop_top: crop_parts.append(f"{self.crop_top*100:.1f}% from top")
                if self.crop_right: crop_parts.append(f"{self.crop_right*100:.1f}% from right")
                if self.crop_bottom: crop_parts.append(f"{self.crop_bottom*100:.1f}% from bottom")
                crop_info = f" Cropped {', '.join(crop_parts)}."

            if self.auto_shape_type and hasattr(self.auto_shape_type, 'name'):
                mask_info = f" Masked as {self.auto_shape_type.name}."

        except (ValueError, AttributeError):
            filename_info = "no embedded image"

        description = (
            f"A PICTURE shape '{self.name}' (ID: {self.shape_id}) displaying: {filename_info}.{mask_info}{crop_info}"
        )

        summary = f"Picture: {filename_info}"

        common_operations = [
            "change image source (replace with new image)",
            "adjust crop properties (crop_left, crop_top, crop_right, crop_bottom)",
            "set mask shape via auto_shape_type property",
            "modify border line properties"
        ]

        return {
            "description": description,
            "summary": summary,
            "common_operations": common_operations
        }


# =============================================================================
# Tree Functionality Mock Classes for FEP-020
# =============================================================================

def create_mock_length(inches=1.0):
    """Create a mock Length object for testing."""
    from unittest.mock import Mock
    length = Mock()
    length.inches = inches
    length.pt = inches * 72
    length.cm = inches * 2.54
    length.mm = inches * 25.4
    length.emu = int(inches * 914400)
    return length


class MockShape(IntrospectionMixin):
    """Mock shape for tree functionality testing."""
    
    def __init__(self, shape_id=42, name="Test Shape", has_text_frame=False):
        super().__init__()
        self.shape_id = shape_id
        self.name = name
        self.has_text_frame = has_text_frame
        self.is_placeholder = False
        
        # Geometry properties
        self.left = create_mock_length(1.0)
        self.top = create_mock_length(2.0)
        self.width = create_mock_length(3.0)
        self.height = create_mock_length(1.5)
        self.rotation = 0.0
        
        # Text content for testing
        self.text = "Sample text content" if has_text_frame else ""
        
        # Mock parent and part
        from unittest.mock import Mock
        self._parent = Mock()
        
    def _get_shape_type_safely(self):
        """Mock shape type access."""
        from unittest.mock import Mock
        shape_type = Mock()
        shape_type.name = "RECTANGLE"
        return shape_type
        
    @property
    def part(self):
        from unittest.mock import Mock
        return Mock()
    
    # Tree functionality implementation matching real BaseShape
    def _to_tree_node_identity(self):
        """Override to provide rich shape identity for tree node representation."""
        identity = {
            "shape_id": self.shape_id,
            "name": self.name,
            "class_name": type(self).__name__,
        }

        # Add shape type if available
        shape_type = self._get_shape_type_safely()
        if shape_type is not None:
            identity["shape_type"] = shape_type.name

        # Add placeholder information if applicable
        if self.is_placeholder:
            try:
                # Mock placeholder format 
                identity["placeholder_type"] = "TITLE"
                identity["placeholder_idx"] = 0
            except (ValueError, AttributeError):
                identity["placeholder_type"] = "UNKNOWN"

        return identity

    def _to_tree_node_geometry(self):
        """Override to provide shape geometry for tree node representation."""
        try:
            return {
                "left": f"{self.left.inches:.2f} in",
                "top": f"{self.top.inches:.2f} in",
                "width": f"{self.width.inches:.2f} in",
                "height": f"{self.height.inches:.2f} in",
                "rotation": f"{self.rotation:.1f}°" if self.rotation != 0 else "0°",
            }
        except Exception:
            # Fallback if geometry access fails
            return None

    def _to_tree_node_content_summary(self):
        """Override to provide shape content summary for tree node representation."""
        # Build summary parts
        summary_parts = []

        # Shape type and name
        shape_type = self._get_shape_type_safely()
        if shape_type:
            summary_parts.append(f"{shape_type.name}")
        else:
            summary_parts.append("Shape")

        # Add name if different from default pattern
        if self.name and not self.name.startswith((type(self).__name__, "Shape")):
            summary_parts.append(f"'{self.name}'")

        # Placeholder context
        if self.is_placeholder:
            summary_parts.append("(placeholder)")

        # Text content hint for shapes with text frames
        if hasattr(self, 'has_text_frame') and self.has_text_frame:
            try:
                if hasattr(self, 'text') and self.text:
                    text_content = self.text.strip()
                    if text_content:
                        # Truncate long text
                        if len(text_content) > 30:
                            text_content = text_content[:27] + "..."
                        summary_parts.append(f"Text: '{text_content}'")
                    else:
                        summary_parts.append("(empty text)")
            except Exception:
                summary_parts.append("(text frame)")

        return " ".join(summary_parts)


class MockPlaceholderFormat(IntrospectionMixin):
    """Mock placeholder format for testing."""
    
    def __init__(self, placeholder_type=PP_PLACEHOLDER.TITLE, idx=0):
        super().__init__()
        self.type = placeholder_type
        self.idx = idx
        
    def to_dict(self, **kwargs):
        """Mock to_dict implementation."""
        return {
            "_object_type": "_PlaceholderFormat",
            "properties": {
                "type": {
                    "_object_type": "PP_PLACEHOLDER_TYPE",
                    "name": self.type.name,
                    "value": int(self.type),
                    "description": ""
                },
                "idx": self.idx
            }
        }


class MockGroupShape(MockShape):
    """Mock group shape for tree functionality testing."""
    
    def __init__(self, shape_id=99, name="Test Group"):
        super().__init__(shape_id, name)
        self.shapes = [MockShape(1, "Child 1"), MockShape(2, "Child 2")]
        
    def _get_shape_type_safely(self):
        """Mock group shape type."""
        from unittest.mock import Mock
        shape_type = Mock()
        shape_type.name = "GROUP"
        return shape_type
        
    # Tree functionality implementation matching real GroupShape
    def get_tree(self, max_depth=2):
        """Generate a hierarchical tree view of this group shape and its contents."""
        access_path = f"group_shape_{self.shape_id}"
        return self._to_tree_node(access_path, max_depth, _current_depth=0)

    def _to_tree_node_content_summary(self):
        """Override to provide group-specific content summary for tree node representation."""
        summary_parts = []

        # Group identifier
        summary_parts.append("Group")

        # Add name if it's meaningful (not default pattern)
        if self.name and not self.name.startswith(("Group", "Grouped")):
            summary_parts.append(f"'{self.name}'")

        # Member shape count
        try:
            shape_count = len(self.shapes)
            if shape_count > 0:
                summary_parts.append(f"({shape_count} shape{'s' if shape_count != 1 else ''})")
            else:
                summary_parts.append("(empty)")
        except Exception:
            summary_parts.append("(unknown contents)")

        return " ".join(summary_parts)

    def _to_tree_node_children(self, access_path, max_depth, current_depth):
        """Override to provide group's member shapes as children."""
        if current_depth > max_depth:
            return None

        children = []

        try:
            # Add member shapes
            for i, shape in enumerate(self.shapes):
                shape_access_path = f"{access_path}.shapes[{i}]"
                if hasattr(shape, '_to_tree_node'):
                    child_node = shape._to_tree_node(shape_access_path, max_depth, current_depth + 1)
                    children.append(child_node)
                else:
                    # Fallback for shapes without tree node support
                    children.append({
                        "_object_type": type(shape).__name__,
                        "_identity": {
                            "shape_id": getattr(shape, 'shape_id', 'unknown'),
                            "class_name": type(shape).__name__
                        },
                        "access_path": shape_access_path,
                        "geometry": None,
                        "content_summary": f"{type(shape).__name__} object",
                        "children": None
                    })

        except Exception:
            # If we can't access shapes, return empty children list
            pass

        return children if children else None


class MockSlide(IntrospectionMixin):
    """Mock slide for tree functionality testing."""
    
    def __init__(self, slide_id=256, name="Test Slide"):
        super().__init__()
        self.slide_id = slide_id
        self.name = name
        self.has_notes_slide = False
        self.follow_master_background = True
        
        # Mock shapes collection
        self.shapes = [
            MockShape(1, "Title 1", has_text_frame=True),
            MockShape(2, "Content 1"),
            MockGroupShape(3, "Group 1")
        ]
        
        # Mock placeholders collection
        self.placeholders = {
            0: MockShape(1, "Title 1", has_text_frame=True),
            1: MockShape(2, "Content 1")
        }
        
        # Mock shapes collection with title property
        from unittest.mock import Mock
        shapes_with_title = Mock()
        shapes_with_title.__len__ = lambda: 3
        shapes_with_title.__iter__ = lambda: iter(self.shapes)
        shapes_with_title.title = self.shapes[0]  # First shape is title
        self.shapes = shapes_with_title
        
        # Mock placeholders collection with len and items
        placeholders_with_methods = Mock()
        placeholders_with_methods.__len__ = lambda: 2
        placeholders_with_methods.items = lambda: self.placeholders.items()
        self.placeholders = placeholders_with_methods
        
        # Mock part for parent relationship
        from unittest.mock import Mock
        self.part = Mock()
        
    @property
    def slide_layout(self):
        """Mock slide layout."""
        from unittest.mock import Mock
        layout = Mock()
        layout.name = "Title Slide"
        return layout
    
    # Tree functionality implementation matching real Slide
    def get_tree(self, max_depth=2):
        """Generate a hierarchical tree view of this slide and its shapes."""
        # For mock, use a simple access path
        access_path = f"slides[slide_id_{self.slide_id}]"
        return self._to_tree_node(access_path, max_depth, _current_depth=0)

    def _to_tree_node_identity(self):
        """Override to provide rich slide identity for tree node representation."""
        identity = {
            "slide_id": self.slide_id,
            "class_name": "Slide",
        }

        if self.name:
            identity["name"] = self.name

        # Add layout information
        try:
            layout = self.slide_layout
            if layout and layout.name:
                identity["layout_name"] = layout.name
        except Exception:
            pass

        return identity

    def _to_tree_node_geometry(self):
        """Override - slides don't have geometry, return None."""
        return None

    def _to_tree_node_content_summary(self):
        """Override to provide slide content summary for tree node representation."""
        summary_parts = []

        # Slide identifier
        if self.name:
            summary_parts.append(f"Slide: '{self.name}'")
        else:
            summary_parts.append(f"Slide {self.slide_id}")

        # Get title if available
        try:
            if hasattr(self.shapes, 'title') and self.shapes.title:
                title_text = getattr(self.shapes.title, 'text', '').strip()
                if title_text:
                    if len(title_text) > 30:
                        title_text = title_text[:27] + "..."
                    summary_parts.append(f"- {title_text}")
        except Exception:
            pass

        # Shape and placeholder counts
        try:
            # Try to get the count from the original shapes if available
            if hasattr(self.shapes, '_original_shapes'):
                shape_count = len(self.shapes._original_shapes)
            else:
                shape_count = len(self.shapes)
        except:
            shape_count = 3  # fallback
        
        try:
            # Try to get the count from the original placeholders if available
            if hasattr(self.placeholders, '_original_placeholders'):
                placeholder_count = len(self.placeholders._original_placeholders)
            else:
                placeholder_count = len(self.placeholders)
        except:
            placeholder_count = 2  # fallback

        counts = []
        if shape_count > 0:
            counts.append(f"{shape_count} shape{'s' if shape_count != 1 else ''}")
        if placeholder_count > 0:
            counts.append(f"{placeholder_count} placeholder{'s' if placeholder_count != 1 else ''}")

        if counts:
            summary_parts.append(f"({', '.join(counts)})")

        return " ".join(summary_parts)

    def _to_tree_node_children(self, access_path, max_depth, current_depth):
        """Override to provide slide's shapes as children."""
        if current_depth > max_depth:
            return None

        children = []

        try:
            # Add shapes
            for i, shape in enumerate(self.shapes):
                shape_access_path = f"{access_path}.shapes[{i}]"
                if hasattr(shape, '_to_tree_node'):
                    child_node = shape._to_tree_node(shape_access_path, max_depth, current_depth + 1)
                    children.append(child_node)
                else:
                    # Fallback for shapes without tree node support
                    children.append({
                        "_object_type": type(shape).__name__,
                        "_identity": {"class_name": type(shape).__name__},
                        "access_path": shape_access_path,
                        "geometry": None,
                        "content_summary": f"{type(shape).__name__} object",
                        "children": None
                    })

        except Exception:
            # If we can't access shapes, return empty children list
            pass

        return children if children else None


class MockPresentation(IntrospectionMixin):
    """Mock presentation for tree functionality testing."""
    
    def __init__(self, title="Test Presentation"):
        super().__init__()
        
        # Mock core properties
        from unittest.mock import Mock
        self.core_properties = Mock()
        self.core_properties.title = title
        
        # Mock slide dimensions
        self.slide_width = create_mock_length(10.0)
        self.slide_height = create_mock_length(7.5)
        
        # Mock slides collection
        self.slides = [
            MockSlide(256, "Slide 1"),
            MockSlide(257, "Slide 2"), 
            MockSlide(258, "Slide 3")
        ]
        
        # Mock slide masters collection
        from unittest.mock import Mock
        self.slide_masters = [Mock()]
        
        # Mock notes master
        self.notes_master = Mock()
        
        # Mock part
        self.part = Mock()
        
        # Add len methods to collections
        slides_with_len = Mock()
        slides_with_len.__len__ = lambda: 3
        slides_with_len.__iter__ = lambda: iter(self.slides)
        slides_with_len.__getitem__ = lambda _, i: self.slides[i]
        self.slides = slides_with_len
        
        masters_with_len = Mock()
        masters_with_len.__len__ = lambda: 1
        masters_with_len.__iter__ = lambda: iter(self.slide_masters)
        self.slide_masters = masters_with_len

    # Tree functionality implementation matching real Presentation
    def get_tree(self, max_depth=2):
        """Generate a hierarchical tree view of this presentation and its slides."""
        return self._to_tree_node("", max_depth, _current_depth=0)

    def _to_tree_node_identity(self):
        """Override to provide rich presentation identity for tree node representation."""
        identity = {
            "class_name": "Presentation",
        }

        # Add title from core properties
        try:
            if self.core_properties and self.core_properties.title:
                identity["title"] = self.core_properties.title
        except Exception:
            pass

        # Add slide and master counts
        try:
            identity["slide_count"] = len(self.slides)
            identity["master_count"] = len(self.slide_masters)
        except Exception:
            pass

        # Add slide dimensions
        try:
            if self.slide_width and self.slide_height:
                identity["slide_width"] = f"{self.slide_width.inches:.2f} in"
                identity["slide_height"] = f"{self.slide_height.inches:.2f} in"
        except Exception:
            pass

        return identity

    def _to_tree_node_geometry(self):
        """Override - presentations don't have geometry, return None."""
        return None

    def _to_tree_node_content_summary(self):
        """Override to provide presentation content summary for tree node representation."""
        summary_parts = []

        # Presentation title
        title = "Untitled Presentation"
        try:
            if self.core_properties and self.core_properties.title:
                title = self.core_properties.title
        except Exception:
            pass

        summary_parts.append(f"Presentation: '{title}'")

        # Slide and master counts
        try:
            slide_count = len(self.slides)
            master_count = len(self.slide_masters)

            counts = []
            if slide_count > 0:
                counts.append(f"{slide_count} slide{'s' if slide_count != 1 else ''}")
            if master_count > 0:
                counts.append(f"{master_count} master{'s' if master_count != 1 else ''}")

            if counts:
                summary_parts.append(f"({', '.join(counts)})")
        except Exception:
            pass

        return " ".join(summary_parts)

    def _to_tree_node_children(self, access_path, max_depth, current_depth):
        """Override to provide presentation's slides as children."""
        if current_depth > max_depth:
            return None

        children = []

        try:
            # Add slides
            for i, slide in enumerate(self.slides):
                slide_access_path = f"slides[{i}]"
                if hasattr(slide, '_to_tree_node'):
                    child_node = slide._to_tree_node(slide_access_path, max_depth, current_depth + 1)
                    children.append(child_node)
                else:
                    # Fallback for slides without tree node support
                    children.append({
                        "_object_type": "Slide",
                        "_identity": {
                            "slide_id": getattr(slide, 'slide_id', 'unknown'),
                            "class_name": "Slide"
                        },
                        "access_path": slide_access_path,
                        "geometry": None,
                        "content_summary": f"Slide {getattr(slide, 'slide_id', 'unknown')}",
                        "children": None
                    })

        except Exception:
            # If we can't access slides, return empty children list
            pass

        return children if children else None
