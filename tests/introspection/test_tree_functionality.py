# tests/introspection/test_tree_functionality.py
"""
Unit tests for FEP-020: Wide-Angle Tree View functionality.

This module tests the get_tree() method and related _to_tree_node() functionality
across container objects (Presentation, Slide, GroupShape) and shape objects.
"""

import unittest
from unittest.mock import Mock, PropertyMock, patch, MagicMock

from pptx.introspection import IntrospectionMixin
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from .mock_helpers import (
    MockPresentation, MockSlide, MockGroupShape, MockShape, MockPlaceholderFormat,
    assert_basic_to_dict_structure, create_mock_length,
)


class TestIntrospectionMixinTreeNode(unittest.TestCase):
    """Test the base tree node functionality in IntrospectionMixin."""

    def setUp(self):
        self.mixin = IntrospectionMixin()

    def test_to_tree_node_basic_structure(self):
        """Test that _to_tree_node returns correct basic structure."""
        result = self.mixin._to_tree_node("test_path", max_depth=1)
        
        # Check required keys
        self.assertIn("_object_type", result)
        self.assertIn("_identity", result)
        self.assertIn("access_path", result)
        self.assertIn("geometry", result)
        self.assertIn("content_summary", result)
        
        # Check values
        self.assertEqual(result["_object_type"], "IntrospectionMixin")
        self.assertEqual(result["access_path"], "test_path")
        self.assertIsInstance(result["_identity"], dict)
        self.assertIsNone(result["geometry"])  # Default implementation
        self.assertIsInstance(result["content_summary"], str)

    def test_to_tree_node_identity_default(self):
        """Test default identity implementation."""
        identity = self.mixin._to_tree_node_identity()
        
        self.assertIn("class_name", identity)
        self.assertIn("memory_address", identity)
        self.assertEqual(identity["class_name"], "IntrospectionMixin")
        self.assertTrue(identity["memory_address"].startswith("0x"))

    def test_to_tree_node_geometry_default(self):
        """Test default geometry implementation returns None."""
        geometry = self.mixin._to_tree_node_geometry()
        self.assertIsNone(geometry)

    def test_to_tree_node_content_summary_default(self):
        """Test default content summary implementation."""
        summary = self.mixin._to_tree_node_content_summary()
        self.assertEqual(summary, "IntrospectionMixin object")

    def test_to_tree_node_children_default(self):
        """Test default children implementation returns None."""
        children = self.mixin._to_tree_node_children("test_path", 2, 1)
        self.assertIsNone(children)

    def test_to_tree_node_depth_limiting(self):
        """Test that tree node respects max_depth limits."""
        # At max depth, should not include children
        result = self.mixin._to_tree_node("test_path", max_depth=1, _current_depth=1)
        self.assertNotIn("children", result)
        
        # Below max depth, should call children method
        result = self.mixin._to_tree_node("test_path", max_depth=2, _current_depth=1)
        # Since default implementation returns None, children won't be added
        self.assertNotIn("children", result)


class TestBaseShapeTreeNode(unittest.TestCase):
    """Test tree node functionality for BaseShape objects."""

    def setUp(self):
        self.shape = MockShape()
        
    def test_tree_node_identity(self):
        """Test BaseShape tree node identity includes shape-specific info."""
        identity = self.shape._to_tree_node_identity()
        
        # Check required shape identity fields
        self.assertIn("shape_id", identity)
        self.assertIn("name", identity)
        self.assertIn("class_name", identity)
        
        # Check values
        self.assertEqual(identity["shape_id"], 42)
        self.assertEqual(identity["name"], "Test Shape")
        self.assertEqual(identity["class_name"], "MockShape")

    @unittest.skip("Complex mocking required - shape type integration tested in live tests")
    def test_tree_node_identity_with_shape_type(self):
        """Test identity includes shape type when available."""
        # Mock shape type
        with patch.object(self.shape, '_get_shape_type_safely', return_value=Mock(name="RECTANGLE")):
            identity = self.shape._to_tree_node_identity()
            self.assertIn("shape_type", identity)
            self.assertEqual(identity["shape_type"], "RECTANGLE")

    @unittest.skip("Complex mocking required - placeholder format integration tested in live tests")
    def test_tree_node_identity_placeholder(self):
        """Test identity includes placeholder info for placeholder shapes."""
        # Make shape a placeholder
        self.shape.is_placeholder = True
        mock_ph_format = MockPlaceholderFormat()
        
        with patch.object(self.shape, 'placeholder_format', mock_ph_format):
            identity = self.shape._to_tree_node_identity()
            
            self.assertIn("placeholder_type", identity)
            self.assertIn("placeholder_idx", identity)
            self.assertEqual(identity["placeholder_type"], "TITLE")
            self.assertEqual(identity["placeholder_idx"], 0)

    def test_tree_node_geometry(self):
        """Test BaseShape geometry information."""
        geometry = self.shape._to_tree_node_geometry()
        
        self.assertIsInstance(geometry, dict)
        self.assertIn("left", geometry)
        self.assertIn("top", geometry)  
        self.assertIn("width", geometry)
        self.assertIn("height", geometry)
        self.assertIn("rotation", geometry)
        
        # Check formatted values
        self.assertTrue(geometry["left"].endswith(" in"))
        self.assertTrue(geometry["top"].endswith(" in"))
        self.assertTrue(geometry["width"].endswith(" in"))
        self.assertTrue(geometry["height"].endswith(" in"))
        self.assertTrue(geometry["rotation"].endswith("°"))

    def test_tree_node_content_summary_basic(self):
        """Test basic content summary for shapes."""
        summary = self.shape._to_tree_node_content_summary()
        
        self.assertIsInstance(summary, str)
        self.assertTrue(len(summary) > 0)
        # Should include shape type or fallback
        self.assertIn("Shape", summary)

    def test_tree_node_content_summary_with_text(self):
        """Test content summary includes text content when available."""
        # Mock text frame capabilities
        self.shape.has_text_frame = True
        
        with patch.object(self.shape, 'text', "Sample text content"):
            summary = self.shape._to_tree_node_content_summary()
            self.assertIn("Sample text content", summary)


class TestSlideTreeFunctionality(unittest.TestCase):
    """Test tree functionality for Slide objects."""

    def setUp(self):
        self.slide = MockSlide()

    def test_get_tree_basic_structure(self):
        """Test slide get_tree returns correct structure."""
        tree = self.slide.get_tree(max_depth=1)
        
        # Check required keys
        self.assertIn("_object_type", tree)
        self.assertIn("_identity", tree)
        self.assertIn("access_path", tree)
        self.assertIn("geometry", tree)
        self.assertIn("content_summary", tree)
        
        # Check values
        self.assertEqual(tree["_object_type"], "MockSlide")
        self.assertTrue(tree["access_path"].startswith("slides["))
        self.assertIsNone(tree["geometry"])  # Slides don't have geometry

    def test_slide_tree_node_identity(self):
        """Test slide identity information."""
        identity = self.slide._to_tree_node_identity()
        
        self.assertIn("slide_id", identity)
        self.assertIn("class_name", identity)
        self.assertEqual(identity["slide_id"], 256)
        self.assertEqual(identity["class_name"], "MockSlide")

    def test_slide_content_summary(self):
        """Test slide content summary includes slide info."""
        summary = self.slide._to_tree_node_content_summary()
        
        self.assertIn("Slide", summary)
        self.assertIn("256", summary)  # slide_id
        # Should include shape/placeholder counts
        self.assertIn("shape", summary.lower())

    def test_slide_children_with_shapes(self):
        """Test slide children include shapes."""
        children = self.slide._to_tree_node_children("slides[0]", max_depth=2, current_depth=1)
        
        self.assertIsInstance(children, list)
        self.assertEqual(len(children), 3)  # MockSlide has 3 shapes
        
        # Check first child structure
        child = children[0]
        self.assertIn("_object_type", child)
        self.assertIn("access_path", child)
        self.assertTrue(child["access_path"].startswith("slides[0].shapes["))

    def test_slide_children_depth_limiting(self):
        """Test slide respects max_depth for children."""
        # At max depth should return None
        children = self.slide._to_tree_node_children("slides[0]", max_depth=1, current_depth=1)
        self.assertIsNone(children)


class TestPresentationTreeFunctionality(unittest.TestCase):
    """Test tree functionality for Presentation objects."""

    def setUp(self):
        self.presentation = MockPresentation()

    def test_get_tree_basic_structure(self):
        """Test presentation get_tree returns correct structure."""
        tree = self.presentation.get_tree(max_depth=1)
        
        # Check required keys
        self.assertIn("_object_type", tree)
        self.assertIn("_identity", tree)
        self.assertIn("access_path", tree)
        self.assertIn("geometry", tree)
        self.assertIn("content_summary", tree)
        
        # Check values
        self.assertEqual(tree["_object_type"], "MockPresentation")
        self.assertEqual(tree["access_path"], "")  # Root has empty path
        self.assertIsNone(tree["geometry"])  # Presentations don't have geometry

    def test_presentation_tree_node_identity(self):
        """Test presentation identity information."""
        identity = self.presentation._to_tree_node_identity()
        
        self.assertIn("class_name", identity)
        self.assertEqual(identity["class_name"], "Presentation")
        
        # Should include counts
        self.assertIn("slide_count", identity)
        self.assertIn("master_count", identity)

    def test_presentation_content_summary(self):
        """Test presentation content summary."""
        summary = self.presentation._to_tree_node_content_summary()
        
        self.assertIn("Presentation", summary)
        # Should include slide count
        self.assertIn("slide", summary.lower())

    def test_presentation_children_with_slides(self):
        """Test presentation children include slides."""
        children = self.presentation._to_tree_node_children("", max_depth=2, current_depth=0)
        
        self.assertIsInstance(children, list)
        self.assertEqual(len(children), 3)  # MockPresentation has 3 slides
        
        # Check first child structure
        child = children[0]
        self.assertIn("_object_type", child)
        self.assertIn("access_path", child)
        self.assertEqual(child["access_path"], "slides[0]")


class TestGroupShapeTreeFunctionality(unittest.TestCase):
    """Test tree functionality for GroupShape objects."""

    def setUp(self):
        self.group = MockGroupShape()

    def test_get_tree_basic_structure(self):
        """Test group get_tree returns correct structure."""
        tree = self.group.get_tree(max_depth=1)
        
        # Check required keys
        self.assertIn("_object_type", tree)
        self.assertIn("_identity", tree)
        self.assertIn("access_path", tree)
        self.assertIn("geometry", tree)
        self.assertIn("content_summary", tree)
        
        # Check values
        self.assertEqual(tree["_object_type"], "MockGroupShape")
        # Access path should be generated for standalone group
        self.assertTrue(tree["access_path"].startswith("group_shape_"))

    def test_group_content_summary(self):
        """Test group content summary includes member count."""
        summary = self.group._to_tree_node_content_summary()
        
        self.assertIn("Group", summary)
        # Should include shape count
        self.assertIn("shape", summary.lower())

    def test_group_children_with_shapes(self):
        """Test group children include member shapes."""
        children = self.group._to_tree_node_children("slides[0].shapes[2]", max_depth=2, current_depth=1)
        
        self.assertIsInstance(children, list)
        self.assertEqual(len(children), 2)  # MockGroupShape has 2 member shapes
        
        # Check first child structure
        child = children[0]
        self.assertIn("_object_type", child)
        self.assertIn("access_path", child)
        self.assertTrue(child["access_path"].startswith("slides[0].shapes[2].shapes["))


class TestTreeDepthAndRecursion(unittest.TestCase):
    """Test depth limiting and recursion behavior in tree functionality."""

    def setUp(self):
        self.presentation = MockPresentation()
        self.slide = MockSlide()
        self.group = MockGroupShape()

    def test_max_depth_zero(self):
        """Test max_depth=0 returns only the root node."""
        tree = self.presentation.get_tree(max_depth=0)
        
        self.assertNotIn("children", tree)

    def test_max_depth_one(self):
        """Test max_depth=1 includes immediate children only."""
        tree = self.presentation.get_tree(max_depth=1)
        
        self.assertIn("children", tree)
        # Children should be slides, but slides shouldn't have children
        for child in tree["children"]:
            self.assertNotIn("children", child)

    def test_max_depth_two(self):
        """Test max_depth=2 includes grandchildren."""
        tree = self.presentation.get_tree(max_depth=2)
        
        self.assertIn("children", tree)
        # First slide should have shapes as children
        slide_child = tree["children"][0]
        if "children" in slide_child and slide_child["children"]:
            # Shape children should not have their own children at this depth
            shape_child = slide_child["children"][0]
            self.assertNotIn("children", shape_child)

    def test_nested_group_recursion(self):
        """Test recursion works correctly with nested groups."""
        # Create a group with nested structure
        nested_group = MockGroupShape()
        nested_group.shapes = [MockShape(), MockGroupShape()]  # Group containing another group
        
        tree = nested_group.get_tree(max_depth=3)
        
        # Should have children (member shapes)
        self.assertIn("children", tree)
        children = tree["children"]
        self.assertEqual(len(children), 2)
        
        # Second child should be the nested group
        nested_child = children[1]
        self.assertEqual(nested_child["_object_type"], "MockGroupShape")


class TestTreeAccessPaths(unittest.TestCase):
    """Test access path generation for tree nodes."""

    def test_presentation_root_path(self):
        """Test presentation has empty root access path."""
        prs = MockPresentation()
        tree = prs.get_tree(max_depth=0)
        self.assertEqual(tree["access_path"], "")

    def test_slide_access_path(self):
        """Test slide access paths are generated correctly."""
        slide = MockSlide()
        
        # Mock the parent presentation relationship
        with patch.object(slide, 'part') as mock_part:
            mock_prs = MockPresentation()
            mock_part.package.presentation_part.presentation = mock_prs
            
            tree = slide.get_tree(max_depth=0)
            self.assertTrue(tree["access_path"].startswith("slides["))
            self.assertTrue(tree["access_path"].endswith("]"))

    def test_shape_access_path_inheritance(self):
        """Test shapes inherit access paths from parents."""
        slide = MockSlide()
        children = slide._to_tree_node_children("slides[0]", max_depth=2, current_depth=1)
        
        if children:
            for i, child in enumerate(children):
                expected_path = f"slides[0].shapes[{i}]"
                self.assertEqual(child["access_path"], expected_path)

    def test_group_member_access_paths(self):
        """Test group member shapes get correct access paths."""
        group = MockGroupShape()
        children = group._to_tree_node_children("slides[0].shapes[2]", max_depth=2, current_depth=1)
        
        if children:
            for i, child in enumerate(children):
                expected_path = f"slides[0].shapes[2].shapes[{i}]"
                self.assertEqual(child["access_path"], expected_path)


if __name__ == "__main__":
    unittest.main()