"""
Test SlideLayout and LayoutPlaceholder introspection functionality - FEP-016

This module tests the SlideLayout.to_dict() and LayoutPlaceholder.to_dict() introspection 
capabilities following the modular testing pattern established for the introspection test suite.
"""

import unittest
from unittest.mock import Mock, PropertyMock, patch

from pptx.shapes.placeholder import LayoutPlaceholder
from pptx.slide import SlideLayout

from .mock_helpers import assert_basic_to_dict_structure


class MockLayoutPlaceholderElement:
    """Mock for layout placeholder XML element."""

    def __init__(self, ph_type="TITLE", ph_idx=0):
        self.ph_type = ph_type
        self.ph_idx = ph_idx
        self.shape_id = 42  # Add missing shape_id attribute
        self.shape_name = f"Mock Placeholder {ph_idx}"  # Fix missing shape_name
        self.has_ph_elm = True  # Fix missing has_ph_elm for is_placeholder property
        # Add geometric properties for BaseShape
        self.x = 914400  # Mock left position (1 inch in EMU)
        self.y = 685800  # Mock top position 
        self.cx = 6858000  # Mock width
        self.cy = 1371600  # Mock height
        self.rot = 0  # Mock rotation (no rotation)


class MockLayoutPlaceholderPart:
    """Mock for layout placeholder part."""

    def __init__(self, layout_name="Title Layout"):
        self.name = layout_name
        self.slide_master = MockSlideMaster()


class MockSlideMaster:
    """Mock for SlideMaster."""

    def __init__(self):
        self.name = "Default Master"
        self.placeholders = MockMasterPlaceholders()


class MockMasterPlaceholders:
    """Mock for master placeholders collection."""

    def __init__(self):
        self.master_placeholder = MockMasterPlaceholder()

    def get(self, ph_type, default=None):
        """Mock get method for master placeholders."""
        return self.master_placeholder if ph_type else default


class MockMasterPlaceholder:
    """Mock for MasterPlaceholder."""

    def __init__(self):
        self.placeholder_format = MockPlaceholderFormat()

    def to_dict(self, **kwargs):
        """Mock to_dict for master placeholder."""
        return {
            "_object_type": "MasterPlaceholder",
            "_identity": {"class_name": "MasterPlaceholder"},
            "properties": {"name": "Master Title Placeholder"}
        }


class MockPlaceholderFormat:
    """Mock for PlaceholderFormat."""

    def __init__(self, idx=0, ph_type="TITLE"):
        self.idx = idx
        self.type = ph_type


class MockSlideLayoutElement:
    """Mock for slide layout XML element."""

    def __init__(self, name="Title Layout"):
        self.cSld = MockCommonSlideData(name)
        self.spTree = Mock()  # Shape tree


class MockCommonSlideData:
    """Mock for CT_CommonSlideData (p:cSld element)."""

    def __init__(self, name="Title Layout"):
        self.name = name
        self.spTree = Mock()  # Shape tree


class MockSlideLayoutPart:
    """Mock for SlideLayoutPart."""

    def __init__(self, layout_name="Title Layout"):
        self.name = layout_name
        self.slide_master = MockSlideMaster()
        self.package = MockPackage()


class MockPackage:
    """Mock for Package."""

    def __init__(self):
        self.presentation_part = MockPresentationPart()


class MockPresentationPart:
    """Mock for PresentationPart."""

    def __init__(self):
        self.presentation = MockPresentation()


class MockPresentation:
    """Mock for Presentation."""

    def __init__(self):
        self.slides = MockSlides()


class MockSlides:
    """Mock for Slides collection."""

    def __init__(self):
        self.mock_slides = [MockSlide(256), MockSlide(257)]

    def __iter__(self):
        return iter(self.mock_slides)


class MockSlide:
    """Mock for Slide."""

    def __init__(self, slide_id):
        self.slide_id = slide_id
        self.slide_layout = None  # Will be set to the layout we're testing

    def to_dict(self, **kwargs):
        """Mock to_dict for slide."""
        return {
            "_object_type": "Slide",
            "_identity": {"slide_id": self.slide_id},
            "properties": {"slide_id": self.slide_id}
        }


class MockBackground:
    """Mock for _Background."""

    def __init__(self):
        self.fill = MockFillFormat()


class MockFillFormat:
    """Mock for FillFormat."""

    def to_dict(self, **kwargs):
        """Mock to_dict for fill format."""
        return {
            "_object_type": "FillFormat",
            "_identity": {"class_name": "FillFormat"},
            "properties": {"type": "SOLID"}
        }


class MockShapes:
    """Mock for shapes collection."""

    def __init__(self, has_placeholders=True):
        self.mock_shape = MockShape(is_placeholder=False)
        self.mock_placeholder = MockLayoutPlaceholder() if has_placeholders else None
        self._shapes = [self.mock_shape] + ([self.mock_placeholder] if self.mock_placeholder else [])

    def __iter__(self):
        return iter(self._shapes)

    def __len__(self):
        return len(self._shapes)


class MockShape:
    """Mock for Shape."""

    def __init__(self, is_placeholder=False):
        self.is_placeholder = is_placeholder
        self.name = "Mock Shape"

    def to_dict(self, **kwargs):
        """Mock to_dict for shape."""
        return {
            "_object_type": "Shape",
            "_identity": {"name": self.name},
            "properties": {"is_placeholder": self.is_placeholder}
        }


class MockLayoutPlaceholders:
    """Mock for LayoutPlaceholders collection."""

    def __init__(self):
        self.mock_placeholders = {0: MockLayoutPlaceholder(), 1: MockLayoutPlaceholder()}

    def items(self):
        """Mock items method."""
        return self.mock_placeholders.items()

    def keys(self):
        """Mock keys method."""
        return self.mock_placeholders.keys()

    def __len__(self):
        return len(self.mock_placeholders)
    
    def __iter__(self):
        """Fix missing __iter__ method for iteration support."""
        return iter(self.mock_placeholders.values())
        
    def __getitem__(self, key):
        """Support indexing access."""
        return self.mock_placeholders[key]


class MockLayoutPlaceholder:
    """Mock for LayoutPlaceholder (used in collections)."""

    def __init__(self, name="Mock Layout Placeholder"):
        self.name = name
        self.is_placeholder = True

    def to_dict(self, **kwargs):
        """Mock to_dict for layout placeholder."""
        return {
            "_object_type": "LayoutPlaceholder",
            "_identity": {"name": self.name},
            "properties": {"name": self.name, "is_placeholder": True}
        }


class TestLayoutPlaceholderIntrospection(unittest.TestCase):
    """Test LayoutPlaceholder.to_dict() functionality."""

    def setUp(self):
        """Set up test fixtures."""
        self.element = MockLayoutPlaceholderElement()
        self.part = MockLayoutPlaceholderPart()
        self.layout_placeholder = LayoutPlaceholder(self.element, self.part)

    def test_basic_to_dict_structure(self):
        """Test basic to_dict structure for LayoutPlaceholder."""
        result = self.layout_placeholder.to_dict()
        assert_basic_to_dict_structure(self, result, "LayoutPlaceholder")

    @patch.object(LayoutPlaceholder, 'placeholder_format', new_callable=PropertyMock)
    def test_to_dict_identity_with_placeholder_format(self, mock_pf):
        """Test _to_dict_identity includes placeholder format details."""
        mock_pf.return_value = MockPlaceholderFormat(idx=1, ph_type="BODY")
        
        result = self.layout_placeholder.to_dict()
        
        identity = result["_identity"]
        self.assertEqual(identity["description"], "Layout placeholder on slide layout")
        self.assertEqual(identity["placeholder_idx"], 1)
        self.assertIn("placeholder_type", identity)

    @patch.object(LayoutPlaceholder, '_base_placeholder', new_callable=PropertyMock)
    def test_to_dict_properties_with_master_placeholder(self, mock_base_ph):
        """Test _to_dict_properties includes master placeholder."""
        mock_base_ph.return_value = MockMasterPlaceholder()
        
        result = self.layout_placeholder.to_dict(max_depth=2)
        
        props = result["properties"]
        self.assertEqual(props["inherits_dimensions"], True)
        self.assertIn("master_placeholder", props)
        self.assertEqual(props["master_placeholder"]["_object_type"], "MasterPlaceholder")

    @patch.object(LayoutPlaceholder, '_base_placeholder', new_callable=PropertyMock)
    def test_to_dict_properties_no_master_placeholder(self, mock_base_ph):
        """Test _to_dict_properties when no master placeholder exists."""
        mock_base_ph.return_value = None
        
        result = self.layout_placeholder.to_dict()
        
        props = result["properties"]
        self.assertEqual(props["inherits_dimensions"], True)
        self.assertIsNone(props["master_placeholder"])

    def test_to_dict_relationships_includes_parent_layout(self):
        """Test _to_dict_relationships includes parent slide layout."""
        result = self.layout_placeholder.to_dict()
        
        rels = result["relationships"]
        self.assertIn("parent_slide_layout_ref", rels)

    @patch.object(LayoutPlaceholder, 'placeholder_format', new_callable=PropertyMock)
    @patch.object(LayoutPlaceholder, '_base_placeholder', new_callable=PropertyMock)
    def test_to_dict_llm_context(self, mock_base_ph, mock_pf):
        """Test _to_dict_llm_context provides descriptive information."""
        mock_pf.return_value = MockPlaceholderFormat(idx=0, ph_type="TITLE")
        mock_base_ph.return_value = MockMasterPlaceholder()
        
        result = self.layout_placeholder.to_dict()
        
        context = result["_llm_context"]
        self.assertIn("Layout placeholder #0", context["description"])
        self.assertIn("inherits from master placeholder", context["description"])
        self.assertIn("common_operations", context)

    def test_to_dict_max_depth_limits_recursion(self):
        """Test that max_depth parameter limits recursion."""
        result = self.layout_placeholder.to_dict(max_depth=1)
        
        # At depth 1, should have basic structure but limited nested content
        assert_basic_to_dict_structure(self, result, "LayoutPlaceholder")

    def test_to_dict_include_private_false(self):
        """Test that include_private=False excludes private attributes."""
        result = self.layout_placeholder.to_dict(include_private=False)
        
        # Should not include private properties
        props = result["properties"]
        private_keys = [k for k in props.keys() if k.startswith("_")]
        # Filter out expected private keys that are part of the introspection system
        unexpected_private = [k for k in private_keys if not k.startswith("_object_type")]
        self.assertEqual(len(unexpected_private), 0)

    def test_to_dict_format_for_llm_false(self):
        """Test that format_for_llm=False excludes LLM context."""
        result = self.layout_placeholder.to_dict(format_for_llm=False)
        
        self.assertNotIn("_llm_context", result)


class TestSlideLayoutIntrospection(unittest.TestCase):
    """Test SlideLayout.to_dict() functionality."""

    def setUp(self):
        """Set up test fixtures."""
        self.element = MockSlideLayoutElement()
        self.part = MockSlideLayoutPart()
        self.slide_layout = SlideLayout(self.element, self.part)

    def test_basic_to_dict_structure(self):
        """Test basic to_dict structure for SlideLayout."""
        result = self.slide_layout.to_dict()
        assert_basic_to_dict_structure(self, result, "SlideLayout")

    def test_to_dict_identity_includes_name(self):
        """Test _to_dict_identity includes layout name."""
        result = self.slide_layout.to_dict()
        
        identity = result["_identity"]
        self.assertEqual(identity["description"], "Slide Layout: 'Title Layout'")
        self.assertEqual(identity["name"], "Title Layout")

    @patch.object(SlideLayout, 'background', new_callable=PropertyMock)
    @patch.object(SlideLayout, 'shapes', new_callable=PropertyMock)
    @patch.object(SlideLayout, 'placeholders', new_callable=PropertyMock)
    def test_to_dict_properties_includes_collections(self, mock_placeholders, mock_shapes, mock_background):
        """Test _to_dict_properties includes all expected collections."""
        mock_background.return_value = MockBackground()
        mock_shapes.return_value = MockShapes()
        mock_placeholders.return_value = MockLayoutPlaceholders()
        
        result = self.slide_layout.to_dict(max_depth=2)
        
        props = result["properties"]
        self.assertIn("background_fill", props)
        self.assertIn("non_placeholder_shapes", props)
        self.assertIn("placeholders", props)
        
        # Verify background fill structure
        self.assertEqual(props["background_fill"]["_object_type"], "FillFormat")
        
        # Verify shapes collection
        self.assertIsInstance(props["non_placeholder_shapes"], list)
        
        # Verify placeholders collection
        self.assertIsInstance(props["placeholders"], list)

    @patch.object(SlideLayout, 'shapes', new_callable=PropertyMock)
    @patch.object(SlideLayout, 'placeholders', new_callable=PropertyMock)
    def test_to_dict_properties_expand_collections_false(self, mock_placeholders, mock_shapes):
        """Test _to_dict_properties with expand_collections=False."""
        mock_shapes.return_value = MockShapes()
        mock_placeholders.return_value = MockLayoutPlaceholders()
        
        result = self.slide_layout.to_dict(expand_collections=False)
        
        props = result["properties"]
        self.assertIn("_collection_summary", props["non_placeholder_shapes"])
        self.assertIn("_collection_summary", props["placeholders"])

    @patch.object(SlideLayout, 'slide_master', new_callable=PropertyMock)
    @patch.object(SlideLayout, 'used_by_slides', new_callable=PropertyMock)
    def test_to_dict_relationships_includes_master_and_slides(self, mock_used_by, mock_master):
        """Test _to_dict_relationships includes slide master and used by slides."""
        mock_master.return_value = MockSlideMaster()
        mock_used_by.return_value = [MockSlide(256), MockSlide(257)]
        
        result = self.slide_layout.to_dict()
        
        rels = result["relationships"]
        self.assertIn("slide_master_ref", rels)
        # Should have either summary or expanded list
        self.assertTrue("used_by_slides_summary" in rels or "used_by_slides" in rels)

    @patch.object(SlideLayout, 'slide_master', new_callable=PropertyMock)
    @patch.object(SlideLayout, 'used_by_slides', new_callable=PropertyMock)
    @patch.object(SlideLayout, 'shapes', new_callable=PropertyMock)
    @patch.object(SlideLayout, 'placeholders', new_callable=PropertyMock)
    def test_to_dict_llm_context(self, mock_placeholders, mock_shapes, mock_used_by, mock_master):
        """Test _to_dict_llm_context provides descriptive information."""
        mock_master.return_value = MockSlideMaster()
        mock_used_by.return_value = [MockSlide(256)]
        mock_shapes.return_value = MockShapes()
        mock_placeholders.return_value = MockLayoutPlaceholders()
        
        result = self.slide_layout.to_dict()
        
        context = result["_llm_context"]
        self.assertIn("Slide Layout 'Title Layout'", context["description"])
        self.assertIn("based on slide master 'Default Master'", context["description"])
        self.assertIn("Used by 1 slide", context["description"])
        self.assertIn("common_operations", context)

    def test_to_dict_max_depth_limits_recursion(self):
        """Test that max_depth parameter limits recursion."""
        result = self.slide_layout.to_dict(max_depth=1)
        
        # At depth 1, should have basic structure but limited nested content
        assert_basic_to_dict_structure(self, result, "SlideLayout")
        
        props = result["properties"]
        # Background fill should show depth exceeded
        self.assertEqual(props["background_fill"]["_depth_exceeded"], True)

    def test_to_dict_include_relationships_false(self):
        """Test that include_relationships=False excludes relationships."""
        result = self.slide_layout.to_dict(include_relationships=False)
        
        self.assertNotIn("relationships", result)

    def test_to_dict_format_for_llm_false(self):
        """Test that format_for_llm=False excludes LLM context."""
        result = self.slide_layout.to_dict(format_for_llm=False)
        
        self.assertNotIn("_llm_context", result)

    @patch.object(SlideLayout, 'used_by_slides', new_callable=PropertyMock)
    def test_to_dict_relationships_many_slides_summary(self, mock_used_by):
        """Test that many used_by_slides are summarized rather than expanded."""
        # Create many slides to test summarization
        many_slides = [MockSlide(i) for i in range(10)]
        mock_used_by.return_value = many_slides
        
        result = self.slide_layout.to_dict()
        
        rels = result["relationships"]
        self.assertIn("used_by_slides_summary", rels)
        self.assertEqual(rels["used_by_slides_summary"], "Used by 10 slide(s)")

    @unittest.skip("Complex dependency chain makes mocking difficult - covered by live tests")
    def test_to_dict_error_handling(self):
        """Test error handling in to_dict methods."""
        # This test would require extensive mocking of the entire dependency chain
        # It's better covered by the live test script where real objects are used
        pass


if __name__ == "__main__":
    unittest.main()