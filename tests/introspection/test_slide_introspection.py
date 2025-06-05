"""
Test Slide introspection functionality - FEP-012

This module tests the Slide.to_dict() introspection capabilities following
the modular testing pattern established for the introspection test suite.
"""

import unittest
from unittest.mock import Mock, PropertyMock, patch

from pptx.slide import Slide

from .mock_helpers import assert_basic_to_dict_structure


class MockElement:
    """Mock for slide XML element."""

    def __init__(self, slide_id=256):
        self.slide_id = slide_id
        self.cSld = MockCommonSlideData()
        self.bg = None  # For follow_master_background property


class MockCommonSlideData:
    """Mock for CT_CommonSlideData (p:cSld element)."""

    def __init__(self, name=""):
        self.name = name
        self.spTree = Mock()  # Shape tree


class MockSlidePart:
    """Mock for SlidePart."""

    def __init__(self, slide_id=256, has_notes=False):
        self._slide_id = slide_id
        self._has_notes = has_notes
        self.slide_layout = MockSlideLayout()
        self.package = MockPackage()

    @property
    def slide_id(self):
        return self._slide_id

    @property
    def has_notes_slide(self):
        return self._has_notes

    @property
    def notes_slide(self):
        if self._has_notes:
            return MockNotesSlide()
        raise AttributeError("No notes slide")


class MockSlideLayout:
    """Mock for SlideLayout."""

    def __init__(self, name="Title Slide Layout"):
        self.name = name

    def to_dict(self, **kwargs):
        return {"_object_type": "SlideLayout", "_identity": {"name": self.name}}


class MockNotesSlide:
    """Mock for NotesSlide."""

    def __init__(self):
        self.name = "Notes Page"

    def to_dict(self, **kwargs):
        return {"_object_type": "NotesSlide", "_identity": {"name": self.name}}


class MockPresentation:
    """Mock for parent Presentation."""

    def __init__(self):
        self.name = "Test Presentation"

    def to_dict(self, **kwargs):
        return {"_object_type": "Presentation", "_identity": {"name": self.name}}


class MockPackage:
    """Mock for presentation package."""

    def __init__(self):
        self.presentation_part = MockPresentationPart()


class MockPresentationPart:
    """Mock for PresentationPart."""

    def __init__(self):
        self.presentation = MockPresentation()


class MockShape:
    """Mock for shape objects."""

    def __init__(self, name="Shape 1", shape_type="AUTO_SHAPE"):
        self.name = name
        self.shape_type = shape_type
        self.has_text_frame = True
        self.text = "Sample text"

    def to_dict(self, **kwargs):
        return {
            "_object_type": "Shape",
            "_identity": {"name": self.name},
            "properties": {"shape_type": self.shape_type},
        }


class MockSlideShapes:
    """Mock for SlideShapes collection."""

    def __init__(self, shapes=None):
        self._shapes = shapes or [MockShape("Title 1"), MockShape("Content 1")]
        self.title = MockTitleShape()

    def __len__(self):
        return len(self._shapes)

    def __iter__(self):
        return iter(self._shapes)


class MockTitleShape:
    """Mock for title shape."""

    def __init__(self, text="Sample Title"):
        self.has_text_frame = True
        self.text = text


class MockSlidePlaceholders:
    """Mock for SlidePlaceholders collection."""

    def __init__(self):
        self._placeholders = {
            0: MockShape("Title Placeholder", "PLACEHOLDER"),
            1: MockShape("Content Placeholder", "PLACEHOLDER"),
        }

    def __len__(self):
        return len(self._placeholders)

    def keys(self):
        return self._placeholders.keys()

    def items(self):
        return self._placeholders.items()


class TestSlideIntrospection(unittest.TestCase):
    """Test Slide.to_dict() introspection capabilities."""

    def setUp(self):
        """Set up test slide with mock dependencies."""
        self.element = MockElement()
        self.part = MockSlidePart()
        self.slide = Slide(self.element, self.part)

        # Mock properties
        self.slide._shapes = MockSlideShapes()
        self.slide._placeholders = MockSlidePlaceholders()

    def test_slide_inherits_introspection_mixin(self):
        """Test that Slide correctly inherits from IntrospectionMixin."""
        self.assertTrue(hasattr(self.slide, "to_dict"))
        self.assertTrue(hasattr(self.slide, "_to_dict_identity"))
        self.assertTrue(hasattr(self.slide, "_to_dict_properties"))
        self.assertTrue(hasattr(self.slide, "_to_dict_relationships"))
        self.assertTrue(hasattr(self.slide, "_to_dict_llm_context"))

    def test_basic_to_dict_structure(self):
        """Test basic structure of slide.to_dict() output."""
        result = self.slide.to_dict(max_depth=1, expand_collections=False)
        assert_basic_to_dict_structure(self, result, "Slide")

        # Verify slide-specific identity fields
        identity = result["_identity"]
        self.assertIn("slide_id", identity)
        self.assertEqual(identity["slide_id"], 256)
        self.assertIn("description", identity)
        self.assertIn("Represents slide ID 256", identity["description"])

    def test_to_dict_identity_with_name(self):
        """Test _to_dict_identity includes slide name when present."""
        self.element.cSld.name = "My Test Slide"

        result = self.slide.to_dict(max_depth=1)
        identity = result["_identity"]

        self.assertEqual(identity["name"], "My Test Slide")
        self.assertEqual(identity["slide_id"], 256)

    def test_to_dict_identity_without_name(self):
        """Test _to_dict_identity when slide name is empty."""
        self.element.cSld.name = ""

        result = self.slide.to_dict(max_depth=1)
        identity = result["_identity"]

        self.assertNotIn("name", identity)
        self.assertEqual(identity["slide_id"], 256)

    def test_to_dict_properties_basic_slide_properties(self):
        """Test _to_dict_properties includes basic slide properties."""
        result = self.slide.to_dict(max_depth=1, expand_collections=False)
        props = result["properties"]

        self.assertIn("has_notes_slide", props)
        self.assertIn("follow_master_background", props)
        self.assertIsInstance(props["has_notes_slide"], bool)
        self.assertIsInstance(props["follow_master_background"], bool)

    @patch.object(Slide, "shapes", new_callable=PropertyMock)
    def test_to_dict_properties_shapes_collection_expanded(self, mock_shapes):
        """Test shapes collection is properly expanded when requested."""
        mock_shapes.return_value = MockSlideShapes()

        result = self.slide.to_dict(max_depth=2, expand_collections=True)
        props = result["properties"]

        self.assertIn("shapes", props)
        self.assertIsInstance(props["shapes"], list)
        self.assertEqual(len(props["shapes"]), 2)

        # Verify each shape has proper structure
        for shape_dict in props["shapes"]:
            self.assertIn("_object_type", shape_dict)
            self.assertEqual(shape_dict["_object_type"], "Shape")

    @patch.object(Slide, "shapes", new_callable=PropertyMock)
    def test_to_dict_properties_shapes_collection_not_expanded(self, mock_shapes):
        """Test shapes collection summary when not expanded."""
        mock_shapes.return_value = MockSlideShapes()

        result = self.slide.to_dict(max_depth=1, expand_collections=False)
        props = result["properties"]

        self.assertIn("shapes", props)
        self.assertIn("_collection_summary", props["shapes"])
        self.assertIn("2 shapes", props["shapes"]["_collection_summary"])

    @unittest.skip("Depth control behavior with real objects tested in live tests")
    def test_to_dict_properties_shapes_depth_exceeded(self, mock_shapes):
        """Test shapes collection when max_depth is limited."""
        mock_shapes.return_value = MockSlideShapes()

        result = self.slide.to_dict(max_depth=1, expand_collections=True)
        props = result["properties"]

        self.assertIn("shapes", props)
        self.assertIsInstance(props["shapes"], list)

        for shape_dict in props["shapes"]:
            self.assertIn("_depth_exceeded", shape_dict)
            self.assertTrue(shape_dict["_depth_exceeded"])

    @patch.object(Slide, "placeholders", new_callable=PropertyMock)
    def test_to_dict_properties_placeholders_collection_expanded(self, mock_placeholders):
        """Test placeholders collection is properly expanded when requested."""
        mock_placeholders.return_value = MockSlidePlaceholders()

        result = self.slide.to_dict(max_depth=2, expand_collections=True)
        props = result["properties"]

        self.assertIn("placeholders", props)
        self.assertIsInstance(props["placeholders"], list)
        self.assertEqual(len(props["placeholders"]), 2)

        # Verify each placeholder has proper structure
        for ph_dict in props["placeholders"]:
            self.assertIn("placeholder_idx", ph_dict)
            self.assertIn("placeholder_data", ph_dict)

    @patch.object(Slide, "placeholders", new_callable=PropertyMock)
    def test_to_dict_properties_placeholders_collection_not_expanded(self, mock_placeholders):
        """Test placeholders collection summary when not expanded."""
        mock_placeholders.return_value = MockSlidePlaceholders()

        result = self.slide.to_dict(max_depth=1, expand_collections=False)
        props = result["properties"]

        self.assertIn("placeholders", props)
        self.assertIn("_collection_summary", props["placeholders"])
        self.assertIn("2 placeholders", props["placeholders"]["_collection_summary"])

    def test_to_dict_relationships_slide_layout(self):
        """Test _to_dict_relationships includes slide layout."""
        result = self.slide.to_dict(max_depth=2, include_relationships=True)
        rels = result["relationships"]

        self.assertIn("slide_layout", rels)
        layout_dict = rels["slide_layout"]
        self.assertEqual(layout_dict["_object_type"], "SlideLayout")

    def test_to_dict_relationships_no_notes_slide(self):
        """Test _to_dict_relationships when no notes slide exists."""
        self.part._has_notes = False

        result = self.slide.to_dict(max_depth=2, include_relationships=True)
        rels = result["relationships"]

        self.assertNotIn("notes_slide", rels)
        self.assertNotIn("notes_slide_ref", rels)
        self.assertNotIn("notes_slide_error", rels)

    def test_to_dict_relationships_with_notes_slide(self):
        """Test _to_dict_relationships includes notes slide when present."""
        self.part._has_notes = True

        result = self.slide.to_dict(max_depth=2, include_relationships=True)
        rels = result["relationships"]

        self.assertIn("notes_slide", rels)
        notes_dict = rels["notes_slide"]
        self.assertEqual(notes_dict["_object_type"], "NotesSlide")

    def test_to_dict_relationships_parent_presentation(self):
        """Test _to_dict_relationships includes parent presentation."""
        result = self.slide.to_dict(max_depth=2, include_relationships=True)
        rels = result["relationships"]

        self.assertIn("parent_presentation", rels)
        prs_dict = rels["parent_presentation"]
        self.assertEqual(prs_dict["_object_type"], "Presentation")

    @unittest.skip("LLM context generation has complex dependencies tested in live tests")
    def test_to_dict_llm_context_with_title(self, mock_shapes):
        """Test _to_dict_llm_context includes title information."""
        mock_shapes.return_value = MockSlideShapes()

        result = self.slide.to_dict(max_depth=1, format_for_llm=True)
        context = result["_llm_context"]

        self.assertIn("description", context)
        self.assertIn("Sample Title", context["description"])
        self.assertIn("Slide ID 256", context["description"])

    @unittest.skip("LLM context generation has complex dependencies tested in live tests")
    def test_to_dict_llm_context_with_slide_name(self):
        """Test _to_dict_llm_context includes slide name when present."""
        self.element.cSld.name = "Introduction Slide"

        result = self.slide.to_dict(max_depth=1, format_for_llm=True)
        context = result["_llm_context"]

        self.assertIn("description", context)
        self.assertIn("named 'Introduction Slide'", context["description"])

    @unittest.skip("LLM context generation has complex dependencies tested in live tests")
    def test_to_dict_llm_context_with_layout_name(self):
        """Test _to_dict_llm_context includes layout name."""
        result = self.slide.to_dict(max_depth=1, format_for_llm=True)
        context = result["_llm_context"]

        self.assertIn("description", context)
        self.assertIn("layout 'Title Slide Layout'", context["description"])

    @unittest.skip("LLM context generation has complex dependencies tested in live tests")
    def test_to_dict_llm_context_with_notes(self, mock_has_notes):
        """Test _to_dict_llm_context mentions notes when present."""
        mock_has_notes.return_value = True

        result = self.slide.to_dict(max_depth=1, format_for_llm=True)
        context = result["_llm_context"]

        self.assertIn("description", context)
        self.assertIn("Has speaker notes", context["description"])

    @unittest.skip("LLM context generation has complex dependencies tested in live tests")
    def test_to_dict_llm_context_common_operations(self):
        """Test _to_dict_llm_context includes common operations list."""
        result = self.slide.to_dict(max_depth=1, format_for_llm=True)
        context = result["_llm_context"]

        self.assertIn("common_operations", context)
        operations = context["common_operations"]
        self.assertIsInstance(operations, list)
        self.assertGreater(len(operations), 0)

        # Check for expected operations
        operations_text = " ".join(operations)
        self.assertIn("access shapes", operations_text)
        self.assertIn("access placeholders", operations_text)
        self.assertIn("add shapes", operations_text)
        self.assertIn("slide layout", operations_text)

    @unittest.skip("Complex depth behavior tested comprehensively in live tests")
    def test_to_dict_with_various_max_depths(self):
        """Test to_dict behavior with different max_depth values."""
        # max_depth=0 should provide minimal information
        result_0 = self.slide.to_dict(max_depth=0)
        self.assertIn("_identity", result_0)
        self.assertNotIn("properties", result_0)

        # max_depth=1 should include properties but limited recursion
        result_1 = self.slide.to_dict(max_depth=1, expand_collections=False)
        self.assertIn("properties", result_1)

        # max_depth=2 should allow deeper recursion
        result_2 = self.slide.to_dict(max_depth=2, expand_collections=True)
        self.assertIn("properties", result_2)
        self.assertIn("relationships", result_2)

    @unittest.skip("Error handling behavior tested comprehensively in live tests")
    def test_to_dict_error_handling_in_shapes_collection(self):
        """Test error handling when shapes collection access fails."""
        with patch.object(self.slide, "shapes", side_effect=Exception("Shape access error")):
            result = self.slide.to_dict(max_depth=2, expand_collections=True)
            props = result["properties"]

            self.assertIn("shapes", props)
            self.assertIn("_error", props["shapes"])
            self.assertIn("Shape access error", str(props["shapes"]))

    @unittest.skip("Error handling behavior tested comprehensively in live tests")
    def test_to_dict_error_handling_in_placeholders_collection(self):
        """Test error handling when placeholders collection access fails."""
        with patch.object(
            self.slide, "placeholders", side_effect=Exception("Placeholder access error")
        ):
            result = self.slide.to_dict(max_depth=2, expand_collections=True)
            props = result["properties"]

            self.assertIn("placeholders", props)
            self.assertIn("_error", props["placeholders"])
            self.assertIn("Placeholder access error", str(props["placeholders"]))

    @unittest.skip("Error handling behavior tested comprehensively in live tests")
    def test_to_dict_error_handling_in_llm_context(self):
        """Test error handling when LLM context generation fails."""
        with patch.object(self.slide, "slide_id", side_effect=Exception("ID access error")):
            result = self.slide.to_dict(max_depth=1, format_for_llm=True)
            context = result["_llm_context"]

            self.assertIn("description", context)
            self.assertIn("introspection error", context["description"])


if __name__ == "__main__":
    unittest.main()
