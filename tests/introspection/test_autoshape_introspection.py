# tests/introspection/test_autoshape_introspection_simple.py

"""
Simplified AutoShape Introspection Tests

These tests focus on the core AutoShape introspection functionality with
minimal mocking to avoid complex property override issues.
"""

import unittest
from unittest.mock import Mock, patch
from pptx.shapes.autoshape import Shape
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE


class TestAutoShapeIntrospectionSimple(unittest.TestCase):
    """Simplified tests for Shape (AutoShape) introspection functionality."""

    def test_auto_shape_type_safe_access(self):
        """Test the _get_auto_shape_type_safely helper method."""
        # Mock CT_Shape that is an autoshape
        mock_ct_shape = Mock()
        mock_ct_shape.is_autoshape = True
        mock_ct_shape.prst = MSO_AUTO_SHAPE_TYPE.RECTANGLE
        mock_ct_shape.shape_id = 42
        mock_ct_shape.shape_name = "Test Rectangle"
        mock_ct_shape.has_ph_elm = False
        
        shape = Shape(mock_ct_shape, None)
        
        # Test safe access
        result = shape._get_auto_shape_type_safely()
        self.assertEqual(result, MSO_AUTO_SHAPE_TYPE.RECTANGLE)

    def test_auto_shape_type_safe_access_not_autoshape(self):
        """Test _get_auto_shape_type_safely with non-AutoShape."""
        # Mock CT_Shape that is not an autoshape
        mock_ct_shape = Mock()
        mock_ct_shape.is_autoshape = False
        mock_ct_shape.is_textbox = True
        mock_ct_shape.has_custom_geometry = False
        mock_ct_shape.shape_id = 42
        mock_ct_shape.shape_name = "Test Text Box"
        mock_ct_shape.has_ph_elm = False
        
        shape = Shape(mock_ct_shape, None)
        
        # Test safe access returns None for non-AutoShape
        result = shape._get_auto_shape_type_safely()
        self.assertIsNone(result)

    def test_shape_type_property(self):
        """Test that Shape.shape_type works correctly for AutoShapes."""
        # Test AutoShape
        mock_ct_shape = Mock()
        mock_ct_shape.is_autoshape = True
        mock_ct_shape.is_textbox = False
        mock_ct_shape.has_custom_geometry = False
        mock_ct_shape.shape_id = 42
        mock_ct_shape.shape_name = "Test Shape"
        mock_ct_shape.has_ph_elm = False
        
        shape = Shape(mock_ct_shape, None)
        self.assertEqual(shape.shape_type, MSO_SHAPE_TYPE.AUTO_SHAPE)
        
        # Test Text Box
        mock_ct_shape.is_autoshape = False
        mock_ct_shape.is_textbox = True
        self.assertEqual(shape.shape_type, MSO_SHAPE_TYPE.TEXT_BOX)
        
        # Test Freeform
        mock_ct_shape.is_autoshape = False
        mock_ct_shape.is_textbox = False
        mock_ct_shape.has_custom_geometry = True
        self.assertEqual(shape.shape_type, MSO_SHAPE_TYPE.FREEFORM)

    def test_identity_override_adds_auto_shape_type(self):
        """Test that _to_dict_identity adds auto_shape_type_details."""
        # Create mock AutoShape
        mock_ct_shape = Mock()
        mock_ct_shape.is_autoshape = True
        mock_ct_shape.prst = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        mock_ct_shape.shape_id = 123
        mock_ct_shape.shape_name = "Test Rounded Rectangle"
        mock_ct_shape.has_ph_elm = False
        
        shape = Shape(mock_ct_shape, None)
        
        # Test identity method directly
        identity = shape._to_dict_identity(
            _visited_ids=set(),
            max_depth=2,
            expand_collections=True,
            format_for_llm=True,
            include_private=False
        )
        
        # Should have auto_shape_type_details for AutoShape
        self.assertIn('auto_shape_type_details', identity)
        auto_shape_type = identity['auto_shape_type_details']
        self.assertEqual(auto_shape_type['_object_type'], 'MSO_AUTO_SHAPE_TYPE')
        self.assertEqual(auto_shape_type['name'], 'ROUNDED_RECTANGLE')

    def test_identity_override_no_auto_shape_type_for_textbox(self):
        """Test that _to_dict_identity doesn't add auto_shape_type_details for text box."""
        # Create mock Text Box
        mock_ct_shape = Mock()
        mock_ct_shape.is_autoshape = False
        mock_ct_shape.is_textbox = True
        mock_ct_shape.has_custom_geometry = False
        mock_ct_shape.shape_id = 456
        mock_ct_shape.shape_name = "Test Text Box"
        mock_ct_shape.has_ph_elm = False
        
        shape = Shape(mock_ct_shape, None)
        
        # Test identity method directly
        identity = shape._to_dict_identity(
            _visited_ids=set(),
            max_depth=2,
            expand_collections=True,
            format_for_llm=True,
            include_private=False
        )
        
        # Should NOT have auto_shape_type_details for text box
        self.assertNotIn('auto_shape_type_details', identity)

    @patch('pptx.shapes.autoshape.Shape.adjustments')
    def test_properties_override_includes_adjustments(self, mock_adjustments):
        """Test that _to_dict_properties includes adjustments."""
        # Setup mock adjustments
        mock_adjustments.__len__ = Mock(return_value=1)
        mock_adjustments.__iter__ = Mock(return_value=iter([0.25]))
        
        # Create shape
        mock_ct_shape = Mock()
        mock_ct_shape.shape_id = 789
        mock_ct_shape.shape_name = "Test Shape"
        shape = Shape(mock_ct_shape, None)
        
        # Test properties method with mocked dependencies
        with patch.object(shape, 'fill') as mock_fill, \
             patch.object(shape, 'line') as mock_line, \
             patch.object(shape, 'text_frame') as mock_text_frame:
            
            # Setup mocks to avoid deep introspection
            mock_fill.to_dict.return_value = {"_object_type": "FillFormat"}
            mock_line.to_dict.return_value = {"_object_type": "LineFormat"}
            mock_text_frame.to_dict.return_value = {"_object_type": "TextFrame"}
            
            props = shape._to_dict_properties(
                include_private=False,
                _visited_ids=set(),
                max_depth=2,
                expand_collections=True,
                format_for_llm=True
            )
            
            # Should include adjustments
            self.assertIn('adjustments', props)
            self.assertEqual(props['adjustments'], [0.25])

    @patch('pptx.shapes.autoshape.Shape.adjustments')
    def test_properties_override_empty_adjustments(self, mock_adjustments):
        """Test that _to_dict_properties excludes empty adjustments."""
        # Setup empty adjustments
        mock_adjustments.__len__ = Mock(return_value=0)
        mock_adjustments.__iter__ = Mock(return_value=iter([]))
        
        # Create shape
        mock_ct_shape = Mock()
        shape = Shape(mock_ct_shape, None)
        
        # Test properties method with mocked dependencies
        with patch.object(shape, 'fill') as mock_fill, \
             patch.object(shape, 'line') as mock_line, \
             patch.object(shape, 'text_frame') as mock_text_frame:
            
            mock_fill.to_dict.return_value = {"_object_type": "FillFormat"}
            mock_line.to_dict.return_value = {"_object_type": "LineFormat"}
            mock_text_frame.to_dict.return_value = {"_object_type": "TextFrame"}
            
            props = shape._to_dict_properties(
                include_private=False,
                _visited_ids=set(),
                max_depth=2,
                expand_collections=True,
                format_for_llm=True
            )
            
            # Should NOT include empty adjustments
            self.assertNotIn('adjustments', props)

    def test_llm_context_override_basic_functionality(self):
        """Test that _to_dict_llm_context provides AutoShape-specific context."""
        # Create mock AutoShape
        mock_ct_shape = Mock()
        mock_ct_shape.is_autoshape = True
        mock_ct_shape.prst = MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW
        mock_ct_shape.shape_id = 999
        mock_ct_shape.shape_name = "Test Arrow"
        mock_ct_shape.has_ph_elm = False
        
        shape = Shape(mock_ct_shape, None)
        
        # Mock dependencies to avoid complex interactions
        with patch.object(shape, 'adjustments') as mock_adjustments, \
             patch.object(shape, 'text_frame') as mock_text_frame:
            
            mock_adjustments.__len__ = Mock(return_value=2)
            mock_text_frame.text = "Arrow pointing right"
            
            context = shape._to_dict_llm_context(
                _visited_ids=set(),
                max_depth=2,
                expand_collections=True,
                format_for_llm=True,
                include_private=False
            )
            
            # Check enhanced description
            description = context['description']
            self.assertIn('AutoShape of type RIGHT_ARROW', description)
            self.assertIn('Test Arrow', description)
            self.assertIn('999', description)
            
            # Check enhanced summary
            summary = context['summary']
            self.assertIn('Contains text:', summary)
            self.assertIn('Has 2 adjustment handle(s)', summary)
            
            # Check AutoShape-specific operations
            operations = context['common_operations']
            self.assertIn('access/modify text_frame', operations)
            self.assertIn('change fill properties', operations)
            self.assertIn('change line properties', operations)
            self.assertIn('modify adjustment values', operations)


if __name__ == '__main__':
    unittest.main()