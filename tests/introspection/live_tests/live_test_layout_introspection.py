#!/usr/bin/env python3
"""
Live Test Script: SlideLayout and LayoutPlaceholder Introspection - FEP-016

This script validates the SlideLayout.to_dict() and LayoutPlaceholder.to_dict() 
implementation using real PowerPoint objects. It demonstrates the comprehensive 
introspection capabilities added in FEP-016.

Usage:
    python live_test_layout_introspection.py

Requirements:
    - python-pptx with FEP-016 implementation
    - Default template presentation (uses built-in template)
"""

import json
import sys
from pptx import Presentation


def print_section(title: str, content: str = None):
    """Print a formatted section header."""
    print(f"\n{'='*80}")
    print(f"  {title}")
    print(f"{'='*80}")
    if content:
        print(content)


def print_json_sample(data, max_lines=30):
    """Print a formatted JSON sample, truncated if too long."""
    json_str = json.dumps(data, indent=2, default=str)
    lines = json_str.split('\n')
    
    if len(lines) <= max_lines:
        print(json_str)
    else:
        print('\n'.join(lines[:max_lines]))
        print(f"... (truncated after {max_lines} lines, {len(lines)} total)")


def test_slide_layout_basic_introspection():
    """Test 1: Basic SlideLayout.to_dict() functionality."""
    print_section("Test 1: Basic SlideLayout Introspection")
    
    try:
        # Create a presentation and get the first slide layout
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        
        print(f"Testing SlideLayout: '{slide_layout.name}'")
        print(f"Layout has {len(slide_layout.placeholders)} placeholders")
        print(f"Layout has {len(slide_layout.shapes)} total shapes")
        
        # Test basic to_dict
        result = slide_layout.to_dict(max_depth=2)
        
        # Verify structure
        assert "_object_type" in result, "Missing _object_type"
        assert result["_object_type"] == "SlideLayout", f"Wrong object type: {result['_object_type']}"
        assert "_identity" in result, "Missing _identity"
        assert "properties" in result, "Missing properties"
        assert "relationships" in result, "Missing relationships"
        assert "_llm_context" in result, "Missing _llm_context"
        
        # Verify identity
        identity = result["_identity"]
        assert "description" in identity, "Missing description in identity"
        assert "name" in identity, "Missing name in identity"
        assert slide_layout.name in identity["description"], "Description doesn't include layout name"
        
        # Verify properties
        props = result["properties"]
        assert "background_fill" in props, "Missing background_fill"
        assert "non_placeholder_shapes" in props, "Missing non_placeholder_shapes"
        assert "placeholders" in props, "Missing placeholders"
        
        # Verify relationships
        rels = result["relationships"]
        assert "slide_master_ref" in rels or "slide_master" in rels, "Missing slide master relationship"
        assert "used_by_slides_summary" in rels or "used_by_slides" in rels, "Missing used_by_slides"
        
        # Verify LLM context
        context = result["_llm_context"]
        assert "description" in context, "Missing description in LLM context"
        assert "common_operations" in context, "Missing common_operations"
        
        print("✓ Basic structure verification passed")
        print("\nSample output (truncated):")
        print_json_sample(result)
        
        return True
        
    except Exception as e:
        print(f"✗ Test failed: {e}")
        return False


def test_layout_placeholder_introspection():
    """Test 2: LayoutPlaceholder.to_dict() functionality."""
    print_section("Test 2: LayoutPlaceholder Introspection")
    
    try:
        # Get a layout placeholder
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        
        print(f"Testing Layout: '{slide_layout.name}'")
        print(f"Available placeholders: {len(slide_layout.placeholders)}")
        
        # Get the first placeholder
        if len(slide_layout.placeholders) == 0:
            print("No placeholders found, using title layout instead")
            slide_layout = prs.slide_layouts[0]
        
        placeholder = next(iter(slide_layout.placeholders))
        placeholder_idx = placeholder.placeholder_format.idx
        
        print(f"Testing LayoutPlaceholder #{placeholder_idx}")
        print(f"Placeholder type: {placeholder.placeholder_format.type}")
        
        # Test placeholder introspection
        result = placeholder.to_dict(max_depth=2)
        
        # Verify structure
        assert "_object_type" in result, "Missing _object_type"
        assert result["_object_type"] == "LayoutPlaceholder", f"Wrong object type: {result['_object_type']}"
        assert "_identity" in result, "Missing _identity"
        assert "properties" in result, "Missing properties"
        assert "relationships" in result, "Missing relationships"
        assert "_llm_context" in result, "Missing _llm_context"
        
        # Verify identity includes placeholder-specific info
        identity = result["_identity"]
        assert "description" in identity, "Missing description"
        assert "Layout placeholder" in identity["description"], "Description doesn't mention layout placeholder"
        
        # Verify properties include layout placeholder specifics
        props = result["properties"]
        assert "inherits_dimensions" in props, "Missing inherits_dimensions"
        assert props["inherits_dimensions"] is True, "inherits_dimensions should be True"
        
        # Verify LLM context
        context = result["_llm_context"]
        assert "description" in context, "Missing description in LLM context"
        assert "Layout placeholder" in context["description"], "LLM context doesn't mention layout placeholder"
        
        print("✓ LayoutPlaceholder structure verification passed")
        print("\nSample output (truncated):")
        print_json_sample(result)
        
        return True
        
    except Exception as e:
        print(f"✗ Test failed: {e}")
        return False


def test_slide_layout_collections():
    """Test 3: SlideLayout collections introspection."""
    print_section("Test 3: SlideLayout Collections")
    
    try:
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content
        
        print(f"Testing collections on layout: '{slide_layout.name}'")
        
        # Test with expand_collections=True
        result_expanded = slide_layout.to_dict(max_depth=2, expand_collections=True)
        
        # Test with expand_collections=False
        result_collapsed = slide_layout.to_dict(max_depth=2, expand_collections=False)
        
        # Verify expanded collections
        props_expanded = result_expanded["properties"]
        if len(slide_layout.placeholders) > 0:
            assert isinstance(props_expanded["placeholders"], list), "Placeholders should be expanded to list"
            assert len(props_expanded["placeholders"]) > 0, "Should have expanded placeholders"
            
            # Check placeholder structure
            first_ph = props_expanded["placeholders"][0]
            assert "placeholder_idx" in first_ph, "Missing placeholder_idx"
            assert "placeholder_data" in first_ph, "Missing placeholder_data"
        
        # Verify collapsed collections
        props_collapsed = result_collapsed["properties"]
        assert "_collection_summary" in props_collapsed["placeholders"], "Should have collection summary"
        
        print("✓ Collections expansion/collapse works correctly")
        print(f"  Expanded placeholders: {len(props_expanded['placeholders']) if isinstance(props_expanded['placeholders'], list) else 'N/A'}")
        print(f"  Collapsed summary: {props_collapsed['placeholders']}")
        
        return True
        
    except Exception as e:
        print(f"✗ Test failed: {e}")
        return False


def test_slide_layout_relationships():
    """Test 4: SlideLayout relationships introspection."""
    print_section("Test 4: SlideLayout Relationships")
    
    try:
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        
        # Add a slide using this layout to test used_by_slides
        slide = prs.slides.add_slide(slide_layout)
        
        print(f"Testing relationships for layout: '{slide_layout.name}'")
        print(f"Used by {len(slide_layout.used_by_slides)} slides")
        
        # Test relationships
        result = slide_layout.to_dict(max_depth=1)
        
        rels = result["relationships"]
        
        # Verify slide master relationship
        assert "slide_master_ref" in rels or "slide_master" in rels, "Missing slide master relationship"
        
        # Verify used_by_slides relationship
        used_by_key = "used_by_slides" if "used_by_slides" in rels else "used_by_slides_summary"
        assert used_by_key in rels, "Missing used_by_slides relationship"
        
        if used_by_key == "used_by_slides_summary":
            assert "1 slide" in rels[used_by_key], "Should show 1 slide usage"
        
        print("✓ Relationship introspection works correctly")
        print(f"  Slide master: {type(rels.get('slide_master_ref', rels.get('slide_master', 'unknown')))}")
        print(f"  Used by slides: {rels[used_by_key]}")
        
        return True
        
    except Exception as e:
        print(f"✗ Test failed: {e}")
        return False


def test_parameter_variations():
    """Test 5: Various parameter combinations."""
    print_section("Test 5: Parameter Variations")
    
    try:
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        
        print("Testing different parameter combinations...")
        
        # Test max_depth variations
        result_depth_1 = slide_layout.to_dict(max_depth=1)
        result_depth_3 = slide_layout.to_dict(max_depth=3)
        
        # Verify depth limiting works
        props_1 = result_depth_1["properties"]
        props_3 = result_depth_3["properties"]
        
        # At depth 1, background_fill should be limited
        assert "_depth_exceeded" in props_1["background_fill"], "Depth 1 should limit background_fill"
        
        print("✓ max_depth parameter works correctly")
        
        # Test include_relationships=False
        result_no_rels = slide_layout.to_dict(include_relationships=False)
        assert "relationships" not in result_no_rels, "Should exclude relationships"
        
        print("✓ include_relationships parameter works correctly")
        
        # Test format_for_llm=False
        result_no_llm = slide_layout.to_dict(format_for_llm=False)
        assert "_llm_context" not in result_no_llm, "Should exclude LLM context"
        
        print("✓ format_for_llm parameter works correctly")
        
        return True
        
    except Exception as e:
        print(f"✗ Test failed: {e}")
        return False


def test_comprehensive_layout_analysis():
    """Test 6: Comprehensive layout analysis."""
    print_section("Test 6: Comprehensive Layout Analysis")
    
    try:
        prs = Presentation()
        
        print("Analyzing all available slide layouts...")
        
        total_layouts = len(prs.slide_layouts)
        successful_layouts = 0
        
        for i, layout in enumerate(prs.slide_layouts):
            try:
                print(f"\nLayout {i+1}/{total_layouts}: '{layout.name}'")
                print(f"  Placeholders: {len(layout.placeholders)}")
                print(f"  Total shapes: {len(layout.shapes)}")
                print(f"  Used by slides: {len(layout.used_by_slides)}")
                
                # Test introspection
                result = layout.to_dict(max_depth=1)
                
                # Basic validation
                assert result["_object_type"] == "SlideLayout"
                assert layout.name in result["_identity"]["description"]
                
                successful_layouts += 1
                print("  ✓ Introspection successful")
                
            except Exception as e:
                print(f"  ✗ Introspection failed: {e}")
        
        print(f"\nSummary: {successful_layouts}/{total_layouts} layouts successfully introspected")
        
        if successful_layouts == total_layouts:
            print("✓ All layouts successfully introspected")
            return True
        else:
            print(f"✗ {total_layouts - successful_layouts} layouts failed")
            return False
        
    except Exception as e:
        print(f"✗ Test failed: {e}")
        return False


def main():
    """Run all live tests for SlideLayout and LayoutPlaceholder introspection."""
    print_section("FEP-016 Live Test: SlideLayout & LayoutPlaceholder Introspection")
    
    tests = [
        test_slide_layout_basic_introspection,
        test_layout_placeholder_introspection,
        test_slide_layout_collections,
        test_slide_layout_relationships,
        test_parameter_variations,
        test_comprehensive_layout_analysis,
    ]
    
    passed = 0
    total = len(tests)
    
    for test_func in tests:
        try:
            if test_func():
                passed += 1
        except Exception as e:
            print(f"✗ Test {test_func.__name__} crashed: {e}")
    
    print_section("Final Results")
    print(f"Tests passed: {passed}/{total}")
    
    if passed == total:
        print("🎉 All tests passed! FEP-016 implementation is working correctly.")
        return 0
    else:
        print(f"❌ {total - passed} test(s) failed. Please review the implementation.")
        return 1


if __name__ == "__main__":
    sys.exit(main())