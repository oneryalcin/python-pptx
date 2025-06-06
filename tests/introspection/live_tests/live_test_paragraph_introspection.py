#!/usr/bin/env python3
"""
Live test script for FEP-010: _Paragraph.to_dict() functionality

This script tests the _Paragraph introspection implementation with real python-pptx objects.
Run this script to validate that the introspection works correctly with actual presentations.

Usage:
    source venv/bin/activate && python live_test_paragraph_introspection.py
"""

import json
import pptx
from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Pt

def test_basic_paragraph_introspection():
    """Test basic paragraph introspection functionality."""
    print("Testing basic paragraph introspection...")
    
    # Create a simple presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Get title textframe
    title_shape = slide.shapes.title
    text_frame = title_shape.text_frame
    
    # Clear and add paragraph
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = "Test paragraph for introspection"
    
    # Test basic to_dict
    try:
        result = p.to_dict()
        print("✓ Basic to_dict() successful")
        print(f"  Object type: {result['_object_type']}")
        print(f"  Text: {result['properties']['text']}")
        print(f"  Alignment: {result['properties']['alignment']}")
        print(f"  Level: {result['properties']['level']}")
        return True
    except Exception as e:
        print(f"✗ Basic to_dict() failed: {e}")
        return False

def test_formatted_paragraph_introspection():
    """Test paragraph introspection with formatting."""
    print("\nTesting formatted paragraph introspection...")
    
    try:
        # Create a simple presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        # Get content textframe
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        
        # Clear and add formatted paragraph
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = "Formatted paragraph with multiple properties"
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        p.level = 1
        p.space_before = Pt(6)
        p.space_after = Pt(12)
        p.line_spacing = 1.5
        
        # Format font
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.bold = True
        
        # Test formatted to_dict
        result = p.to_dict()
        print("✓ Formatted paragraph to_dict() successful")
        print(f"  Alignment: {result['properties']['alignment']}")
        print(f"  Level: {result['properties']['level']}")
        print(f"  Line spacing: {result['properties']['line_spacing']}")
        print(f"  Font object type: {result['properties']['font']['_object_type']}")
        
        # Test LLM context
        if '_llm_context' in result:
            print(f"  LLM Summary: {result['_llm_context']['summary'][:100]}...")
            
        return True
    except Exception as e:
        print(f"✗ Formatted paragraph to_dict() failed: {e}")
        return False

def test_paragraph_with_multiple_runs():
    """Test paragraph introspection with multiple runs."""
    print("\nTesting paragraph with multiple runs...")
    
    try:
        # Create a simple presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        # Get content textframe
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        
        # Clear and add paragraph with multiple runs
        text_frame.clear()
        p = text_frame.paragraphs[0]
        
        # Add first run
        run1 = p.add_run()
        run1.text = "First run with "
        run1.font.bold = True
        
        # Add second run
        run2 = p.add_run()
        run2.text = "second run in "
        run2.font.italic = True
        
        # Add third run  
        run3 = p.add_run()
        run3.text = "same paragraph"
        run3.font.name = "Times New Roman"
        
        # Test multi-run to_dict
        result = p.to_dict(expand_collections=True)
        print("✓ Multi-run paragraph to_dict() successful")
        print(f"  Number of runs: {len(result['properties']['runs'])}")
        
        for i, run_dict in enumerate(result['properties']['runs']):
            if 'properties' in run_dict:
                run_text = run_dict['properties']['text']
                print(f"    Run {i+1}: '{run_text}'")
            
        return True
    except Exception as e:
        print(f"✗ Multi-run paragraph to_dict() failed: {e}")
        return False

def test_max_depth_control():
    """Test max_depth parameter controls recursion."""
    print("\nTesting max_depth control...")
    
    try:
        # Create a simple presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        # Get title textframe
        title_shape = slide.shapes.title
        text_frame = title_shape.text_frame
        
        # Clear and add paragraph
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = "Depth control test"
        
        # Test with max_depth=1 (should truncate font and runs)
        result_shallow = p.to_dict(max_depth=1)
        print("✓ max_depth=1 successful")
        
        font_truncated = result_shallow['properties']['font'].get('_depth_exceeded', False)
        print(f"  Font truncated: {font_truncated}")
        
        # Test with max_depth=3 (should expand everything)
        result_deep = p.to_dict(max_depth=3)
        print("✓ max_depth=3 successful")
        
        font_expanded = 'properties' in result_deep['properties']['font']
        print(f"  Font expanded: {font_expanded}")
        
        return True
    except Exception as e:
        print(f"✗ max_depth control failed: {e}")
        return False

def test_error_handling():
    """Test error handling in introspection."""
    print("\nTesting error handling...")
    
    try:
        # Create a simple presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        # Get title textframe
        title_shape = slide.shapes.title
        text_frame = title_shape.text_frame
        
        # Clear and add paragraph
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = "Error handling test"
        
        # Test to_dict with various parameters
        result = p.to_dict(
            include_relationships=True,
            max_depth=2,
            include_private=False,
            expand_collections=True,
            format_for_llm=True
        )
        
        print("✓ Error handling test successful")
        print(f"  Structure preserved: {'_object_type' in result}")
        print(f"  Properties exist: {'properties' in result}")
        print(f"  Relationships exist: {'relationships' in result}")
        print(f"  LLM context exists: {'_llm_context' in result}")
        
        return True
    except Exception as e:
        print(f"✗ Error handling test failed: {e}")
        return False

def demonstrate_full_output():
    """Demonstrate complete paragraph introspection output."""
    print("\nDemonstrating full introspection output...")
    
    try:
        # Create a rich paragraph example
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        # Get content textframe
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        
        # Clear and create rich paragraph
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = "This is a richly formatted paragraph for demonstration purposes."
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        p.level = 1
        p.space_before = Pt(12)
        p.line_spacing = 1.2
        
        # Format font
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.bold = True
        
        # Get full introspection
        result = p.to_dict(
            include_relationships=True,
            max_depth=3,
            expand_collections=True,
            format_for_llm=True
        )
        
        print("✓ Full introspection successful")
        print("\nSample output (formatted):")
        
        # Pretty print key sections
        print(f"Identity: {result['_identity']['description']}")
        print(f"Text: {result['properties']['text']}")
        print(f"Alignment: {result['properties']['alignment']}")
        print(f"Level: {result['properties']['level']}")
        print(f"LLM Summary: {result['_llm_context']['summary']}")
        
        return True
    except Exception as e:
        print(f"✗ Full demonstration failed: {e}")
        return False

def main():
    """Run all live tests for paragraph introspection."""
    print("FEP-010: _Paragraph.to_dict() Live Tests")
    print("="*50)
    
    tests = [
        test_basic_paragraph_introspection,
        test_formatted_paragraph_introspection,
        test_paragraph_with_multiple_runs,
        test_max_depth_control,
        test_error_handling,
        demonstrate_full_output,
    ]
    
    passed = 0
    total = len(tests)
    
    for test_func in tests:
        try:
            if test_func():
                passed += 1
        except Exception as e:
            print(f"✗ Test {test_func.__name__} crashed: {e}")
    
    print("\n" + "="*50)
    print(f"Results: {passed}/{total} tests passed")
    
    if passed == total:
        print("🎉 All tests passed! FEP-010 implementation is working correctly.")
        return 0
    else:
        print("❌ Some tests failed. Please review the implementation.")
        return 1

if __name__ == "__main__":
    exit(main())