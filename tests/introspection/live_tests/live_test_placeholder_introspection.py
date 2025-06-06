#!/usr/bin/env python3
"""
FEP-014: PlaceholderFormat.to_dict() & Richer Placeholder Details - Live Test Script

This script validates the implementation of FEP-014 using real python-pptx objects
to demonstrate _PlaceholderFormat introspection capabilities and enhanced BaseShape
placeholder handling.

Key Features Tested:
1. _PlaceholderFormat.to_dict() basic functionality
2. Enhanced BaseShape placeholder details using to_dict()
3. Various placeholder types (TITLE, BODY, PICTURE, etc.)
4. LLM context generation for placeholder descriptions
5. Integration between PlaceholderFormat and BaseShape introspection

Usage:
    python live_test_placeholder_introspection.py
"""

import sys
import json
from pathlib import Path

# Add src to path for imports
sys.path.insert(0, str(Path(__file__).parent / "src"))

try:
    import pptx
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    from pptx.util import Inches
except ImportError as e:
    print(f"Import error: {e}")
    print("Make sure you've installed the package: pip install -e .")
    sys.exit(1)


def test_placeholder_format_introspection():
    """Test PlaceholderFormat.to_dict() functionality."""
    print("=" * 80)
    print("Testing PlaceholderFormat.to_dict() Functionality")
    print("=" * 80)
    
    # Create a presentation and access placeholders
    prs = Presentation()
    
    # Get slide with default layout (has TITLE and BODY placeholders)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    print(f"Created slide with {len(slide.shapes)} shapes")
    
    placeholder_count = 0
    for i, shape in enumerate(slide.shapes):
        if shape.is_placeholder:
            placeholder_count += 1
            print(f"\n--- Placeholder {placeholder_count}: Shape {i} ---")
            
            # Get placeholder format
            ph_format = shape.placeholder_format
            print(f"Placeholder Type: {ph_format.type}")
            print(f"Placeholder Index: {ph_format.idx}")
            
            # Test PlaceholderFormat.to_dict()
            print("\n1. PlaceholderFormat.to_dict() Output:")
            ph_dict = ph_format.to_dict(max_depth=3, format_for_llm=True)
            
            # Print formatted JSON
            print(json.dumps(ph_dict, indent=2, default=str)[:1000] + "...")
            
            # Verify expected structure
            assert "_object_type" in ph_dict
            assert ph_dict["_object_type"] == "_PlaceholderFormat"
            assert "properties" in ph_dict
            assert "idx" in ph_dict["properties"]
            assert "type" in ph_dict["properties"]
            assert "_llm_context" in ph_dict
            
            print("✓ PlaceholderFormat.to_dict() structure validated")
            
            # Test LLM context
            llm_context = ph_dict["_llm_context"]
            print(f"\n2. LLM Context Description: {llm_context.get('description', 'N/A')}")
            print(f"   Common Operations: {len(llm_context.get('common_operations', []))} operations")
            
            # Verify type serialization
            type_dict = ph_dict["properties"]["type"]
            if isinstance(type_dict, dict):
                print(f"3. Type Serialization: {type_dict.get('name', 'N/A')} (value: {type_dict.get('value', 'N/A')})")
            else:
                print(f"3. Type Serialization: {type_dict}")
    
    print(f"\nFound and tested {placeholder_count} placeholders")
    return placeholder_count > 0


def test_enhanced_baseshape_placeholder_details():
    """Test enhanced BaseShape placeholder details using PlaceholderFormat.to_dict()."""
    print("\n" + "=" * 80)
    print("Testing Enhanced BaseShape Placeholder Details")
    print("=" * 80)
    
    # Create presentation with various placeholder types
    prs = Presentation()
    
    # Test with different slide layouts to get different placeholder types
    layout_names = ["Title Slide", "Title and Content", "Section Header", "Two Content", "Comparison"]
    
    for layout_idx in range(min(5, len(prs.slide_layouts))):
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        
        print(f"\n--- Layout {layout_idx}: {len(slide.shapes)} shapes ---")
        
        for i, shape in enumerate(slide.shapes):
            if shape.is_placeholder:
                print(f"\nShape {i} - Placeholder Analysis:")
                
                # Test BaseShape.to_dict() placeholder handling
                shape_dict = shape.to_dict(max_depth=3, include_relationships=False, format_for_llm=True)
                
                # Check identity section
                identity = shape_dict["_identity"]
                assert identity["is_placeholder"] is True
                assert "placeholder_details" in identity
                
                # Verify placeholder_details is from PlaceholderFormat.to_dict()
                placeholder_details = identity["placeholder_details"]
                assert isinstance(placeholder_details, dict)
                assert "_object_type" in placeholder_details
                assert placeholder_details["_object_type"] == "_PlaceholderFormat"
                
                print(f"   Type: {placeholder_details['properties']['type']['name']}")
                print(f"   Index: {placeholder_details['properties']['idx']}")
                print(f"   Description: {placeholder_details['_identity']['description']}")
                
                # Test LLM context integration
                shape_llm_context = shape_dict["_llm_context"]
                description = shape_llm_context.get("description", "")
                print(f"   Shape LLM Description: {description[:100]}...")
                
                # Verify placeholder information is in shape description
                assert "placeholder" in description.lower()
                
                print("✓ Enhanced placeholder details validated")
                
                # Break after first placeholder to avoid too much output
                break
    
    print("\n✓ Enhanced BaseShape placeholder details tested successfully")
    return True


def test_various_placeholder_types():
    """Test introspection with various placeholder types."""
    print("\n" + "=" * 80)
    print("Testing Various Placeholder Types")
    print("=" * 80)
    
    # Create a presentation
    prs = Presentation()
    
    # Test different layouts to get different placeholder types
    placeholder_types_found = set()
    
    for layout_idx in range(min(8, len(prs.slide_layouts))):
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                placeholder_types_found.add(ph_type.name)
                
                print(f"\nPlaceholder Type: {ph_type.name}")
                
                # Test PlaceholderFormat introspection
                ph_dict = shape.placeholder_format.to_dict(max_depth=2)
                
                # Verify type-specific details
                type_name = ph_dict["properties"]["type"]["name"]
                assert type_name == ph_type.name
                
                # Check description includes type
                description = ph_dict["_identity"]["description"]
                assert ph_type.name in description
                
                print(f"   Index: {ph_dict['properties']['idx']}")
                print(f"   Description: {description}")
                print(f"✓ {ph_type.name} placeholder introspection validated")
    
    print(f"\nFound placeholder types: {sorted(placeholder_types_found)}")
    print(f"Total unique types tested: {len(placeholder_types_found)}")
    
    return len(placeholder_types_found) > 0


def test_error_handling():
    """Test error handling in placeholder introspection."""
    print("\n" + "=" * 80)
    print("Testing Error Handling")
    print("=" * 80)
    
    # Create presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Test with non-placeholder shape
    non_placeholder_shapes = [shape for shape in slide.shapes if not shape.is_placeholder]
    
    if non_placeholder_shapes:
        shape = non_placeholder_shapes[0]
        print(f"Testing non-placeholder shape: {shape.name}")
        
        # This should not have placeholder details
        shape_dict = shape.to_dict(max_depth=2, include_relationships=False)
        identity = shape_dict["_identity"]
        
        assert identity["is_placeholder"] is False
        assert "placeholder_details" not in identity
        
        print("✓ Non-placeholder shape correctly handled")
    else:
        print("No non-placeholder shapes found to test")
    
    # Test with placeholder shape
    placeholder_shapes = [shape for shape in slide.shapes if shape.is_placeholder]
    
    if placeholder_shapes:
        shape = placeholder_shapes[0]
        print(f"\nTesting placeholder shape: {shape.name}")
        
        # This should have proper placeholder details
        shape_dict = shape.to_dict(max_depth=2, include_relationships=False)
        identity = shape_dict["_identity"]
        
        assert identity["is_placeholder"] is True
        assert "placeholder_details" in identity
        
        # Verify the placeholder details structure
        placeholder_details = identity["placeholder_details"]
        assert isinstance(placeholder_details, dict)
        assert "_object_type" in placeholder_details
        
        print("✓ Placeholder shape correctly handled")
    
    return True


def test_performance_and_depth():
    """Test performance and depth handling."""
    print("\n" + "=" * 80)
    print("Testing Performance and Depth Handling")
    print("=" * 80)
    
    # Create presentation with multiple slides
    prs = Presentation()
    
    # Add multiple slides with placeholders
    for i in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    print(f"Created presentation with {len(prs.slides)} slides")
    
    # Test with different max_depth values
    for depth in [1, 2, 3]:
        print(f"\n--- Testing max_depth={depth} ---")
        
        slide = prs.slides[0]
        placeholder_shapes = [shape for shape in slide.shapes if shape.is_placeholder]
        
        if placeholder_shapes:
            shape = placeholder_shapes[0]
            
            # Test PlaceholderFormat with different depths
            ph_dict = shape.placeholder_format.to_dict(max_depth=depth)
            print(f"PlaceholderFormat depth {depth}: {len(str(ph_dict))} chars")
            
            # Test BaseShape with different depths
            shape_dict = shape.to_dict(max_depth=depth, include_relationships=False)
            print(f"BaseShape depth {depth}: {len(str(shape_dict))} chars")
            
            # Verify basic structure is always present
            assert "_object_type" in ph_dict
            assert "properties" in ph_dict
            
            print(f"✓ Depth {depth} handling validated")
    
    return True


def main():
    """Run all placeholder introspection tests."""
    print("FEP-014: PlaceholderFormat.to_dict() & Enhanced BaseShape - Live Test")
    print("=" * 80)
    
    try:
        # Test 1: Basic PlaceholderFormat introspection
        test1_passed = test_placeholder_format_introspection()
        
        # Test 2: Enhanced BaseShape placeholder details
        test2_passed = test_enhanced_baseshape_placeholder_details()
        
        # Test 3: Various placeholder types
        test3_passed = test_various_placeholder_types()
        
        # Test 4: Error handling
        test4_passed = test_error_handling()
        
        # Test 5: Performance and depth
        test5_passed = test_performance_and_depth()
        
        # Summary
        print("\n" + "=" * 80)
        print("TEST SUMMARY")
        print("=" * 80)
        print(f"1. PlaceholderFormat.to_dict() Basic Functionality: {'✓ PASS' if test1_passed else '✗ FAIL'}")
        print(f"2. Enhanced BaseShape Placeholder Details: {'✓ PASS' if test2_passed else '✗ FAIL'}")
        print(f"3. Various Placeholder Types: {'✓ PASS' if test3_passed else '✗ FAIL'}")
        print(f"4. Error Handling: {'✓ PASS' if test4_passed else '✗ FAIL'}")
        print(f"5. Performance and Depth Handling: {'✓ PASS' if test5_passed else '✗ FAIL'}")
        
        all_passed = all([test1_passed, test2_passed, test3_passed, test4_passed, test5_passed])
        
        if all_passed:
            print("\n🎉 ALL TESTS PASSED! FEP-014 implementation is working correctly.")
            print("\nKey achievements:")
            print("✓ _PlaceholderFormat now inherits from IntrospectionMixin")
            print("✓ _PlaceholderFormat.to_dict() provides comprehensive placeholder details")
            print("✓ BaseShape.to_dict() uses PlaceholderFormat.to_dict() for richer placeholder information")
            print("✓ LLM context generation for placeholder descriptions")
            print("✓ Proper error handling and depth management")
            return 0
        else:
            print("\n❌ SOME TESTS FAILED! Please check the implementation.")
            return 1
            
    except Exception as e:
        print(f"\n❌ ERROR during testing: {e}")
        import traceback
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    sys.exit(main())