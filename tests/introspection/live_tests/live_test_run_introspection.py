#!/usr/bin/env python3
"""
Live Testing Script for FEP-009: _Run.to_dict() Introspection

This script demonstrates the _Run introspection functionality by creating
actual PowerPoint objects and testing the to_dict() output.

This script should be deleted before creating the PR, but can be referenced
in PR comments for engineers to review and run manually.
"""

import json
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_TEXT_UNDERLINE_TYPE


def test_basic_run_introspection():
    """Test basic run introspection functionality."""
    print("\n" + "="*60)
    print("TEST 1: Basic Run Introspection")
    print("="*60)
    
    # Create a presentation with text
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    title = slide.shapes.title
    title.text = "Test Slide"
    
    # Add text to content placeholder
    content = slide.shapes.placeholders[1]
    content.text = "This is a sample text with multiple runs."
    
    # Get the first paragraph and its first run
    paragraph = content.text_frame.paragraphs[0]
    run = paragraph.runs[0]
    
    # Test to_dict functionality
    result = run.to_dict()
    
    print(f"Run text: '{run.text}'")
    print(f"Object type: {result['_object_type']}")
    print(f"Description: {result['_identity']['description']}")
    print(f"Text property: {result['properties']['text']}")
    print(f"Hyperlink address: {result['properties']['hyperlink_address']}")
    print(f"LLM Summary: {result['_llm_context']['summary']}")
    
    print("\nFull JSON output:")
    print(json.dumps(result, indent=2)[:500] + "..." if len(json.dumps(result)) > 500 else json.dumps(result, indent=2))


def test_run_with_formatting():
    """Test run with various font formatting."""
    print("\n" + "="*60)
    print("TEST 2: Run with Font Formatting")
    print("="*60)
    
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    content = slide.shapes.placeholders[1]
    content.text = "Formatted text"
    
    # Get run and apply formatting
    run = content.text_frame.paragraphs[0].runs[0]
    
    # Apply various formatting
    run.font.name = "Times New Roman"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.italic = True
    run.font.underline = MSO_TEXT_UNDERLINE_TYPE.SINGLE_LINE
    
    result = run.to_dict()
    
    print(f"Run text: '{run.text}'")
    print(f"LLM Summary: {result['_llm_context']['summary']}")
    
    # Check font properties
    font_props = result['properties']['font']['properties']
    print(f"Font name: {font_props['name']}")
    print(f"Font size: {font_props['size']['pt']}pt")
    print(f"Bold: {font_props['bold']}")
    print(f"Italic: {font_props['italic']}")
    
    print("\nFont section of JSON:")
    print(json.dumps(result['properties']['font'], indent=2))


def test_run_with_hyperlink():
    """Test run with hyperlink."""
    print("\n" + "="*60)
    print("TEST 3: Run with Hyperlink")
    print("="*60)
    
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    content = slide.shapes.placeholders[1]
    content.text = "Click here for more info"
    
    # Get run and add hyperlink
    run = content.text_frame.paragraphs[0].runs[0]
    run.hyperlink.address = "https://python-pptx.readthedocs.io"
    
    result = run.to_dict()
    
    print(f"Run text: '{run.text}'")
    print(f"Hyperlink address: {result['properties']['hyperlink_address']}")
    print(f"LLM Summary: {result['_llm_context']['summary']}")
    
    # Check relationships
    if 'relationships' in result and 'hyperlink' in result['relationships']:
        hyperlink_rel = result['relationships']['hyperlink']
        print(f"Relationship rId: {hyperlink_rel['rId']}")
        print(f"Target URL: {hyperlink_rel['target_url']}")
        print(f"Is external: {hyperlink_rel['is_external']}")
    
    print("\nRelationships section:")
    print(json.dumps(result.get('relationships', {}), indent=2))


def test_multiple_runs():
    """Test multiple runs in a paragraph."""
    print("\n" + "="*60)
    print("TEST 4: Multiple Runs in Paragraph")
    print("="*60)
    
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    content = slide.shapes.placeholders[1]
    
    # Clear existing content and add multiple runs
    p = content.text_frame.paragraphs[0]
    p.clear()  # Clear content properly
    
    # Add first run
    run1 = p.add_run()
    run1.text = "Normal text, "
    
    # Add second run with formatting
    run2 = p.add_run()
    run2.text = "bold text, "
    run2.font.bold = True
    
    # Add third run with hyperlink
    run3 = p.add_run()
    run3.text = "hyperlinked text"
    run3.hyperlink.address = "https://example.com"
    
    print("Testing each run:")
    for i, run in enumerate(p.runs):
        print(f"\nRun {i+1}:")
        result = run.to_dict()
        print(f"  Text: '{run.text}'")
        print(f"  Summary: {result['_llm_context']['summary']}")
        if result['properties']['hyperlink_address']:
            print(f"  Hyperlink: {result['properties']['hyperlink_address']}")


def test_max_depth_control():
    """Test max_depth parameter control."""
    print("\n" + "="*60)
    print("TEST 5: Max Depth Control")
    print("="*60)
    
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    content = slide.shapes.placeholders[1]
    content.text = "Test depth control"
    
    run = content.text_frame.paragraphs[0].runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(12)
    run.font.bold = True
    
    # Test with different max_depth values
    for depth in [1, 2, 3]:
        print(f"\nMax depth = {depth}:")
        result = run.to_dict(max_depth=depth)
        
        if depth == 1:
            print(f"  Font object: {result['properties']['font']}")
        else:
            font_props = result['properties']['font']['properties']
            print(f"  Font name: {font_props['name']}")
            print(f"  Font size: {font_props['size']}")


def test_format_for_llm_flag():
    """Test format_for_llm flag."""
    print("\n" + "="*60)
    print("TEST 6: format_for_llm Flag")
    print("="*60)
    
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    content = slide.shapes.placeholders[1]
    content.text = "Testing LLM flag"
    
    run = content.text_frame.paragraphs[0].runs[0]
    
    # Test with format_for_llm=True
    result_llm = run.to_dict(format_for_llm=True)
    print("With format_for_llm=True:")
    print(f"  Has _llm_context: {'_llm_context' in result_llm}")
    if '_llm_context' in result_llm:
        print(f"  LLM Summary: {result_llm['_llm_context']['summary']}")
    
    # Test with format_for_llm=False
    result_no_llm = run.to_dict(format_for_llm=False)
    print("\nWith format_for_llm=False:")
    print(f"  Has _llm_context: {'_llm_context' in result_no_llm}")
    print(f"  Keys present: {list(result_no_llm.keys())}")


def main():
    """Run all tests."""
    print("FEP-009: _Run.to_dict() Introspection - Live Testing")
    print("This script demonstrates the _Run introspection functionality")
    
    try:
        test_basic_run_introspection()
        test_run_with_formatting()
        test_run_with_hyperlink()
        test_multiple_runs()
        test_max_depth_control()
        test_format_for_llm_flag()
        
        print("\n" + "="*60)
        print("ALL TESTS COMPLETED SUCCESSFULLY")
        print("="*60)
        
    except Exception as e:
        print(f"\nTEST FAILED: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()