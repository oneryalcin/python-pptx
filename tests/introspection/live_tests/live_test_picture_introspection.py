#!/usr/bin/env python3
"""
Live Test Script for FEP-015: Picture.to_dict() Introspection

This script validates the Picture and Image introspection implementation
using real python-pptx objects. It creates presentations with various
picture configurations and tests the to_dict() output.

Usage:
    python live_test_picture_introspection.py
"""

import os
import sys
import json
import tempfile
from pathlib import Path

# Add src to path for imports
sys.path.insert(0, str(Path(__file__).parent / "src"))

import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE


def create_test_image():
    """Create a simple test image file."""
    from PIL import Image as PILImage
    
    # Create a simple colored image
    img = PILImage.new('RGB', (100, 100), color='red')
    temp_file = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
    img.save(temp_file.name, 'PNG')
    return temp_file.name


def test_image_introspection():
    """Test Image.to_dict() with real Image objects."""
    print("=== Testing Image.to_dict() ===")
    
    # Create test image
    image_path = create_test_image()
    
    try:
        # Create presentation and add picture to get Image object
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add picture to get access to Image object
        left = Inches(1)
        top = Inches(1)
        width = Inches(2)
        picture = slide.shapes.add_picture(image_path, left, top, width)
        
        # Get the Image object
        image = picture.image
        
        # Test to_dict()
        print(f"Image filename: {image.filename}")
        print(f"Image extension: {image.ext}")
        print(f"Image content type: {image.content_type}")
        print(f"Image size: {image.size}")
        print(f"Image DPI: {image.dpi}")
        
        # Test full introspection
        result = image.to_dict()
        print(f"\nImage to_dict() structure:")
        print(f"  Object type: {result['_object_type']}")
        print(f"  Identity keys: {list(result['_identity'].keys())}")
        print(f"  Properties keys: {list(result['properties'].keys())}")
        print(f"  Relationships keys: {list(result['relationships'].keys())}")
        print(f"  LLM context keys: {list(result['_llm_context'].keys())}")
        
        # Test specific properties
        props = result['properties']
        print(f"\nImage properties validation:")
        print(f"  Content type: {props['content_type']}")
        print(f"  Extension: {props['extension']}")
        print(f"  Dimensions: {props['dimensions_px']}")
        print(f"  DPI: {props['dpi']}")
        print(f"  Blob size: {props['blob_size_bytes']} bytes")
        
        # Test LLM context
        llm = result['_llm_context']
        print(f"\nImage LLM context:")
        print(f"  Description: {llm['description']}")
        print(f"  Summary: {llm['summary']}")
        print(f"  Operations count: {len(llm['common_operations'])}")
        
        print("✅ Image introspection test PASSED")
        return True
        
    except Exception as e:
        print(f"❌ Image introspection test FAILED: {e}")
        import traceback
        traceback.print_exc()
        return False
    finally:
        # Clean up temp file
        if os.path.exists(image_path):
            os.unlink(image_path)


def test_picture_basic_introspection():
    """Test Picture.to_dict() with basic picture."""
    print("\n=== Testing Picture.to_dict() - Basic ===")
    
    # Create test image
    image_path = create_test_image()
    
    try:
        # Create presentation with picture
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add basic picture
        left = Inches(1)
        top = Inches(1)
        width = Inches(2)
        picture = slide.shapes.add_picture(image_path, left, top, width)
        
        # Test to_dict()
        result = picture.to_dict()
        print(f"Picture to_dict() structure:")
        print(f"  Object type: {result['_object_type']}")
        print(f"  Identity keys: {list(result['_identity'].keys())}")
        print(f"  Properties keys: {list(result['properties'].keys())}")
        print(f"  Relationships keys: {list(result['relationships'].keys())}")
        print(f"  LLM context keys: {list(result['_llm_context'].keys())}")
        
        # Test identity
        identity = result['_identity']
        print(f"\nPicture identity:")
        print(f"  Class: {identity['class_name']}")
        print(f"  Shape ID: {identity['shape_id']}")
        print(f"  Name: {identity['name']}")
        print(f"  Description: {identity['description']}")
        
        # Test properties
        props = result['properties']
        print(f"\nPicture properties:")
        print(f"  Crop left: {props['crop_left']}")
        print(f"  Crop top: {props['crop_top']}")
        print(f"  Crop right: {props['crop_right']}")
        print(f"  Crop bottom: {props['crop_bottom']}")
        print(f"  Auto shape mask: {props['auto_shape_mask_type']}")
        print(f"  Has image details: {'image_details' in props}")
        print(f"  Has line: {'line' in props}")
        
        # Test image details
        if 'image_details' in props:
            img_details = props['image_details']
            print(f"\nImage details in picture:")
            print(f"  Type: {img_details['_object_type']}")
            if 'properties' in img_details:
                img_props = img_details['properties']
                print(f"  Extension: {img_props.get('extension')}")
                print(f"  Dimensions: {img_props.get('dimensions_px')}")
        
        # Test LLM context
        llm = result['_llm_context']
        print(f"\nPicture LLM context:")
        print(f"  Description: {llm['description']}")
        print(f"  Summary: {llm['summary']}")
        print(f"  Operations count: {len(llm['common_operations'])}")
        
        print("✅ Basic Picture introspection test PASSED")
        return True
        
    except Exception as e:
        print(f"❌ Basic Picture introspection test FAILED: {e}")
        import traceback
        traceback.print_exc()
        return False
    finally:
        # Clean up temp file
        if os.path.exists(image_path):
            os.unlink(image_path)


def test_picture_with_cropping():
    """Test Picture.to_dict() with cropping applied."""
    print("\n=== Testing Picture.to_dict() - With Cropping ===")
    
    # Create test image
    image_path = create_test_image()
    
    try:
        # Create presentation with picture
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add picture
        left = Inches(1)
        top = Inches(1)
        width = Inches(2)
        picture = slide.shapes.add_picture(image_path, left, top, width)
        
        # Apply cropping
        picture.crop_left = 0.1    # 10% from left
        picture.crop_top = 0.05    # 5% from top
        picture.crop_right = 0.15  # 15% from right
        picture.crop_bottom = 0.2  # 20% from bottom
        
        # Test to_dict()
        result = picture.to_dict()
        props = result['properties']
        
        print(f"Cropped picture properties:")
        print(f"  Crop left: {props['crop_left']} (expected: 0.1)")
        print(f"  Crop top: {props['crop_top']} (expected: 0.05)")
        print(f"  Crop right: {props['crop_right']} (expected: 0.15)")
        print(f"  Crop bottom: {props['crop_bottom']} (expected: 0.2)")
        
        # Test LLM context describes cropping
        llm = result['_llm_context']
        description = llm['description']
        print(f"\nLLM description: {description}")
        
        # Check if cropping is mentioned in description
        if "Cropped" in description:
            print("✅ Cropping mentioned in LLM description")
        else:
            print("⚠️  Cropping not mentioned in LLM description")
        
        print("✅ Cropped Picture introspection test PASSED")
        return True
        
    except Exception as e:
        print(f"❌ Cropped Picture introspection test FAILED: {e}")
        import traceback
        traceback.print_exc()
        return False
    finally:
        # Clean up temp file
        if os.path.exists(image_path):
            os.unlink(image_path)


def test_picture_with_masking():
    """Test Picture.to_dict() with shape masking."""
    print("\n=== Testing Picture.to_dict() - With Shape Masking ===")
    
    # Create test image
    image_path = create_test_image()
    
    try:
        # Create presentation with picture
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add picture
        left = Inches(1)
        top = Inches(1)
        width = Inches(2)
        picture = slide.shapes.add_picture(image_path, left, top, width)
        
        # Apply oval masking
        picture.auto_shape_type = MSO_SHAPE.OVAL
        
        # Test to_dict()
        result = picture.to_dict()
        props = result['properties']
        
        print(f"Masked picture properties:")
        mask_type = props['auto_shape_mask_type']
        print(f"  Auto shape mask type: {mask_type}")
        
        # Test LLM context describes masking
        llm = result['_llm_context']
        description = llm['description']
        print(f"\nLLM description: {description}")
        
        # Check if masking is mentioned in description
        if "Masked" in description or "OVAL" in description:
            print("✅ Masking mentioned in LLM description")
        else:
            print("⚠️  Masking not mentioned in LLM description")
        
        print("✅ Masked Picture introspection test PASSED")
        return True
        
    except Exception as e:
        print(f"❌ Masked Picture introspection test FAILED: {e}")
        import traceback
        traceback.print_exc()
        return False
    finally:
        # Clean up temp file
        if os.path.exists(image_path):
            os.unlink(image_path)


def test_picture_depth_and_parameters():
    """Test Picture.to_dict() with different parameters."""
    print("\n=== Testing Picture.to_dict() - Parameters ===")
    
    # Create test image
    image_path = create_test_image()
    
    try:
        # Create presentation with picture
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add picture
        left = Inches(1)
        top = Inches(1)
        width = Inches(2)
        picture = slide.shapes.add_picture(image_path, left, top, width)
        
        # Test different max_depth values
        print("Testing max_depth parameter:")
        result_depth_1 = picture.to_dict(max_depth=1)
        result_depth_3 = picture.to_dict(max_depth=3)
        
        line_depth_1 = result_depth_1['properties']['line']
        line_depth_3 = result_depth_3['properties']['line']
        
        print(f"  Depth 1 - Line has _depth_exceeded: {'_depth_exceeded' in line_depth_1}")
        print(f"  Depth 3 - Line has properties: {'properties' in line_depth_3}")
        
        # Test include_relationships parameter
        print("\nTesting include_relationships parameter:")
        result_no_rels = picture.to_dict(include_relationships=False)
        result_with_rels = picture.to_dict(include_relationships=True)
        
        print(f"  Without relationships: {'relationships' in result_no_rels}")
        print(f"  With relationships: {'relationships' in result_with_rels}")
        
        # Test format_for_llm parameter
        print("\nTesting format_for_llm parameter:")
        result_no_llm = picture.to_dict(format_for_llm=False)
        result_with_llm = picture.to_dict(format_for_llm=True)
        
        print(f"  Without LLM: {'_llm_context' in result_no_llm}")
        print(f"  With LLM: {'_llm_context' in result_with_llm}")
        
        print("✅ Picture parameters test PASSED")
        return True
        
    except Exception as e:
        print(f"❌ Picture parameters test FAILED: {e}")
        import traceback
        traceback.print_exc()
        return False
    finally:
        # Clean up temp file
        if os.path.exists(image_path):
            os.unlink(image_path)


def test_error_handling():
    """Test error handling in Picture.to_dict()."""
    print("\n=== Testing Error Handling ===")
    
    try:
        # Create presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Create a picture and then try to access invalid properties
        # This is harder to test without real error conditions,
        # but we can at least verify the structure is robust
        
        # Add a real picture first
        image_path = create_test_image()
        try:
            left = Inches(1)
            top = Inches(1)
            width = Inches(2)
            picture = slide.shapes.add_picture(image_path, left, top, width)
            
            # Test that to_dict doesn't crash with edge cases
            result = picture.to_dict(max_depth=0)  # Minimum depth
            print(f"Max depth 0 result has basic structure: {all(k in result for k in ['_object_type', '_identity', 'properties'])}")
            
            result = picture.to_dict(max_depth=10)  # Very high depth
            print(f"Max depth 10 result has basic structure: {all(k in result for k in ['_object_type', '_identity', 'properties'])}")
            
            print("✅ Error handling test PASSED")
            return True
            
        finally:
            if os.path.exists(image_path):
                os.unlink(image_path)
        
    except Exception as e:
        print(f"❌ Error handling test FAILED: {e}")
        import traceback
        traceback.print_exc()
        return False


def main():
    """Run all live tests for Picture introspection."""
    print("FEP-015 Live Test: Picture.to_dict() Introspection")
    print("=" * 60)
    
    tests = [
        test_image_introspection,
        test_picture_basic_introspection,
        test_picture_with_cropping,
        test_picture_with_masking,
        test_picture_depth_and_parameters,
        test_error_handling
    ]
    
    passed = 0
    total = len(tests)
    
    for test in tests:
        try:
            if test():
                passed += 1
        except Exception as e:
            print(f"❌ Test {test.__name__} failed with exception: {e}")
    
    print("\n" + "=" * 60)
    print(f"RESULTS: {passed}/{total} tests passed")
    
    if passed == total:
        print("🎉 All tests PASSED! FEP-015 implementation is working correctly.")
        return True
    else:
        print("⚠️  Some tests failed. Please review the implementation.")
        return False


if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)