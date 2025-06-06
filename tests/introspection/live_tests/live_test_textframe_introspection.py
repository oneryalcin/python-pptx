#!/usr/bin/env python3
"""
Live Test Script for FEP-011: TextFrame.to_dict() Introspection

This script demonstrates and validates the TextFrame introspection functionality
using real python-pptx objects. It covers various text frame configurations
and formatting scenarios to ensure the implementation works correctly.

Usage:
    python live_test_textframe_introspection.py

Requirements:
    - python-pptx package with FEP-011 implementation
    - Sample presentation files (created by the script)
"""

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT
    from pptx.util import Inches, Pt
except ImportError as e:
    print(f"Error importing python-pptx: {e}")
    print("Please ensure python-pptx is installed and available in your Python path.")
    sys.exit(1)


def create_test_presentation():
    """Create a test presentation with various text frame configurations."""
    prs = Presentation()
    
    # Slide 1: Basic text frame
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    title_shape = slide.shapes.title
    title_shape.text = "TextFrame Introspection Test"
    
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.text = "This is a basic text frame with default settings."
    
    # Slide 2: Text frame with custom margins and properties
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Add a text box with custom properties
    left = Inches(1)
    top = Inches(1)
    width = Inches(6)
    height = Inches(4)
    textbox = slide2.shapes.add_textbox(left, top, width, height)
    tf2 = textbox.text_frame
    
    # Set various text frame properties
    tf2.margin_left = Inches(0.2)
    tf2.margin_right = Inches(0.2)
    tf2.margin_top = Inches(0.1)
    tf2.margin_bottom = Inches(0.1)
    tf2.word_wrap = True
    tf2.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf2.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    tf2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    tf2.level = 1
    
    # Add multiple paragraphs with different formatting
    p1 = tf2.paragraphs[0]
    p1.text = "First paragraph with center alignment"
    p1.font.name = "Arial"
    p1.font.size = Pt(14)
    p1.font.bold = True
    
    p2 = tf2.add_paragraph()
    p2.text = "Second paragraph with different formatting"
    p2.font.name = "Calibri"
    p2.font.size = Pt(12)
    p2.font.italic = True
    p2.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    
    p3 = tf2.add_paragraph()
    p3.text = "Third paragraph with multiple runs"
    r1 = p3.runs[0]
    r1.text = "Bold text "
    r1.font.bold = True
    
    r2 = p3.add_run()
    r2.text = "and italic text "
    r2.font.italic = True
    
    r3 = p3.add_run()
    r3.text = "and normal text."
    
    return prs


def test_basic_introspection():
    """Test basic TextFrame introspection functionality."""
    print("=" * 60)
    print("TEST 1: Basic TextFrame Introspection")
    print("=" * 60)
    
    prs = create_test_presentation()
    slide = prs.slides[0]
    
    # Get the content text frame
    content_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            content_shape = shape
            break
    
    if content_shape is None:
        print("❌ Could not find a shape with text frame")
        return False
    
    text_frame = content_shape.text_frame
    
    # Test that TextFrame has to_dict method
    if not hasattr(text_frame, 'to_dict'):
        print("❌ TextFrame does not have to_dict method")
        return False
    
    try:
        result = text_frame.to_dict()
        print("✅ TextFrame.to_dict() executed successfully")
        
        # Check basic structure
        required_keys = ["_object_type", "_identity", "properties"]
        for key in required_keys:
            if key not in result:
                print(f"❌ Missing required key: {key}")
                return False
        
        print("✅ Basic structure validation passed")
        
        # Check object type
        if result["_object_type"] != "TextFrame":
            print(f"❌ Incorrect object type: {result['_object_type']}")
            return False
        
        print("✅ Object type validation passed")
        
        # Print condensed result for inspection
        print("\nCondensed result structure:")
        print(f"Object type: {result['_object_type']}")
        print(f"Identity keys: {list(result['_identity'].keys())}")
        print(f"Properties keys: {list(result['properties'].keys())}")
        if "relationships" in result:
            print(f"Relationships keys: {list(result['relationships'].keys())}")
        if "_llm_context" in result:
            print(f"LLM context keys: {list(result['_llm_context'].keys())}")
        
        return True
        
    except Exception as e:
        print(f"❌ Error during introspection: {e}")
        import traceback
        traceback.print_exc()
        return False


def test_properties_introspection():
    """Test detailed properties introspection."""
    print("\n" + "=" * 60)
    print("TEST 2: Properties Introspection")
    print("=" * 60)
    
    prs = create_test_presentation()
    slide = prs.slides[1]  # Second slide with custom text frame
    
    # Find the textbox shape
    textbox = None
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame and shape.shape_type.name == 'TEXT_BOX':
            textbox = shape
            break
    
    if textbox is None:
        print("❌ Could not find textbox shape")
        return False
    
    text_frame = textbox.text_frame
    
    try:
        result = text_frame.to_dict()
        props = result["properties"]
        
        # Test essential properties
        essential_props = [
            "text", "paragraphs", "margin_left", "margin_top", 
            "margin_right", "margin_bottom", "vertical_anchor",
            "word_wrap", "auto_size", "alignment", "level", "font"
        ]
        
        missing_props = []
        for prop in essential_props:
            if prop not in props:
                missing_props.append(prop)
        
        if missing_props:
            print(f"❌ Missing properties: {missing_props}")
            return False
        
        print("✅ All essential properties present")
        
        # Test specific property values
        print(f"Text content length: {len(props['text'])}")
        print(f"Number of paragraphs: {len(props['paragraphs']) if isinstance(props['paragraphs'], list) else 'not expanded'}")
        print(f"Auto-size: {props['auto_size']['name'] if props['auto_size'] else None}")
        print(f"Vertical anchor: {props['vertical_anchor']['name'] if props['vertical_anchor'] else None}")
        print(f"Word wrap: {props['word_wrap']}")
        print(f"Alignment: {props['alignment']['name'] if props['alignment'] else None}")
        print(f"Level: {props['level']}")
        
        # Test margin formatting
        for margin in ["margin_left", "margin_top", "margin_right", "margin_bottom"]:
            margin_data = props[margin]
            if not isinstance(margin_data, dict) or "_object_type" not in margin_data:
                print(f"❌ Margin {margin} not properly formatted")
                return False
        
        print("✅ Margin properties properly formatted")
        
        return True
        
    except Exception as e:
        print(f"❌ Error during properties introspection: {e}")
        import traceback
        traceback.print_exc()
        return False


def test_paragraphs_collection():
    """Test paragraphs collection introspection."""
    print("\n" + "=" * 60)
    print("TEST 3: Paragraphs Collection Introspection")
    print("=" * 60)
    
    prs = create_test_presentation()
    slide = prs.slides[1]  # Second slide with multiple paragraphs
    
    # Find the textbox with multiple paragraphs
    textbox = None
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame and len(shape.text_frame.paragraphs) > 1:
            textbox = shape
            break
    
    if textbox is None:
        print("❌ Could not find textbox with multiple paragraphs")
        return False
    
    text_frame = textbox.text_frame
    
    try:
        # Test with expanded collections
        result_expanded = text_frame.to_dict(expand_collections=True, max_depth=3)
        paragraphs_expanded = result_expanded["properties"]["paragraphs"]
        
        if not isinstance(paragraphs_expanded, list):
            print("❌ Paragraphs not expanded to list")
            return False
        
        print(f"✅ Paragraphs expanded to list with {len(paragraphs_expanded)} items")
        
        # Check first paragraph structure
        first_para = paragraphs_expanded[0]
        if "_object_type" not in first_para or first_para["_object_type"] != "_Paragraph":
            print("❌ First paragraph missing object type")
            return False
        
        if "properties" not in first_para or "text" not in first_para["properties"]:
            print("❌ First paragraph missing text property")
            return False
        
        print(f"✅ First paragraph text: '{first_para['properties']['text'][:50]}...'")
        
        # Test with non-expanded collections
        result_summary = text_frame.to_dict(expand_collections=False)
        paragraphs_summary = result_summary["properties"]["paragraphs"]
        
        if not isinstance(paragraphs_summary, dict) or "_collection_summary" not in paragraphs_summary:
            print("❌ Paragraphs summary not properly formatted")
            return False
        
        print(f"✅ Paragraphs summary: {paragraphs_summary['_collection_summary']}")
        
        return True
        
    except Exception as e:
        print(f"❌ Error during paragraphs collection introspection: {e}")
        import traceback
        traceback.print_exc()
        return False


def test_depth_control():
    """Test max_depth parameter control."""
    print("\n" + "=" * 60)
    print("TEST 4: Depth Control")
    print("=" * 60)
    
    prs = create_test_presentation()
    slide = prs.slides[0]
    
    # Get a text frame
    text_frame = None
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            text_frame = shape.text_frame
            break
    
    if text_frame is None:
        print("❌ Could not find text frame")
        return False
    
    try:
        # Test with depth 1 - font should be depth exceeded
        result_depth1 = text_frame.to_dict(max_depth=1)
        font_data = result_depth1["properties"]["font"]
        
        if not isinstance(font_data, dict) or not font_data.get("_depth_exceeded"):
            print("❌ Font depth control not working at depth 1")
            return False
        
        print("✅ Depth 1: Font marked as depth exceeded")
        
        # Test with depth 3 - font should be expanded
        result_depth3 = text_frame.to_dict(max_depth=3)
        font_data = result_depth3["properties"]["font"]
        
        if isinstance(font_data, dict) and font_data.get("_depth_exceeded"):
            print("❌ Font still marked as depth exceeded at depth 3")
            return False
        
        print("✅ Depth 3: Font properly expanded")
        
        return True
        
    except Exception as e:
        print(f"❌ Error during depth control testing: {e}")
        import traceback
        traceback.print_exc()
        return False


def test_llm_context():
    """Test LLM context generation."""
    print("\n" + "=" * 60)
    print("TEST 5: LLM Context Generation")
    print("=" * 60)
    
    prs = create_test_presentation()
    slide = prs.slides[1]  # Slide with formatted text frame
    
    # Find the textbox shape
    textbox = None
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame and shape.shape_type.name == 'TEXT_BOX':
            textbox = shape
            break
    
    if textbox is None:
        print("❌ Could not find textbox shape")
        return False
    
    text_frame = textbox.text_frame
    
    try:
        result = text_frame.to_dict(format_for_llm=True)
        
        if "_llm_context" not in result:
            print("❌ LLM context not present")
            return False
        
        llm_context = result["_llm_context"]
        required_llm_keys = ["description", "summary", "common_operations"]
        
        for key in required_llm_keys:
            if key not in llm_context:
                print(f"❌ Missing LLM context key: {key}")
                return False
        
        print("✅ All required LLM context keys present")
        
        # Check content quality
        description = llm_context["description"]
        if "TextFrame" not in description:
            print("❌ Description doesn't mention TextFrame")
            return False
        
        if "paragraph(s)" not in description:
            print("❌ Description doesn't mention paragraphs")
            return False
        
        print(f"✅ Description: {description[:100]}...")
        
        # Check operations
        operations = llm_context["common_operations"]
        if not isinstance(operations, list) or len(operations) == 0:
            print("❌ Common operations not properly listed")
            return False
        
        print(f"✅ Common operations: {len(operations)} operations listed")
        print(f"   Sample operations: {operations[:3]}")
        
        return True
        
    except Exception as e:
        print(f"❌ Error during LLM context testing: {e}")
        import traceback
        traceback.print_exc()
        return False


def test_relationships():
    """Test relationships section."""
    print("\n" + "=" * 60)
    print("TEST 6: Relationships Section")
    print("=" * 60)
    
    prs = create_test_presentation()
    slide = prs.slides[0]
    
    # Get a text frame
    text_frame = None
    parent_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            text_frame = shape.text_frame
            parent_shape = shape
            break
    
    if text_frame is None:
        print("❌ Could not find text frame")
        return False
    
    try:
        result = text_frame.to_dict()
        
        if "relationships" not in result:
            print("❌ Relationships section not present")
            return False
        
        relationships = result["relationships"]
        
        if "parent_shape" not in relationships:
            print("❌ Parent shape relationship not present")
            return False
        
        parent_data = relationships["parent_shape"]
        
        # Should have basic structure
        if not isinstance(parent_data, dict):
            print("❌ Parent shape data not a dictionary")
            return False
        
        print("✅ Parent shape relationship present")
        print(f"   Parent shape data type: {type(parent_data)}")
        
        return True
        
    except Exception as e:
        print(f"❌ Error during relationships testing: {e}")
        import traceback
        traceback.print_exc()
        return False


def run_all_tests():
    """Run all live tests and return overall success."""
    print("FEP-011: TextFrame Introspection Live Testing")
    print("=" * 60)
    
    tests = [
        ("Basic Introspection", test_basic_introspection),
        ("Properties Introspection", test_properties_introspection),
        ("Paragraphs Collection", test_paragraphs_collection),
        ("Depth Control", test_depth_control),
        ("LLM Context", test_llm_context),
        ("Relationships", test_relationships),
    ]
    
    results = []
    
    for test_name, test_func in tests:
        try:
            success = test_func()
            results.append((test_name, success))
        except Exception as e:
            print(f"❌ Test '{test_name}' failed with exception: {e}")
            results.append((test_name, False))
    
    # Summary
    print("\n" + "=" * 60)
    print("TEST SUMMARY")
    print("=" * 60)
    
    passed = sum(1 for _, success in results if success)
    total = len(results)
    
    for test_name, success in results:
        status = "✅ PASS" if success else "❌ FAIL"
        print(f"{status} {test_name}")
    
    print(f"\nOverall: {passed}/{total} tests passed")
    
    if passed == total:
        print("🎉 All tests passed! TextFrame introspection is working correctly.")
        return True
    else:
        print("❌ Some tests failed. Please review the implementation.")
        return False


if __name__ == "__main__":
    success = run_all_tests()
    sys.exit(0 if success else 1)