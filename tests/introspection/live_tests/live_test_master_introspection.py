#!/usr/bin/env python3
"""
Live Test Script for FEP-017: SlideMaster.to_dict() & MasterPlaceholder.to_dict()

This script validates the introspection capabilities of SlideMaster and MasterPlaceholder objects
by testing on real python-pptx objects.

Tests:
1. Basic SlideMaster.to_dict() functionality
2. MasterPlaceholder.to_dict() functionality  
3. SlideMaster collections (shapes, placeholders)
4. SlideMaster relationships (slide_layouts, theme_part)
5. Parameter variations (max_depth, expand_collections, etc.)
6. Comprehensive master analysis

Usage:
    python live_test_master_introspection.py

Requirements:
    - python-pptx with FEP-017 implementation
    - Access to default presentation templates
"""

import sys
import json
import traceback
from pathlib import Path

# Add src to path for imports
sys.path.insert(0, str(Path(__file__).parent / "src"))

from pptx import Presentation


def print_test_header(test_name: str, test_num: int):
    """Print a formatted test header."""
    print(f"\n{'='*80}")
    print(f"TEST {test_num}: {test_name}")
    print(f"{'='*80}")


def print_test_result(success: bool, message: str):
    """Print test result with formatting."""
    status = "✅ PASS" if success else "❌ FAIL"
    print(f"\n{status}: {message}")


def safe_dict_access(obj, *keys):
    """Safely access nested dictionary keys."""
    current = obj
    for key in keys:
        if isinstance(current, dict) and key in current:
            current = current[key]
        else:
            return None
    return current


def test_1_basic_slide_master():
    """Test 1: Basic SlideMaster.to_dict() functionality."""
    print_test_header("Basic SlideMaster.to_dict() Functionality", 1)
    
    try:
        # Create a new presentation to get a clean master
        prs = Presentation()
        slide_master = prs.slide_masters[0]
        
        # Test basic to_dict call
        result = slide_master.to_dict(max_depth=2)
        
        # Validate structure
        assert isinstance(result, dict), "Result should be a dictionary"
        assert "_object_type" in result, "Should have _object_type"
        assert result["_object_type"] == "SlideMaster", "Should be SlideMaster type"
        
        # Validate identity
        identity = result.get("_identity", {})
        assert "description" in identity, "Should have description in identity"
        assert "Slide Master" in identity["description"], "Description should mention slide master"
        
        # Validate properties
        props = result.get("properties", {})
        assert "background_fill" in props, "Should have background_fill"
        assert "shapes" in props, "Should have shapes"
        assert "placeholders" in props, "Should have placeholders"
        assert "color_map" in props, "Should have color_map"
        assert "text_styles_summary" in props, "Should have text_styles_summary"
        
        # Validate relationships
        rels = result.get("relationships", {})
        assert "slide_layouts_summary" in rels or "slide_layouts" in rels, "Should have slide layouts info"
        
        # Validate LLM context
        llm_context = result.get("_llm_context", {})
        assert "description" in llm_context, "Should have LLM description"
        assert "common_operations" in llm_context, "Should have common operations"
        
        print(f"SlideMaster Identity: {identity.get('description', 'N/A')}")
        print(f"Layout Count: {rels.get('slide_layouts_summary', 'N/A')}")
        print(f"Background Fill Type: {safe_dict_access(props, 'background_fill', '_object_type')}")
        print(f"Color Map Present: {props.get('color_map') is not None}")
        print(f"Text Styles Present: {safe_dict_access(props, 'text_styles_summary', 'present')}")
        
        print_test_result(True, "Basic SlideMaster.to_dict() works correctly")
        return True
        
    except Exception as e:
        print(f"Error details: {str(e)}")
        print(f"Traceback: {traceback.format_exc()}")
        print_test_result(False, f"Basic SlideMaster.to_dict() failed: {str(e)}")
        return False


def test_2_master_placeholder():
    """Test 2: MasterPlaceholder.to_dict() functionality."""
    print_test_header("MasterPlaceholder.to_dict() Functionality", 2)
    
    try:
        # Create a presentation and get a master placeholder
        prs = Presentation()
        slide_master = prs.slide_masters[0]
        
        # Get a master placeholder
        master_placeholder = None
        for placeholder in slide_master.placeholders:
            master_placeholder = placeholder
            break
            
        if not master_placeholder:
            print_test_result(False, "No master placeholders found to test")
            return False
            
        # Test to_dict call
        result = master_placeholder.to_dict(max_depth=2)
        
        # Validate structure
        assert isinstance(result, dict), "Result should be a dictionary"
        assert "_object_type" in result, "Should have _object_type"
        assert result["_object_type"] == "MasterPlaceholder", "Should be MasterPlaceholder type"
        
        # Validate identity
        identity = result.get("_identity", {})
        assert "description" in identity, "Should have description in identity"
        assert "Master Placeholder" in identity["description"], "Description should mention master placeholder"
        
        # Validate properties  
        props = result.get("properties", {})
        assert "is_master_placeholder" in props, "Should have is_master_placeholder flag"
        assert props["is_master_placeholder"] is True, "Should be marked as master placeholder"
        assert "inheritance_role" in props, "Should have inheritance_role"
        
        # Validate relationships
        rels = result.get("relationships", {})
        # Note: parent_slide_master might be available depending on implementation
        
        # Validate LLM context
        llm_context = result.get("_llm_context", {})
        assert "description" in llm_context, "Should have LLM description"
        assert "role" in llm_context, "Should have role description"
        assert "inheritance_explanation" in llm_context, "Should have inheritance explanation"
        
        print(f"Placeholder Identity: {identity.get('description', 'N/A')}")
        print(f"Placeholder Type: {identity.get('placeholder_type', 'N/A')}")
        print(f"Inheritance Role: {props.get('inheritance_role', 'N/A')}")
        print(f"LLM Description: {llm_context.get('description', 'N/A')[:100]}...")
        
        print_test_result(True, "MasterPlaceholder.to_dict() works correctly")
        return True
        
    except Exception as e:
        print(f"Error details: {str(e)}")
        print(f"Traceback: {traceback.format_exc()}")
        print_test_result(False, f"MasterPlaceholder.to_dict() failed: {str(e)}")
        return False


def test_3_slide_master_collections():
    """Test 3: SlideMaster collections (shapes, placeholders) with expand_collections."""
    print_test_header("SlideMaster Collections with expand_collections", 3)
    
    try:
        prs = Presentation()
        slide_master = prs.slide_masters[0]
        
        # Test with expand_collections=True
        result = slide_master.to_dict(max_depth=2, expand_collections=True)
        props = result.get("properties", {})
        
        # Validate expanded shapes
        shapes = props.get("shapes", [])
        assert isinstance(shapes, list), "Shapes should be a list when expanded"
        
        # Validate expanded placeholders
        placeholders = props.get("placeholders", [])
        assert isinstance(placeholders, list), "Placeholders should be a list when expanded"
        
        # Check placeholder structure
        if placeholders:
            first_placeholder = placeholders[0]
            assert "placeholder_type_key" in first_placeholder, "Should have placeholder_type_key"
            assert "_object_type" in first_placeholder, "Should have _object_type"
            assert first_placeholder["_object_type"] == "MasterPlaceholder", "Should be MasterPlaceholder"
        
        # Test with expand_collections=False
        result_collapsed = slide_master.to_dict(max_depth=2, expand_collections=False)
        props_collapsed = result_collapsed.get("properties", {})
        
        shapes_collapsed = props_collapsed.get("shapes", "")
        placeholders_collapsed = props_collapsed.get("placeholders", "")
        
        assert isinstance(shapes_collapsed, str), "Shapes should be string summary when collapsed"
        assert isinstance(placeholders_collapsed, str), "Placeholders should be string summary when collapsed"
        assert "Collection of" in shapes_collapsed, "Should contain collection summary"
        assert "Collection of" in placeholders_collapsed, "Should contain collection summary"
        
        print(f"Expanded shapes count: {len(shapes)}")
        print(f"Expanded placeholders count: {len(placeholders)}")
        print(f"Collapsed shapes summary: {shapes_collapsed}")
        print(f"Collapsed placeholders summary: {placeholders_collapsed}")
        
        print_test_result(True, "SlideMaster collections work correctly")
        return True
        
    except Exception as e:
        print(f"Error details: {str(e)}")
        print(f"Traceback: {traceback.format_exc()}")
        print_test_result(False, f"SlideMaster collections failed: {str(e)}")
        return False


def test_4_slide_master_relationships():
    """Test 4: SlideMaster relationships (slide_layouts, theme_part, parent_presentation)."""
    print_test_header("SlideMaster Relationships", 4)
    
    try:
        prs = Presentation()
        slide_master = prs.slide_masters[0]
        
        # Test relationships with expand_collections=True
        result = slide_master.to_dict(max_depth=2, expand_collections=True, include_relationships=True)
        rels = result.get("relationships", {})
        
        # Validate slide layouts
        assert "slide_layouts" in rels or "slide_layouts_summary" in rels, "Should have slide layouts info"
        
        if "slide_layouts" in rels:
            slide_layouts = rels["slide_layouts"]
            assert isinstance(slide_layouts, list), "Slide layouts should be a list when expanded"
            if slide_layouts:
                first_layout = slide_layouts[0]
                assert "_object_type" in first_layout, "Layout should have _object_type"
                assert first_layout["_object_type"] == "SlideLayout", "Should be SlideLayout"
        
        # Validate theme part reference
        assert "theme_part_ref" in rels, "Should have theme_part_ref"
        theme_ref = rels["theme_part_ref"]
        if isinstance(theme_ref, dict):
            assert "partname" in theme_ref, "Theme ref should have partname"
            assert "_object_type" in theme_ref, "Theme ref should have _object_type"
        
        # Test with collapsed relationships
        result_collapsed = slide_master.to_dict(max_depth=1, expand_collections=False)
        rels_collapsed = result_collapsed.get("relationships", {})
        
        if "slide_layouts_summary" in rels_collapsed:
            summary = rels_collapsed["slide_layouts_summary"]
            assert isinstance(summary, str), "Summary should be string"
            assert "slide layout" in summary.lower(), "Summary should mention slide layouts"
        
        print(f"Slide layouts: {rels.get('slide_layouts_summary', len(rels.get('slide_layouts', [])))}")
        print(f"Theme part: {rels.get('theme_part_ref', 'N/A')}")
        print(f"Parent presentation: {'present' if 'parent_presentation' in rels else 'not included'}")
        
        print_test_result(True, "SlideMaster relationships work correctly")
        return True
        
    except Exception as e:
        print(f"Error details: {str(e)}")
        print(f"Traceback: {traceback.format_exc()}")
        print_test_result(False, f"SlideMaster relationships failed: {str(e)}")
        return False


def test_5_parameter_variations():
    """Test 5: Parameter variations (max_depth, include_relationships, format_for_llm)."""
    print_test_header("Parameter Variations", 5)
    
    try:
        prs = Presentation()
        slide_master = prs.slide_masters[0]
        
        # Test different max_depth values
        result_depth_0 = slide_master.to_dict(max_depth=0)
        result_depth_1 = slide_master.to_dict(max_depth=1)
        result_depth_2 = slide_master.to_dict(max_depth=2)
        
        # Depth 0 should be truncated at the top level
        assert "_truncated" in result_depth_0, "Should be truncated at depth 0"
        assert "properties" not in result_depth_0, "Should not have properties at depth 0"
        
        # Test include_relationships=False
        result_no_rels = slide_master.to_dict(include_relationships=False)
        assert "relationships" not in result_no_rels, "Should not include relationships when disabled"
        
        # Test format_for_llm=True vs False
        result_llm_true = slide_master.to_dict(format_for_llm=True)
        result_llm_false = slide_master.to_dict(format_for_llm=False)
        
        assert "_llm_context" in result_llm_true, "Should have LLM context when format_for_llm=True"
        # Note: format_for_llm=False might still include LLM context, depends on implementation
        
        print(f"Depth 0 result: {result_depth_0}")
        print(f"No relationships keys: {list(result_no_rels.keys())}")
        print(f"LLM context present (True): {'_llm_context' in result_llm_true}")
        print(f"LLM context present (False): {'_llm_context' in result_llm_false}")
        
        print_test_result(True, "Parameter variations work correctly")
        return True
        
    except Exception as e:
        print(f"Error details: {str(e)}")
        print(f"Traceback: {traceback.format_exc()}")
        print_test_result(False, f"Parameter variations failed: {str(e)}")
        return False


def test_6_comprehensive_master_analysis():
    """Test 6: Comprehensive analysis of slide master structure."""
    print_test_header("Comprehensive Master Analysis", 6)
    
    try:
        prs = Presentation()
        slide_master = prs.slide_masters[0]
        
        # Get comprehensive data
        result = slide_master.to_dict(max_depth=3, expand_collections=True, include_relationships=True)
        
        # Analyze structure
        properties = result.get("properties", {})
        relationships = result.get("relationships", {})
        llm_context = result.get("_llm_context", {})
        
        # Extract key metrics
        color_map = properties.get("color_map", {})
        text_styles = properties.get("text_styles_summary", {})
        
        placeholders = properties.get("placeholders", [])
        placeholder_types = []
        if isinstance(placeholders, list):
            for ph in placeholders:
                ph_type = ph.get("placeholder_type_key", "Unknown")
                placeholder_types.append(str(ph_type))
        
        shapes = properties.get("shapes", [])
        shape_count = len(shapes) if isinstance(shapes, list) else "Unknown"
        
        # Print comprehensive analysis
        print(f"\n📊 COMPREHENSIVE SLIDE MASTER ANALYSIS")
        print(f"{'─'*60}")
        print(f"Master Name: {result.get('_identity', {}).get('name', 'Default Master')}")
        print(f"Object Type: {result.get('_object_type', 'Unknown')}")
        print(f"Description: {result.get('_identity', {}).get('description', 'N/A')}")
        
        print(f"\n🎨 DESIGN PROPERTIES")
        print(f"Background Fill: {safe_dict_access(properties, 'background_fill', '_object_type')}")
        print(f"Color Map: {'Present' if color_map else 'Not present'}")
        if color_map:
            print(f"  - Background colors: bg1={color_map.get('bg1')}, bg2={color_map.get('bg2')}")
            print(f"  - Text colors: tx1={color_map.get('tx1')}, tx2={color_map.get('tx2')}")
        
        print(f"Text Styles: {'Present' if text_styles.get('present') else 'Not present'}")
        if text_styles.get('available_styles'):
            print(f"  - Available: {', '.join(text_styles['available_styles'])}")
        
        print(f"\n📐 CONTENT STRUCTURE")
        print(f"Non-placeholder shapes: {shape_count}")
        print(f"Master placeholders: {len(placeholder_types)}")
        if placeholder_types:
            print(f"  - Types: {', '.join(placeholder_types)}")
        
        print(f"\n🔗 RELATIONSHIPS")
        layouts_info = relationships.get("slide_layouts_summary", relationships.get("slide_layouts", []))
        if isinstance(layouts_info, list):
            print(f"Slide layouts: {len(layouts_info)} layouts")
        else:
            print(f"Slide layouts: {layouts_info}")
            
        theme_part = relationships.get("theme_part_ref", "N/A")
        if isinstance(theme_part, dict):
            print(f"Theme part: {theme_part.get('partname', 'Unknown')}")
        else:
            print(f"Theme part: {theme_part}")
        
        print(f"\n🤖 LLM CONTEXT")
        print(f"Description: {llm_context.get('description', 'N/A')}")
        print(f"Role: {llm_context.get('role', 'N/A')}")
        
        operations = llm_context.get('common_operations', [])
        if operations:
            print(f"Common operations ({len(operations)}):")
            for i, op in enumerate(operations[:3], 1):  # Show first 3
                print(f"  {i}. {op}")
            if len(operations) > 3:
                print(f"  ... and {len(operations) - 3} more")
        
        print_test_result(True, "Comprehensive master analysis completed successfully")
        return True
        
    except Exception as e:
        print(f"Error details: {str(e)}")
        print(f"Traceback: {traceback.format_exc()}")
        print_test_result(False, f"Comprehensive master analysis failed: {str(e)}")
        return False


def main():
    """Run all live tests for FEP-017."""
    print("🧪 FEP-017 Live Test Suite: SlideMaster & MasterPlaceholder Introspection")
    print(f"{'='*80}")
    
    # Track test results
    tests = [
        test_1_basic_slide_master,
        test_2_master_placeholder,
        test_3_slide_master_collections,
        test_4_slide_master_relationships,
        test_5_parameter_variations,
        test_6_comprehensive_master_analysis,
    ]
    
    results = []
    for test in tests:
        try:
            success = test()
            results.append(success)
        except Exception as e:
            print(f"❌ Test {test.__name__} crashed: {str(e)}")
            results.append(False)
    
    # Print summary
    passed = sum(results)
    total = len(results)
    
    print(f"\n{'='*80}")
    print(f"📊 TEST SUMMARY")
    print(f"{'='*80}")
    print(f"Tests passed: {passed}/{total}")
    print(f"Tests failed: {total - passed}/{total}")
    print(f"Success rate: {(passed/total)*100:.1f}%")
    
    if passed == total:
        print(f"🎉 ALL TESTS PASSED! FEP-017 implementation is working correctly.")
        return 0
    else:
        print(f"❌ Some tests failed. Please review the implementation.")
        return 1


if __name__ == "__main__":
    exit_code = main()
    sys.exit(exit_code)