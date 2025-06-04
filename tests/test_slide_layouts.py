"""Unit tests for slide layout creation and manipulation."""

import pytest
from io import BytesIO

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.shapes.placeholder import LayoutPlaceholder
from pptx.slide import SlideLayout
from pptx.util import Inches


class TestSlideLayouts:
    """Test suite for features related to slide layouts."""

    def test_add_blank_layout(self, prs_with_master):
        """Test adding a new blank layout to a slide master."""
        slide_master = prs_with_master.slide_masters[0]
        initial_count = len(slide_master.slide_layouts)

        new_layout = slide_master.slide_layouts.add_layout(
            name="Test Blank Layout", base_type="blank"
        )

        assert isinstance(new_layout, SlideLayout)
        assert len(slide_master.slide_layouts) == initial_count + 1
        assert new_layout.name == "Test Blank Layout"
        assert new_layout.layout_type == "blank"
        assert slide_master.slide_layouts[-1] == new_layout
        # Check if the name is actually set on the cSld element
        assert new_layout._element.cSld.name == "Test Blank Layout"
        # Check if the type is set on the sldLayout element
        assert new_layout._element.type == "blank"


    def test_add_placeholders_to_layout(self, prs_with_master):
        """Test adding various placeholders to a newly created slide layout."""
        slide_master = prs_with_master.slide_masters[0]
        layout = slide_master.slide_layouts.add_layout(
            name="Layout With Placeholders", base_type="custom"
        )

        # Add a TITLE placeholder
        title_ph = layout.placeholders.add_placeholder(
            ph_type=PP_PLACEHOLDER.TITLE,
            left=Inches(1.0), top=Inches(0.5),
            width=Inches(8.0), height=Inches(1.0),
            idx=0
        )
        assert isinstance(title_ph, LayoutPlaceholder)
        assert len(layout.placeholders) == 1
        assert title_ph.placeholder_format.type == PP_PLACEHOLDER.TITLE
        assert title_ph.placeholder_format.idx == 0
        assert title_ph.left == Inches(1.0)
        assert title_ph.top == Inches(0.5)
        assert title_ph.width == Inches(8.0)
        assert title_ph.height == Inches(1.0)

        # Add a BODY placeholder
        body_ph = layout.placeholders.add_placeholder(
            ph_type=PP_PLACEHOLDER.BODY,
            left=Inches(1.0), top=Inches(2.0),
            width=Inches(8.0), height=Inches(4.0),
            idx=1
        )
        assert len(layout.placeholders) == 2
        assert isinstance(body_ph, LayoutPlaceholder)
        assert body_ph.placeholder_format.type == PP_PLACEHOLDER.BODY
        assert body_ph.placeholder_format.idx == 1
        assert body_ph.left == Inches(1.0)
        assert body_ph.top == Inches(2.0)

        # Check if placeholders are retrievable by index from the collection
        # (assuming LayoutPlaceholders is ordered by insertion or idx)
        # Iteration order of LayoutPlaceholders is based on XML order,
        # which should be insertion order here.
        placeholders_list = list(layout.placeholders)
        assert placeholders_list[0] == title_ph
        assert placeholders_list[1] == body_ph

        # Test LayoutPlaceholders.get(idx=...)
        assert layout.placeholders.get(idx=0) == title_ph
        assert layout.placeholders.get(idx=1) == body_ph


    def test_layout_persistence(self, prs_with_master):
        """Test that added layouts and their placeholders persist after save and load."""
        prs = prs_with_master  # Use the fixture for a clean presentation
        master = prs.slide_masters[0]

        layout_name = "Persistent Layout"
        layout_type = "customTypeForTest"
        title_idx, pic_idx = 10, 12
        title_left, title_top = Inches(0.5), Inches(0.5)
        pic_width, pic_height = Inches(3.0), Inches(4.0)

        layout = master.slide_layouts.add_layout(name=layout_name, base_type=layout_type)
        layout.placeholders.add_placeholder(
            PP_PLACEHOLDER.TITLE, title_left, title_top, Inches(9), Inches(1), idx=title_idx
        )
        layout.placeholders.add_placeholder(
            PP_PLACEHOLDER.PICTURE, Inches(1), Inches(2), pic_width, pic_height, idx=pic_idx
        )

        stream = BytesIO()
        prs.save(stream)
        stream.seek(0)

        prs_reloaded = Presentation(stream)
        master_reloaded = prs_reloaded.slide_masters[0]

        layout_reloaded = None
        for l in master_reloaded.slide_layouts:
            if l.name == layout_name:
                layout_reloaded = l
                break

        assert layout_reloaded is not None, f"Layout named '{layout_name}' not found after reload."
        assert layout_reloaded.name == layout_name
        assert layout_reloaded.layout_type == layout_type
        assert len(layout_reloaded.placeholders) == 2

        title_ph_reloaded = layout_reloaded.placeholders.get(idx=title_idx)
        assert title_ph_reloaded is not None
        assert title_ph_reloaded.placeholder_format.type == PP_PLACEHOLDER.TITLE
        assert title_ph_reloaded.placeholder_format.idx == title_idx
        assert title_ph_reloaded.left == title_left
        assert title_ph_reloaded.top == title_top

        pic_ph_reloaded = layout_reloaded.placeholders.get(idx=pic_idx)
        assert pic_ph_reloaded is not None
        assert pic_ph_reloaded.placeholder_format.type == PP_PLACEHOLDER.PICTURE
        assert pic_ph_reloaded.placeholder_format.idx == pic_idx
        assert pic_ph_reloaded.width == pic_width
        assert pic_ph_reloaded.height == pic_height

# --- Fixtures ---

@pytest.fixture
def prs_with_master(request):
    """Provides a Presentation instance with at least one slide master."""
    return Presentation()

# Further tests could include:
# - Adding layouts with all different base_type values.
# - Adding all placeholder types.
# - Testing behavior with multiple slide masters.
# - Testing placeholder name generation more extensively.
# - Testing ID uniqueness for p:sldLayoutId elements if possible to inspect.
# - Testing removal of layouts (if/when implemented).
# - Testing layouts with no placeholders.
# - Testing layouts with overlapping placeholders.
# - Edge cases for placeholder dimensions/positions.
# - Max number of layouts / placeholders.
# - Layout names with special characters.
# - Interaction with existing layouts in a presentation.
# - Correct `rId` generation and usage (harder to test directly at this level).
# - Correct `id` attribute for `p:sldLayoutId` (harder to test directly at this level).
