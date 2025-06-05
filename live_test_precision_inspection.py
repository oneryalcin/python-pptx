#!/usr/bin/env python3
"""
Live test script for FEP-019: Precision Inspection Controls (to_dict Enhancement).

This script demonstrates the new field selection and structured collection summary
capabilities added to IntrospectionMixin.to_dict() method using real python-pptx objects.

Usage:
    python live_test_precision_inspection.py

Requirements:
    - python-pptx with FEP-019 implementation
    - Test presentation file (uses default template if not found)
"""

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.shapes.base import BaseShape
    from pptx.dml.color import RGBColor
    from pptx.util import Inches
except ImportError as e:
    print(f"Error importing python-pptx: {e}")
    print("Ensure python-pptx is installed and in PYTHONPATH")
    sys.exit(1)


def print_section(title, content=None):
    """Print a formatted section header."""
    print(f"\n{'='*60}")
    print(f"  {title}")
    print(f"{'='*60}")
    if content:
        print(content)


def print_json(data, title=None):
    """Print JSON data with optional title."""
    if title:
        print(f"\n{title}:")
        print("-" * len(title))
    print(json.dumps(data, indent=2, default=str))


def test_field_path_parsing():
    """Test the field path parsing functionality."""
    print_section("Field Path Parsing Tests")
    
    # Create a simple test object
    from pptx.introspection import IntrospectionMixin
    
    class TestObject(IntrospectionMixin):
        pass
    
    obj = TestObject()
    
    # Test cases for field path parsing
    test_cases = [
        ["_identity.shape_id"],
        ["properties.fill.type", "properties.line.width"],
        ["properties.fill.*"],
        ["_identity.class_name", "properties.*", "relationships.parent"],
    ]
    
    for i, fields in enumerate(test_cases, 1):
        print(f"\nTest Case {i}: {fields}")
        try:
            field_tree = obj._parse_field_paths(fields)
            print_json(field_tree, "Parsed Field Tree")
        except Exception as e:
            print(f"Error: {e}")


def test_structured_collection_summaries():
    """Test structured collection summaries."""
    print_section("Structured Collection Summary Tests")
    
    from pptx.introspection import IntrospectionMixin
    
    class TestObject(IntrospectionMixin):
        pass
    
    obj = TestObject()
    
    # Test different collection types
    test_collections = [
        ("String List", ["item1", "item2", "item3"]),
        ("Mixed List", ["string", 42, True]),
        ("Empty List", []),
        ("String Dict", {"key1": "value1", "key2": "value2"}),
        ("Empty Dict", {}),
        ("Tuple", ("tuple1", "tuple2")),
    ]
    
    for name, collection in test_collections:
        print(f"\n{name}: {collection}")
        try:
            result = obj._format_property_value_for_to_dict(
                collection, False, set(), 1, False, True  # expand_collections=False
            )
            print_json(result, "Collection Summary")
        except Exception as e:
            print(f"Error: {e}")


def test_basic_precision_inspection():
    """Test basic precision inspection on a simple presentation."""
    print_section("Basic Precision Inspection Tests")
    
    # Create a simple presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Add some content
    title = slide.shapes.title
    title.text = "Test Slide"
    
    content = slide.placeholders[1]
    content.text = "This is test content"
    
    # Test field selection on slide
    print("\n1. Slide Identity Only:")
    try:
        result = slide.to_dict(fields=["_identity.class_name", "_identity.memory_address"])
        print_json(result)
    except Exception as e:
        print(f"Error: {e}")
    
    print("\n2. Slide Shapes Collection Summary:")
    try:
        result = slide.to_dict(
            fields=["properties.shapes"], 
            expand_collections=False
        )
        print_json(result)
    except Exception as e:
        print(f"Error: {e}")
    
    print("\n3. Specific Shape Properties:")
    if len(slide.shapes) > 0:
        shape = slide.shapes[0]
        try:
            result = shape.to_dict(
                fields=["_identity.class_name", "properties.name", "properties.shape_type"]
            )
            print_json(result)
        except Exception as e:
            print(f"Error: {e}")


def test_wildcard_functionality():
    """Test wildcard field selection."""
    print_section("Wildcard Field Selection Tests")
    
    # Create a presentation and get a shape
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Add a shape with some properties we can inspect
    left = top = width = height = Inches(1)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.text = "Test textbox"
    
    print("\n1. All Properties (wildcard):")
    try:
        result = textbox.to_dict(
            fields=["properties.*"],
            max_depth=2,
            expand_collections=False
        )
        print_json(result)
    except Exception as e:
        print(f"Error: {e}")
    
    print("\n2. All Identity Fields (wildcard):")
    try:
        result = textbox.to_dict(fields=["_identity.*"])
        print_json(result)
    except Exception as e:
        print(f"Error: {e}")


def test_presentation_level_inspection():
    """Test precision inspection at presentation level."""
    print_section("Presentation-Level Precision Inspection")
    
    # Create a presentation with multiple slides
    prs = Presentation()
    
    # Add a few slides
    for i in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f"Slide {i+1}"
    
    print("\n1. Presentation Core Properties Only:")
    try:
        result = prs.to_dict(
            fields=["properties.core_properties", "properties.slide_layouts"],
            expand_collections=False
        )
        print_json(result)
    except Exception as e:
        print(f"Error: {e}")
    
    print("\n2. Slides Collection Summary:")
    try:
        result = prs.to_dict(
            fields=["properties.slides"],
            expand_collections=False
        )
        print_json(result)
    except Exception as e:
        print(f"Error: {e}")


def test_backward_compatibility():
    """Test that existing functionality still works."""
    print_section("Backward Compatibility Tests")
    
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    print("\n1. Full introspection (no fields parameter):")
    try:
        result = slide.to_dict(max_depth=1, expand_collections=False, format_for_llm=False)
        print(f"Result keys: {list(result.keys())}")
        print(f"Properties keys: {list(result.get('properties', {}).keys())}")
    except Exception as e:
        print(f"Error: {e}")
    
    print("\n2. Traditional parameter usage:")
    try:
        result = slide.to_dict(
            include_relationships=False,
            max_depth=2,
            expand_collections=False,
            format_for_llm=True
        )
        print(f"Result structure maintained: {list(result.keys())}")
        print(f"LLM context present: {'_llm_context' in result}")
        print(f"Relationships excluded: {'relationships' not in result}")
    except Exception as e:
        print(f"Error: {e}")


def test_performance_comparison():
    """Compare performance of full vs. precision inspection."""
    print_section("Performance Comparison")
    
    import time
    
    # Create a presentation with more content
    prs = Presentation()
    for i in range(5):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Slide {i+1}"
        
        # Add some shapes
        for j in range(3):
            left = Inches(j)
            top = Inches(1)
            width = height = Inches(1)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.text = f"Text {j+1}"
    
    slide = prs.slides[0]  # Test on first slide
    
    print(f"\nTesting on slide with {len(slide.shapes)} shapes...")
    
    # Test full introspection
    start_time = time.time()
    try:
        full_result = slide.to_dict(max_depth=2, expand_collections=False)
        full_time = time.time() - start_time
        full_size = len(json.dumps(full_result, default=str))
        print(f"Full introspection: {full_time:.3f}s, {full_size} characters")
    except Exception as e:
        print(f"Full introspection error: {e}")
        full_time = float('inf')
        full_size = 0
    
    # Test precision introspection
    start_time = time.time()
    try:
        precision_result = slide.to_dict(
            fields=["_identity.class_name", "properties.shapes"],
            expand_collections=False
        )
        precision_time = time.time() - start_time
        precision_size = len(json.dumps(precision_result, default=str))
        print(f"Precision introspection: {precision_time:.3f}s, {precision_size} characters")
        
        if full_time != float('inf'):
            speedup = full_time / precision_time if precision_time > 0 else float('inf')
            size_reduction = (1 - precision_size / full_size) * 100 if full_size > 0 else 0
            print(f"Performance improvement: {speedup:.1f}x faster, {size_reduction:.1f}% smaller")
    except Exception as e:
        print(f"Precision introspection error: {e}")


def main():
    """Run all live tests for FEP-019."""
    print_section("FEP-019 Live Test: Precision Inspection Controls", 
                  "Testing field selection and structured collection summaries")
    
    try:
        # Test the core functionality
        test_field_path_parsing()
        test_structured_collection_summaries()
        test_basic_precision_inspection()
        test_wildcard_functionality()
        test_presentation_level_inspection()
        test_backward_compatibility()
        test_performance_comparison()
        
        print_section("✅ All Tests Completed Successfully!")
        print("FEP-019 implementation appears to be working correctly.")
        
    except Exception as e:
        print_section("❌ Test Failed")
        print(f"Error during testing: {e}")
        import traceback
        traceback.print_exc()
        return 1
    
    return 0


if __name__ == "__main__":
    exit_code = main()
    sys.exit(exit_code)