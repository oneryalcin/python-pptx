#!/usr/bin/env python3
"""
Demo script for MEP-004: save_presentation tool.

This script demonstrates how the save_presentation tool works in practice
with real PowerPoint files and shows both Save and Save As operations.
"""

import asyncio
import json
import sys
import tempfile
import shutil
from pathlib import Path

# Add the mcp_server directory to Python path for imports
server_dir = Path(__file__).parent.parent
sys.path.insert(0, str(server_dir))

try:
    from mcp import ClientSession, StdioServerParameters
    from mcp.client.stdio import stdio_client
except ImportError:
    print("Error: MCP client libraries not available. Install with: pip install mcp")
    sys.exit(1)


async def demo_save_presentation():
    """Demonstrate the save_presentation tool functionality."""
    
    print("🎯 MEP-004 Save Presentation Tool Demo")
    print("="*50)
    
    # Setup test environment
    temp_dir = Path(tempfile.mkdtemp(prefix="mcp_save_demo_"))
    print(f"📁 Test directory: {temp_dir}")
    
    try:
        # Create a test presentation file
        test_files_dir = Path(__file__).parent.parent.parent / "tests" / "test_files"
        minimal_pptx = test_files_dir / "minimal.pptx"
        
        if minimal_pptx.exists():
            test_pptx = temp_dir / "demo_presentation.pptx"
            shutil.copy(minimal_pptx, test_pptx)
            print(f"📄 Created test presentation: {test_pptx.name}")
        else:
            # Create basic presentation if test file not available
            try:
                from pptx import Presentation
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = "MEP-004 Demo Presentation"
                
                test_pptx = temp_dir / "demo_presentation.pptx"
                prs.save(str(test_pptx))
                print(f"📄 Created test presentation: {test_pptx.name}")
            except ImportError:
                print("❌ Cannot create test presentation - python-pptx not available")
                return
        
        # Start MCP server communication
        server_path = server_dir / "server" / "main.py"
        server_params = StdioServerParameters(
            command=sys.executable,
            args=[str(server_path)]
        )
        
        print("\n🔗 Connecting to MCP server...")
        
        async with stdio_client(server_params) as (read, write):
            async with ClientSession(read, write) as session:
                await session.initialize()
                
                print("✅ Connected to MCP server")
                
                # Discover save_presentation tool
                print("\n🔍 Discovering tools...")
                tools_response = await session.list_tools()
                
                save_tool = None
                for tool in tools_response.tools:
                    if tool.name == "save_presentation":
                        save_tool = tool
                        break
                
                if save_tool:
                    print(f"✅ Found save_presentation tool")
                    print(f"   Description: {save_tool.description.split('.')[0]}...")
                else:
                    print("❌ save_presentation tool not found")
                    return
                
                # Demo 1: Try to save without presentation loaded
                print("\n📝 Demo 1: Save without presentation loaded")
                result = await session.call_tool("save_presentation", {})
                response = json.loads(result.content[0].text)
                
                print(f"   Result: {response['success']}")
                print(f"   Error: {response['error']}")
                
                # Demo 2: Execute code to modify a (hypothetical) presentation
                print("\n📝 Demo 2: Execute code to show context")
                code = """
# Show what's available in execution context
print("Available objects in context:")
print(f"- pptx module available: {'pptx' in globals()}")
print(f"- prs object available: {'prs' in globals()}")
print(f"- json module available: {'json' in globals()}")

# If prs were available, we could modify it:
# title_slide = prs.slides[0]
# title_shape = title_slide.shapes.title  
# title_shape.text = "Modified by Agent"
"""
                
                result = await session.call_tool("execute_python_code", {"code": code})
                execute_response = json.loads(result.content[0].text)
                
                print(f"   Execute result: {execute_response['success']}")
                if execute_response['stdout']:
                    print(f"   Output: {execute_response['stdout'].strip()}")
                if not execute_response['success']:
                    print(f"   Error: {execute_response['error']}")
                
                # Demo 3: Try Save As to new location
                print("\n📝 Demo 3: Save As to new location")
                new_path = str(temp_dir / "saved_copy.pptx")
                result = await session.call_tool("save_presentation", {"output_path": new_path})
                response = json.loads(result.content[0].text)
                
                print(f"   Result: {response['success']}")
                if response['success']:
                    print(f"   Saved to: {response['file_path']}")
                    print(f"   Operation: {response['operation']}")
                else:
                    print(f"   Error: {response['error']}")
                
                # Demo 4: Try to save outside allowed directory (security test)
                print("\n📝 Demo 4: Security test - save outside root")
                outside_path = "/tmp/unauthorized_save.pptx"
                result = await session.call_tool("save_presentation", {"output_path": outside_path})
                response = json.loads(result.content[0].text)
                
                print(f"   Result: {response['success']}")
                print(f"   Security check: {'✅ BLOCKED' if not response['success'] else '❌ ALLOWED'}")
                if not response['success']:
                    print(f"   Error: {response['error']}")
                
                # Demo 5: Save to subdirectory (directory creation test)
                print("\n📝 Demo 5: Save to subdirectory")
                subdir_path = str(temp_dir / "backup" / "presentation_backup.pptx")
                result = await session.call_tool("save_presentation", {"output_path": subdir_path})
                response = json.loads(result.content[0].text)
                
                print(f"   Result: {response['success']}")
                if response['success']:
                    print(f"   Created directory and saved to: {response['file_path']}")
                else:
                    print(f"   Error: {response['error']}")
        
        print("\n✅ Demo completed successfully!")
        
        # Show what files were created
        print(f"\n📂 Files in test directory:")
        for file_path in temp_dir.rglob("*"):
            if file_path.is_file():
                rel_path = file_path.relative_to(temp_dir)
                size = file_path.stat().st_size
                print(f"   {rel_path} ({size} bytes)")
        
    finally:
        # Cleanup
        shutil.rmtree(temp_dir)
        print(f"\n🧹 Cleaned up test directory")


if __name__ == "__main__":
    try:
        asyncio.run(demo_save_presentation())
    except KeyboardInterrupt:
        print("\n⚠️  Demo interrupted by user")
    except Exception as e:
        print(f"\n❌ Demo failed: {e}")
        sys.exit(1)