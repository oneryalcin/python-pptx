#!/usr/bin/env python3
"""
Live test script for MEP-003: Root and Resource Management implementation.

This script validates the actual MCP server functionality by simulating
a real MCP client interaction with roots and resources. It tests:
- Server startup with resource capabilities
- Resource discovery (list_resources)
- Resource reading (read_resource with get_tree() output)
- Simplified execute_python_code tool (no file_path parameter)

Usage:
    python mcp_server/tests/live_test_mep003.py

Requirements:
    - Virtual environment must be activated
    - MCP dependencies must be installed
    - Test presentation file in the project directory
"""

import asyncio
import json
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any, Dict, List

# Add the project root to Python path for imports
PROJECT_ROOT = Path(__file__).parent.parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

# Import MCP client functionality for testing
try:
    from mcp import types
    from mcp.client.session import ClientSession
    from mcp.client.stdio import StdioServerParameters, stdio_client
except ImportError as e:
    print(f"Error: Could not import MCP client libraries: {e}")
    print("Please ensure MCP dependencies are installed with:")
    print("  pip install -r requirements-dev.txt")
    sys.exit(1)


class MEP003Tester:
    """Test harness for MEP-003 functionality."""

    def __init__(self):
        self.server_path = PROJECT_ROOT / "mcp_server" / "server" / "main.py"
        self.test_results: List[Dict[str, Any]] = []
        self.test_pptx_path = None

    def log_test(self, test_name: str, success: bool, details: str = "", error: str = ""):
        """Log a test result."""
        result = {
            "test": test_name,
            "success": success,
            "details": details,
            "error": error
        }
        self.test_results.append(result)
        
        status = "✅ PASS" if success else "❌ FAIL"
        print(f"{status} {test_name}")
        if details:
            print(f"     Details: {details}")
        if error:
            print(f"     Error: {error}")

    def create_test_presentation(self) -> Path:
        """Create a minimal test presentation file."""
        try:
            import pptx
            
            # Create a temporary .pptx file
            temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
            temp_path = Path(temp_file.name)
            temp_file.close()
            
            # Create a minimal presentation
            prs = pptx.Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            title.text = "MEP-003 Test Presentation"
            
            prs.save(temp_path)
            self.test_pptx_path = temp_path
            return temp_path
            
        except ImportError:
            # If pptx is not available, create a dummy file for testing
            temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
            temp_path = Path(temp_file.name)
            temp_file.write(b"dummy pptx content for testing")
            temp_file.close()
            self.test_pptx_path = temp_path
            return temp_path
        except Exception as e:
            raise RuntimeError(f"Failed to create test presentation: {e}")

    def cleanup_test_presentation(self):
        """Clean up the test presentation file."""
        if self.test_pptx_path and self.test_pptx_path.exists():
            self.test_pptx_path.unlink()

    async def test_server_startup(self) -> bool:
        """Test that the server starts without errors."""
        try:
            server_params = StdioServerParameters(
                command=sys.executable,
                args=[str(self.server_path)]
            )
            
            async with stdio_client(server_params) as (read, write):
                async with ClientSession(read, write) as session:
                    await session.initialize()
                    
                    self.log_test(
                        "Server Startup",
                        True,
                        "Server started and initialized successfully"
                    )
                    return True
                    
        except Exception as e:
            self.log_test(
                "Server Startup",
                False,
                error=f"Failed to start server: {str(e)}"
            )
            return False

    async def test_resource_capabilities(self) -> bool:
        """Test that the server advertises resource capabilities."""
        try:
            server_params = StdioServerParameters(
                command=sys.executable,
                args=[str(self.server_path)]
            )
            
            async with stdio_client(server_params) as (read, write):
                async with ClientSession(read, write) as session:
                    result = await session.initialize()
                    
                    # Check if resources capability is advertised
                    capabilities = result.capabilities
                    has_resources = hasattr(capabilities, 'resources')
                    
                    self.log_test(
                        "Resource Capabilities",
                        has_resources,
                        f"Server capabilities: {capabilities}" if has_resources else "",
                        "Server does not advertise resources capability" if not has_resources else ""
                    )
                    return has_resources
                    
        except Exception as e:
            self.log_test(
                "Resource Capabilities",
                False,
                error=f"Failed to check capabilities: {str(e)}"
            )
            return False

    async def test_list_resources_empty(self) -> bool:
        """Test listing resources when no presentation is loaded."""
        try:
            server_params = StdioServerParameters(
                command=sys.executable,
                args=[str(self.server_path)]
            )
            
            async with stdio_client(server_params) as (read, write):
                async with ClientSession(read, write) as session:
                    await session.initialize()
                    
                    # Try to list resources (should be empty initially)
                    try:
                        resources = await session.list_resources()
                        resource_count = len(resources.resources)
                        
                        self.log_test(
                            "List Resources (Empty)",
                            resource_count == 0,
                            f"Found {resource_count} resources (expected 0)",
                            f"Expected no resources, but found {resource_count}" if resource_count > 0 else ""
                        )
                        return resource_count == 0
                        
                    except Exception as inner_e:
                        self.log_test(
                            "List Resources (Empty)",
                            False,
                            error=f"Failed to list resources: {str(inner_e)}"
                        )
                        return False
                    
        except Exception as e:
            self.log_test(
                "List Resources (Empty)",
                False,
                error=f"Failed to connect to server: {str(e)}"
            )
            return False

    async def test_execute_python_code_no_presentation(self) -> bool:
        """Test execute_python_code when no presentation is loaded."""
        try:
            server_params = StdioServerParameters(
                command=sys.executable,
                args=[str(self.server_path)]
            )
            
            async with stdio_client(server_params) as (read, write):
                async with ClientSession(read, write) as session:
                    await session.initialize()
                    
                    # Try to execute Python code without a loaded presentation
                    try:
                        result = await session.call_tool(
                            "execute_python_code",
                            arguments={"code": "print('test')"}
                        )
                        
                        result_data = json.loads(result.content[0].text)
                        success = not result_data.get("success", True)
                        expected_error = "No PowerPoint presentation loaded"
                        has_expected_error = expected_error in result_data.get("error", "")
                        
                        test_success = success and has_expected_error
                        
                        self.log_test(
                            "Execute Python Code (No Presentation)",
                            test_success,
                            f"Tool correctly failed with expected error message",
                            f"Unexpected response: {result_data}" if not test_success else ""
                        )
                        return test_success
                        
                    except Exception as inner_e:
                        self.log_test(
                            "Execute Python Code (No Presentation)",
                            False,
                            error=f"Failed to call tool: {str(inner_e)}"
                        )
                        return False
                    
        except Exception as e:
            self.log_test(
                "Execute Python Code (No Presentation)",
                False,
                error=f"Failed to connect to server: {str(e)}"
            )
            return False

    async def test_manual_root_simulation(self) -> bool:
        """
        Test simulated root functionality by manually calling _set_client_roots.
        Note: This is a simulation since FastMCP may not support roots callback yet.
        """
        try:
            # Create test presentation
            test_pptx = self.create_test_presentation()
            
            # Import server functions for manual testing
            from mcp_server.server.session import set_client_roots, get_session
            from mcp import types
            
            # Simulate setting client roots
            roots = [types.Root(uri=f"file://{test_pptx}")]
            set_client_roots(roots)
            
            # Check if presentation was loaded
            session = get_session()
            presentation_loaded = session.loaded_presentation_path is not None
            
            self.log_test(
                "Manual Root Simulation",
                presentation_loaded,
                f"Presentation loaded from: {session.loaded_presentation_path}" if presentation_loaded else "",
                "Failed to load presentation from root" if not presentation_loaded else ""
            )
            
            return presentation_loaded
            
        except Exception as e:
            self.log_test(
                "Manual Root Simulation",
                False,
                error=f"Failed to simulate roots: {str(e)}"
            )
            return False

    async def run_all_tests(self) -> Dict[str, Any]:
        """Run all MEP-003 tests."""
        print("🧪 Starting MEP-003: Root and Resource Management Tests")
        print("=" * 60)
        
        try:
            # Test 1: Server startup
            await self.test_server_startup()
            
            # Test 2: Resource capabilities
            await self.test_resource_capabilities()
            
            # Test 3: List resources (empty)
            await self.test_list_resources_empty()
            
            # Test 4: Execute Python code without presentation
            await self.test_execute_python_code_no_presentation()
            
            # Test 5: Manual root simulation
            await self.test_manual_root_simulation()
            
        except Exception as e:
            print(f"❌ Test suite failed with error: {e}")
        
        finally:
            # Cleanup
            self.cleanup_test_presentation()
        
        # Generate summary
        total_tests = len(self.test_results)
        passed_tests = sum(1 for result in self.test_results if result["success"])
        failed_tests = total_tests - passed_tests
        
        print("\n" + "=" * 60)
        print("📊 Test Summary")
        print(f"Total tests: {total_tests}")
        print(f"Passed: {passed_tests}")
        print(f"Failed: {failed_tests}")
        print(f"Success rate: {(passed_tests/total_tests)*100:.1f}%")
        
        if failed_tests > 0:
            print("\n❌ Failed tests:")
            for result in self.test_results:
                if not result["success"]:
                    print(f"  - {result['test']}: {result['error']}")
        
        return {
            "total": total_tests,
            "passed": passed_tests,
            "failed": failed_tests,
            "success_rate": (passed_tests/total_tests)*100,
            "details": self.test_results
        }


async def main():
    """Main test runner."""
    tester = MEP003Tester()
    
    try:
        results = await tester.run_all_tests()
        
        # Exit with appropriate code
        if results["failed"] == 0:
            print("\n🎉 All tests passed!")
            sys.exit(0)
        else:
            print(f"\n⚠️  {results['failed']} test(s) failed.")
            sys.exit(1)
            
    except KeyboardInterrupt:
        print("\n⏹️  Tests interrupted by user.")
        sys.exit(130)
    except Exception as e:
        print(f"\n💥 Test runner crashed: {e}")
        sys.exit(1)


if __name__ == "__main__":
    asyncio.run(main())