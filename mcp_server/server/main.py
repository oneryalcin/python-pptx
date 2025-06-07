"""
MCP server for python-pptx agentic toolkit.

This server provides AI agents with access to python-pptx library capabilities
through the Model Context Protocol (MCP). Implements MEP-004: Unified Save and Save As Tool.

Features:
- Automatic presentation loading from client roots
- Resource discovery and tree-based content reading
- Simplified execute_python_code tool (no file_path required)
- Unified save_presentation tool with security validation
- Session-based state management for multi-client support
"""

import contextlib
import io
import json
import threading
import time
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional
from urllib.parse import urlparse

from mcp.server import FastMCP
from mcp import types

try:
    import pptx
except ImportError:
    pptx = None

# Initialize the FastMCP server
mcp = FastMCP("pptx-agent-server")

# Define the path to the info document
_INFO_DOC_PATH = Path(__file__).parent.parent / "llm_info.md"


@dataclass
class SessionContext:
    """Represents a client session with its state."""
    session_id: str
    created_at: float
    last_accessed: float
    client_roots: List[types.Root]
    loaded_presentation: Optional[pptx.Presentation] = None
    loaded_presentation_path: Optional[Path] = None
    
    def update_access(self):
        """Update the last accessed timestamp."""
        self.last_accessed = time.time()
    
    def load_presentation(self, file_path: Path) -> bool:
        """Load a presentation from the given file path."""
        if pptx is None:
            return False
            
        try:
            self.loaded_presentation = pptx.Presentation(file_path)
            self.loaded_presentation_path = file_path
            return True
        except Exception:
            self.loaded_presentation = None
            self.loaded_presentation_path = None
            return False


# Session management
_session_store: Dict[str, SessionContext] = {}
_session_lock = threading.Lock()
_current_session_id = threading.local()


def _get_session_id() -> str:
    """
    Get the current session ID. Since FastMCP doesn't provide explicit session context,
    we use a thread-local approach. For now, we create a default session per thread.
    This is a temporary solution until we can implement proper session handling.
    """
    if not hasattr(_current_session_id, 'value'):
        # Create a new session for this thread
        session_id = str(uuid.uuid4())
        _current_session_id.value = session_id
        
        with _session_lock:
            if session_id not in _session_store:
                _session_store[session_id] = SessionContext(
                    session_id=session_id,
                    created_at=time.time(),
                    last_accessed=time.time(),
                    client_roots=[]
                )
    
    return _current_session_id.value


def _get_session() -> SessionContext:
    """Get the current session context."""
    session_id = _get_session_id()
    
    with _session_lock:
        session = _session_store.get(session_id)
        if session:
            session.update_access()
            return session
        else:
            # This shouldn't happen, but create a new session if needed
            session = SessionContext(
                session_id=session_id,
                created_at=time.time(),
                last_accessed=time.time(),
                client_roots=[]
            )
            _session_store[session_id] = session
            return session


def _cleanup_expired_sessions(max_age: float = 3600) -> None:
    """Remove sessions that haven't been used for max_age seconds."""
    current_time = time.time()
    
    with _session_lock:
        expired_keys = [
            session_id for session_id, session in _session_store.items()
            if current_time - session.last_accessed > max_age
        ]
        
        for session_id in expired_keys:
            del _session_store[session_id]


def _set_client_roots(roots: List[types.Root]) -> None:
    """Store the client-provided roots and attempt to load a presentation for the current session."""
    session = _get_session()
    
    # Update session with new roots and clear any existing presentation
    session.client_roots = roots
    session.loaded_presentation = None
    session.loaded_presentation_path = None
    
    # Scan roots for .pptx files and load the first one found
    for root in roots:
        try:
            # Parse the root URI to get the file path
            parsed = urlparse(str(root.uri))
            if parsed.scheme == 'file':
                root_path = Path(parsed.path)
                if root_path.exists():
                    # Search for .pptx files in this root
                    if root_path.is_file() and root_path.suffix.lower() == '.pptx':
                        # Root points directly to a .pptx file
                        if session.load_presentation(root_path):
                            break
                    elif root_path.is_dir():
                        # Search directory for .pptx files
                        pptx_files = list(root_path.glob('*.pptx'))
                        if pptx_files:
                            if session.load_presentation(pptx_files[0]):
                                break
        except Exception:
            # Skip invalid roots
            continue


def _load_presentation(file_path: Path) -> bool:
    """Load a presentation from the given file path for the current session."""
    session = _get_session()
    return session.load_presentation(file_path)


def _validate_output_path(output_path: Path) -> tuple[bool, str]:
    """
    Validate that the output path is within one of the client-provided roots.
    
    Args:
        output_path: The path to validate
        
    Returns:
        tuple[bool, str]: (is_valid, error_message_if_invalid)
    """
    session = _get_session()
    
    if not session.client_roots:
        return False, "No client roots configured. Cannot save files."
    
    # Resolve the output path to handle relative paths and symlinks
    try:
        resolved_output = output_path.resolve()
    except (OSError, ValueError) as e:
        return False, f"Invalid output path: {str(e)}"
    
    # Check if the output path is within any of the client roots
    for root in session.client_roots:
        try:
            # Parse the root URI to get the file path
            parsed = urlparse(str(root.uri))
            if parsed.scheme == 'file':
                root_path = Path(parsed.path).resolve()
                
                # Debug prints (will remove after testing)
                # print(f"DEBUG: root.uri = {root.uri}")
                # print(f"DEBUG: parsed.path = {parsed.path}")
                # print(f"DEBUG: root_path = {root_path}")
                # print(f"DEBUG: resolved_output = {resolved_output}")
                
                # Check if the resolved output path is within this root
                try:
                    relative_path = resolved_output.relative_to(root_path)
                    # print(f"DEBUG: relative_path = {relative_path}")
                    return True, ""  # Path is valid - within this root
                except ValueError as e:
                    # print(f"DEBUG: relative_to failed: {e}")
                    # Path is not within this root, continue checking other roots
                    continue
        except Exception as e:
            # print(f"DEBUG: Exception in root processing: {e}")
            # Skip invalid roots
            continue
    
    return False, f"Output path '{output_path}' is not within any configured client root directory."


@mcp.tool()
async def get_info() -> str:
    """
    You MUST call this tool first before generating any Python code. It provides essential context and examples for interacting with the python-pptx library.
    
    This tool provides the essential "onboarding manual" for AI agents working with
    the python-pptx library. It explains the two-phase workflow (discover, inspect, act)
    and provides concrete code examples.
    
    Returns:
        str: Markdown content with instructions and examples
    """
    try:
        with open(_INFO_DOC_PATH, encoding="utf-8") as f:
            content = f.read()

        if not content.strip():
            return "# Error: Information document is empty."

        return content

    except FileNotFoundError:
        return f"# Error: Information document not found at {_INFO_DOC_PATH}"
    except PermissionError:
        return "# Error: Permission denied reading information document."
    except UnicodeDecodeError:
        return "# Error: Could not decode information document (encoding issue)."
    except Exception as e:
        return f"# Error: Could not read information document.\n\n{str(e)}"


@mcp.tool()
async def execute_python_code(code: str) -> str:
    """
    Execute Python code with the currently loaded PowerPoint presentation available as 'prs'.
    
    The presentation is automatically loaded from the client-provided roots. The loaded
    presentation object is available as 'prs' in the execution context.
    Captures stdout, stderr, and any exceptions during execution.
    
    Args:
        code: Python code to execute
        
    Returns:
        JSON string with execution results including stdout, stderr, and any errors
    """
    start_time = time.time()
    
    # Clean up expired sessions periodically
    _cleanup_expired_sessions()
    
    # Get the current session
    session = _get_session()
    
    # Validate python-pptx is available
    if pptx is None:
        return json.dumps({
            "success": False,
            "stdout": "",
            "stderr": "",
            "error": "python-pptx library is not available",
            "execution_time": time.time() - start_time
        })
    
    # Check if a presentation is loaded in this session
    if session.loaded_presentation is None:
        return json.dumps({
            "success": False,
            "stdout": "",
            "stderr": "",
            "error": "No PowerPoint presentation loaded. Ensure a .pptx file is available in the client roots.",
            "execution_time": time.time() - start_time
        })
    
    # Use the session's pre-loaded presentation
    prs = session.loaded_presentation
    
    # Prepare execution context
    exec_globals = {
        "__builtins__": __builtins__,
        "prs": prs,
        "pptx": pptx,  # Make pptx module available too
        "Path": Path,   # Useful for file operations
        "print": print, # Ensure print works
        "json": json,   # Make json available for output formatting
    }
    
    # Capture stdout and stderr
    stdout_capture = io.StringIO()
    stderr_capture = io.StringIO()
    
    try:
        with contextlib.redirect_stdout(stdout_capture), \
             contextlib.redirect_stderr(stderr_capture):
            # Execute the provided code
            exec(code, exec_globals)
            
        return json.dumps({
            "success": True,
            "stdout": stdout_capture.getvalue(),
            "stderr": stderr_capture.getvalue(),
            "error": None,
            "execution_time": time.time() - start_time
        })
        
    except SyntaxError as e:
        return json.dumps({
            "success": False,
            "stdout": stdout_capture.getvalue(),
            "stderr": stderr_capture.getvalue(),
            "error": f"Syntax error in Python code: {str(e)}",
            "execution_time": time.time() - start_time
        })
        
    except Exception as e:
        return json.dumps({
            "success": False,
            "stdout": stdout_capture.getvalue(),
            "stderr": stderr_capture.getvalue(),
            "error": f"Runtime error: {str(e)}",
            "execution_time": time.time() - start_time
        })


@mcp.tool()
async def save_presentation(output_path: Optional[str] = None) -> str:
    """
    Save the currently loaded PowerPoint presentation to disk. Supports both 'Save' and 'Save As' operations.
    
    If output_path is None, overwrites the original file (Save operation).
    If output_path is provided, saves to the new location (Save As operation).
    All output paths must be within the client-configured root directories for security.
    
    Args:
        output_path: Optional path where to save the presentation. If None, saves to original location.
        
    Returns:
        JSON string with save operation results including success status and file path
    """
    start_time = time.time()
    
    # Clean up expired sessions periodically
    _cleanup_expired_sessions()
    
    # Get the current session
    session = _get_session()
    
    # Validate python-pptx is available
    if pptx is None:
        return json.dumps({
            "success": False,
            "operation": "save",
            "file_path": None,
            "error": "python-pptx library is not available",
            "execution_time": time.time() - start_time
        })
    
    # Check if a presentation is loaded in this session
    if session.loaded_presentation is None:
        return json.dumps({
            "success": False,
            "operation": "save",
            "file_path": None,
            "error": "No PowerPoint presentation loaded. Ensure a .pptx file is available in the client roots.",
            "execution_time": time.time() - start_time
        })
    
    # Determine the target file path
    if output_path is None:
        # Save operation: use the original file path
        if session.loaded_presentation_path is None:
            return json.dumps({
                "success": False,
                "operation": "save",
                "file_path": None,
                "error": "Cannot determine original file path for save operation.",
                "execution_time": time.time() - start_time
            })
        target_path = session.loaded_presentation_path
        operation = "save"
    else:
        # Save As operation: use the provided output path
        target_path = Path(output_path)
        operation = "save_as"
    
    # Validate the target path is within client roots
    is_valid, validation_error = _validate_output_path(target_path)
    if not is_valid:
        return json.dumps({
            "success": False,
            "operation": operation,
            "file_path": str(target_path),
            "error": f"Security validation failed: {validation_error}",
            "execution_time": time.time() - start_time
        })
    
    # Ensure the target directory exists
    try:
        target_path.parent.mkdir(parents=True, exist_ok=True)
    except Exception as e:
        return json.dumps({
            "success": False,
            "operation": operation,
            "file_path": str(target_path),
            "error": f"Failed to create target directory: {str(e)}",
            "execution_time": time.time() - start_time
        })
    
    # Attempt to save the presentation
    try:
        session.loaded_presentation.save(str(target_path))
        
        # Update session state if this was a Save As operation
        if operation == "save_as":
            session.loaded_presentation_path = target_path
        
        return json.dumps({
            "success": True,
            "operation": operation,
            "file_path": str(target_path),
            "error": None,
            "execution_time": time.time() - start_time
        })
        
    except PermissionError:
        return json.dumps({
            "success": False,
            "operation": operation,
            "file_path": str(target_path),
            "error": "Permission denied: Cannot write to target file. Check file permissions and ensure the file is not open in another application.",
            "execution_time": time.time() - start_time
        })
        
    except Exception as e:
        return json.dumps({
            "success": False,
            "operation": operation,
            "file_path": str(target_path),
            "error": f"Failed to save presentation: {str(e)}",
            "execution_time": time.time() - start_time
        })


@mcp.resource("pptx://presentation")
async def get_presentation_tree() -> str:
    """Get the tree structure of the currently loaded PowerPoint presentation."""
    # Clean up expired sessions periodically
    _cleanup_expired_sessions()
    
    # Get the current session
    session = _get_session()
    
    if session.loaded_presentation is None or session.loaded_presentation_path is None:
        return json.dumps({
            "error": "No presentation loaded",
            "message": "Ensure a .pptx file is available in the client roots.",
            "session_id": session.session_id  # Include session ID for debugging
        }, indent=2)
    
    # Return get_tree() output if available, otherwise return basic info
    try:
        if hasattr(session.loaded_presentation, 'get_tree'):
            tree_data = session.loaded_presentation.get_tree()
            return json.dumps(tree_data, indent=2)
        else:
            # Fallback: basic presentation info
            info = {
                "type": "presentation",
                "slide_count": len(session.loaded_presentation.slides),
                "file_path": str(session.loaded_presentation_path),
                "session_id": session.session_id,
                "note": "get_tree() method not available. Please ensure you have the latest python-pptx with introspection features."
            }
            return json.dumps(info, indent=2)
    except Exception as e:
        # If get_tree() fails, return error info
        error_info = {
            "type": "presentation", 
            "slide_count": len(session.loaded_presentation.slides),
            "file_path": str(session.loaded_presentation_path),
            "session_id": session.session_id,
            "error": f"Failed to get tree data: {str(e)}"
        }
        return json.dumps(error_info, indent=2)


if __name__ == "__main__":
    # Run the server using the stdio transport
    mcp.run(transport='stdio')
