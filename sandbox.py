import rich
from pptx import Presentation

prs = Presentation("/tmp/single_slide.pptx")
# prs.core_properties.title = "My Presentation"
# prs.core_properties.author = "AI Assistant"
# slide = prs.slides.add_slide(prs.slide_layouts[0])
# slide.shapes.title.text = "Welcome Slide"

result = prs.to_dict(max_depth=3
                     , expand_collections=False)

rich.print(result)